#!/usr/bin/env python3
"""Build the editable 10-slide Korean technical-conference deck.

Text, annotations, architecture blocks, timelines, bars, arrows, and badges are
native PowerPoint objects. Raster inputs are limited to the decorative cover
and native KiCad render/crops. The same scene operations render slide PNGs and
the PDF, keeping review artifacts aligned with the PPTX.
"""

from __future__ import annotations

import csv
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
import shutil
import subprocess
import textwrap
import zipfile

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt
from reportlab import rl_config
from reportlab.pdfgen import canvas as pdf_canvas

rl_config.invariant = 1


ROOT = Path(__file__).resolve().parents[2]
FINAL = ROOT / "presentation" / "final"
ASSETS = FINAL / "assets"
SLIDES = FINAL / "slides"
CACHE = FINAL / ".work" / "deck_cache"
PPTX_PATH = FINAL / "presentation.pptx"
PDF_PATH = FINAL / "presentation.pdf"

W, H = 1920, 1080
PX_PER_IN = 144
FONT_REGULAR = "/usr/share/fonts/truetype/nanum/NanumSquareR.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/nanum/NanumSquareB.ttf"
PPT_FONT = "NanumSquare"

BG = "#07111F"
BG2 = "#0B1627"
BG3 = "#102238"
WHITE = "#F4F7FB"
MUTED = "#91A5BA"
CYAN = "#36D7E7"
TEAL = "#28B6A6"
BLUE = "#4C8DFF"
AMBER = "#F1B44C"
RED = "#FF6B6B"
GRAY = "#5E7188"
GRID = "#20364D"


def rgb(value: str) -> tuple[int, int, int]:
    value = value.lstrip("#")
    return tuple(int(value[index : index + 2], 16) for index in (0, 2, 4))


def ppt_rgb(value: str) -> RGBColor:
    return RGBColor(*rgb(value))


def inch(px: float) -> Inches:
    return Inches(px / PX_PER_IN)


def font(size_pt: float, bold: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REGULAR, max(10, round(size_pt * 2)))


def wrap_for_width(draw: ImageDraw.ImageDraw, value: str, face: ImageFont.FreeTypeFont, width: int) -> str:
    lines: list[str] = []
    for paragraph in value.split("\n"):
        words = paragraph.split()
        if not words:
            lines.append("")
            continue
        current = words[0]
        for word in words[1:]:
            trial = f"{current} {word}"
            if draw.textbbox((0, 0), trial, font=face)[2] <= width:
                current = trial
            else:
                lines.append(current)
                current = word
        lines.append(current)
    return "\n".join(lines)


class SlideCanvas:
    def __init__(self, prs: Presentation, title: str, notes: str):
        self.prs = prs
        self.slide = prs.slides.add_slide(prs.slide_layouts[6])
        self.image = Image.new("RGB", (W, H), rgb(BG))
        self.draw = ImageDraw.Draw(self.image)
        self.title = title
        self.notes = notes
        fill = self.slide.background.fill
        fill.solid()
        fill.fore_color.rgb = ppt_rgb(BG)

    def rect(self, x: float, y: float, w: float, h: float, fill: str, *, stroke: str | None = None, stroke_width: float = 1.0, radius: int = 0, alpha: int = 255) -> None:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shape = self.slide.shapes.add_shape(shape_type, inch(x), inch(y), inch(w), inch(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(fill)
        shape.fill.transparency = round((255 - alpha) / 255 * 100)
        if stroke:
            shape.line.color.rgb = ppt_rgb(stroke)
            shape.line.width = Pt(stroke_width)
        else:
            shape.line.fill.background()
        box = (round(x), round(y), round(x + w), round(y + h))
        if alpha == 0:
            if radius:
                self.draw.rounded_rectangle(box, radius=radius, fill=None, outline=rgb(stroke) if stroke else None, width=max(1, round(stroke_width * 2)))
            else:
                self.draw.rectangle(box, fill=None, outline=rgb(stroke) if stroke else None, width=max(1, round(stroke_width * 2)))
        elif alpha < 255:
            overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
            overlay_draw = ImageDraw.Draw(overlay)
            fill_rgba = (*rgb(fill), alpha)
            outline_rgba = (*rgb(stroke), 255) if stroke else None
            if radius:
                overlay_draw.rounded_rectangle(box, radius=radius, fill=fill_rgba, outline=outline_rgba, width=max(1, round(stroke_width * 2)))
            else:
                overlay_draw.rectangle(box, fill=fill_rgba, outline=outline_rgba, width=max(1, round(stroke_width * 2)))
            self.image = Image.alpha_composite(self.image.convert("RGBA"), overlay).convert("RGB")
            self.draw = ImageDraw.Draw(self.image)
        elif radius:
            self.draw.rounded_rectangle(box, radius=radius, fill=rgb(fill), outline=rgb(stroke) if stroke else None, width=max(1, round(stroke_width * 2)))
        else:
            self.draw.rectangle(box, fill=rgb(fill), outline=rgb(stroke) if stroke else None, width=max(1, round(stroke_width * 2)))

    def line(self, x1: float, y1: float, x2: float, y2: float, color: str, *, width: float = 2.0) -> None:
        shape = self.slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, inch(x1), inch(y1), inch(x2), inch(y2))
        shape.line.color.rgb = ppt_rgb(color)
        shape.line.width = Pt(width)
        self.draw.line((x1, y1, x2, y2), fill=rgb(color), width=max(1, round(width * 2)))

    def circle(self, x: float, y: float, d: float, fill: str, *, stroke: str | None = None, stroke_width: float = 1.0) -> None:
        shape = self.slide.shapes.add_shape(MSO_SHAPE.OVAL, inch(x), inch(y), inch(d), inch(d))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(fill)
        if stroke:
            shape.line.color.rgb = ppt_rgb(stroke)
            shape.line.width = Pt(stroke_width)
        else:
            shape.line.fill.background()
        self.draw.ellipse((x, y, x + d, y + d), fill=rgb(fill), outline=rgb(stroke) if stroke else None, width=max(1, round(stroke_width * 2)))

    def chevron(self, x: float, y: float, w: float, h: float, fill: str) -> None:
        shape = self.slide.shapes.add_shape(MSO_SHAPE.CHEVRON, inch(x), inch(y), inch(w), inch(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(fill)
        shape.line.fill.background()
        points = [(x, y), (x + w * 0.65, y), (x + w, y + h / 2), (x + w * 0.65, y + h), (x, y + h), (x + w * 0.35, y + h / 2)]
        self.draw.polygon(points, fill=rgb(fill))

    def text(self, value: str, x: float, y: float, w: float, h: float, *, size: float = 22, color: str = WHITE, bold: bool = False, align: str = "left", valign: str = "top", margin: float = 0) -> None:
        # Insert deterministic line breaks with the same font metrics used by
        # the PNG renderer.  PowerPoint/LibreOffice therefore never decide a
        # different wrap point, which was the main source of overflow in v1.
        face = font(size, bold)
        wrapped = wrap_for_width(self.draw, value, face, max(1, round((w - margin * 2) * 0.90)))
        line_count = max(1, wrapped.count("\n") + 1)
        h = max(h, line_count * size * 2 * 1.08 + margin * 2 + 4)
        if y + h > H:
            y = H - h
        box = self.slide.shapes.add_textbox(inch(x), inch(y), inch(w), inch(h))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = False
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = inch(margin)
        frame.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}[valign]
        paragraph = frame.paragraphs[0]
        paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
        paragraph.space_before = paragraph.space_after = Pt(0)
        run = paragraph.add_run()
        run.text = wrapped
        run.font.name = PPT_FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = ppt_rgb(color)
        bbox = self.draw.multiline_textbbox((0, 0), wrapped, font=face, spacing=round(size * 0.35), align=align)
        text_h = bbox[3] - bbox[1]
        py = y + margin
        if valign == "middle":
            py = y + (h - text_h) / 2
        elif valign == "bottom":
            py = y + h - text_h - margin
        if align == "center":
            px = x + w / 2
            anchor = "ma"
        elif align == "right":
            px = x + w - margin
            anchor = "ra"
        else:
            px = x + margin
            anchor = "la"
        self.draw.multiline_text((px, py), wrapped, font=face, fill=rgb(color), spacing=round(size * 0.35), align=align, anchor=anchor)

    def image_file(self, path: Path, x: float, y: float, w: float, h: float, *, mode: str = "cover") -> None:
        source = Image.open(path).convert("RGB")
        target_size = (max(1, round(w)), max(1, round(h)))
        if mode == "cover":
            rendered = ImageOps.fit(source, target_size, method=Image.Resampling.LANCZOS)
        else:
            rendered = Image.new("RGB", target_size, rgb(BG))
            thumb = source.copy()
            thumb.thumbnail(target_size, Image.Resampling.LANCZOS)
            rendered.paste(thumb, ((target_size[0] - thumb.width) // 2, (target_size[1] - thumb.height) // 2))
        CACHE.mkdir(parents=True, exist_ok=True)
        cache_path = CACHE / f"{len(self.slide.shapes):03d}_{path.stem}_{round(w)}x{round(h)}.png"
        rendered.save(cache_path)
        self.slide.shapes.add_picture(str(cache_path), inch(x), inch(y), inch(w), inch(h))
        self.image.paste(rendered, (round(x), round(y)))

    def save(self, index: int) -> Path:
        self.slide.notes_slide.notes_text_frame.text = self.notes
        path = SLIDES / f"slide_{index:02d}.png"
        self.image.save(path, quality=95)
        return path


@dataclass(frozen=True)
class SlideMeta:
    title: str
    section: str
    conclusion: str
    duration: str
    notes: str


SLIDE_META = (
    SlideMeta("정적 큐의 Tail을 줄이는 Work Stealing", "RESEARCH QUESTION", "Gemma 3 1B 작업에서 Tail 감소와 원격 이동 비용이 교차하는 조건을 묻는다.", "0:35", """기억할 문장: 이 논문은 Work Stealing이 항상 빠르다는 주장이 아니라, 실제 LLM 작업에서 언제 Tail 감소가 이동 비용보다 큰지를 묻습니다.

오늘 말씀드릴 질문은 단순히 메모리 대역폭이 중요한가가 아닙니다. 실제 Gemma 3 1B 작업을 여러 연산 클러스터에 나눌 때, 정적 소유권 때문에 일부 큐가 길어지고 다른 클러스터가 쉬는 상황을 어떻게 줄일 것인가입니다.

제안 구조는 K26 연산 SoC와 외부 Memory FPGA를 결합하고, Multi-Queue FCFS 위에 지역성 비용을 고려한 Work Stealing을 적용합니다. 핵심은 무조건 작업을 옮기는 것이 아니라, Tail을 줄이는 이득이 링크·메모리 이동 비용보다 큰 조건을 찾는 것입니다.

먼저 정적 큐에서 왜 Tail이 생기는지 보겠습니다."""),
    SlideMeta("연구 질문: 지역성과 부하 균형을 함께 얻을 수 있는가", "MOTIVATION", "정적 소유권의 locality 이점과 queue skew의 Tail 비용을 함께 측정해야 한다.", "0:55", """기억할 문장: 연구 공백은 locality와 load balance를 따로 보는 것이 아니라 p95·p99와 원격 byte를 같은 작업에서 연결하는 데 있습니다.

S1 정적 로컬 큐에서는 각 TileJob이 home cluster의 FCFS 큐에 들어갑니다. 이 방식은 원격 가중치 이동을 만들지 않는 장점이 있지만, skew가 생기면 Cluster 0의 큐만 길어지고 나머지 클러스터는 할 일이 없어집니다.

평균 지연만 보면 이 현상이 약하게 보일 수 있습니다. 하지만 사용자 체감과 시스템의 느린 요청을 결정하는 p95와 p99에서는 긴 queue wait가 그대로 드러납니다. 따라서 이 연구의 목표는 평균 처리량보다 Tail을 줄이면서 이동 비용을 통제하는 것입니다.

이 문제를 풀기 위해 계산 소유권과 메모리 친화도를 분리한 구조를 설계했습니다."""),
    SlideMeta("연구 대상 구조와 구현 경계", "SYSTEM MODEL", "K26의 계산 소유권과 Memory FPGA의 데이터 친화도를 TileJob에서 분리한다.", "1:10", """기억할 문장: 논문의 시스템 모델은 계산 위치와 데이터 위치를 분리하지만, 닫힌 DDR 응답 경로까지 구현했다고 주장하지 않습니다.

왼쪽 K26에는 TileScheduler, payload store, 네 개의 compute cluster와 signed INT8 16×4 MatVec가 있습니다. TileJob은 작업 identity, K/N 범위, weight·activation 주소, preferred channel과 link bundle을 보존합니다.

오른쪽 Memory FPGA 후보는 네 DDR3L 채널과 네 link bundle을 통해 weight service를 분리하는 구조입니다. 스케줄러는 어느 cluster가 계산할지 정하고, channel affinity는 데이터 이동 비용을 계산하는 기준이 됩니다.

현재 runnable RTL은 scheduler에서 payload store와 MatVec 결과까지입니다. DDR response와 link receive의 완전한 폐루프는 다음 검증 단계이며, 오늘의 중심은 이 구조 위에서 scheduling 조건을 분석한 결과입니다.

이제 유휴 클러스터가 어떤 규칙으로 작업을 가져오는지 보겠습니다."""),
    SlideMeta("지역성 비용을 반영한 Work Stealing", "METHOD", "S3는 job age에서 원격 이동 비용을 뺀 점수가 양수인 작업만 훔친다.", "1:15", """기억할 문장: S3는 가장 오래된 작업을 그대로 가져오지 않고, 대기 이득에서 데이터 이동 비용을 뺀 선택을 합니다.

첫째, local queue가 비면 cluster가 유휴 상태임을 확인합니다. 둘째, 다른 victim queue에서 실행 가능한 작업을 탐색합니다. 셋째, job age에서 원격 weight, activation, reduction owner와 bundle mismatch 비용을 뺀 locality score를 계산합니다.

넷째, score가 양수인 eligible TileJob을 선택해 유휴 cluster로 옮깁니다. 마지막으로 job identity를 유지해 정확히 한 번만 completion이 발생하는지 검사합니다.

비교 기준인 S2는 가장 오래된 eligible job을 바로 가져옵니다. S3는 이동량을 줄이기 위해 지역성 비용을 함께 보므로 Tail과 remote traffic 사이의 절충이 생깁니다.

이 다섯 단계는 정책 개념을 설명합니다. 분석 모델은 victim queue의 모든 eligible job을 탐색하지만 현재 RTL은 victim head만 검사하며 ownership과 score 정의도 다릅니다. 따라서 model과 RTL이 알고리즘적으로 동일하다고 해석하면 안 됩니다.

정책을 비교하려면 실제 모델 작업을 동일한 TileJob ledger로 만들어야 합니다."""),
    SlideMeta("실제 Gemma graph에서 평가 작업을 만든다", "WORKLOAD", "실제 projection 순서를 보존한 5,856개 TileJob과 조건 탐색용 stress를 분리한다.", "1:00", """기억할 문장: 실제 Gemma replay는 연구의 현실성을, 별도 synthetic stress는 효과가 발생하는 조건을 검증합니다.

해시로 고정한 Gemma 3 1B ONNX graph에는 7,837개 node가 있습니다. 여기서 attention의 q·k·v·o, MLP의 gate·up·down, lm_head를 필터링하면 token당 183개 projection이 됩니다.

Decode-32 조건에서는 183 곱하기 32, 즉 5,856개의 coarse TileJob ledger를 만듭니다. 이 ledger로 실제 Gemma replay를 만들고, 별도 stress로 효과 조건을 분리해 확인했습니다.

다음 세 슬라이드는 이 한 개 replay를 반복하지 않습니다. 정책이 유효한 조건을 보기 위해 별도로 생성한 1,000-job 치우침 스트레스와 5개 seed를 사용합니다. 모든 비교에서 input, dispatch, completion ID가 같고 중복과 누락이 0인지 확인했습니다.

다음 슬라이드는 전체 2,000개 event를 축소한 그림이 아니라, 실제 차이가 시작되는 동일 시간 구간을 확대합니다."""),
    SlideMeta("동일 조건에서 S1과 S3의 실행을 비교한다", "EXPERIMENT DESIGN", "같은 job stream에서 정책만 바꾸면 유휴 구간이 steal 실행으로 전환된다.", "1:05", """기억할 문장: 공정한 비교를 위해 같은 job identity와 resource 조건을 고정하고 scheduler 정책만 바꿨습니다.

여기부터는 Gemma replay와 별개의 synthetic scheduler 실험입니다. 각 workload는 1,000 job이고 seed 19, 23, 29, 31, 43을 사용합니다. full-overlap analytical model이며 물리 cycle이 아닙니다. 타임라인은 seed 23의 S1과 S3, 총 2,000개 event 중 동일한 41k에서 43k 구간을 확대했습니다.

S1에서는 Cluster 0만 파란 로컬 작업을 처리하고 Cluster 1부터 3은 유휴입니다. S3에서는 J172, J162, J143이 각각 C0에서 C1, C2, C3으로 이동합니다. 회색선은 큐 대기, 어두운 색 구간은 링크·메모리 준비, 밝은 끝 구간은 compute입니다.

전체 5개 seed 중앙값에서 skew p95는 S1 대비 18.12% 줄었습니다. 그 대가로 S3에는 약 1.45 MiB의 원격 가중치 이동이 생깁니다. 즉 유휴를 줄인 효과와 이동 비용을 함께 봐야 합니다.

이제 p95뿐 아니라 p99도 같은 방향인지 확인하겠습니다."""),
    SlideMeta("실제 Gemma replay에서도 Tail이 감소했다", "RESULT", "실제 replay의 p95·p99 감소를 확인하고 별도 stress로 효과 조건을 분리했다.", "1:20", """기억할 문장: 실제 Gemma replay의 감소와 synthetic skew의 조건 분석이 같은 방향을 보이지만, 두 수치를 같은 실험으로 합치면 안 됩니다.

먼저 실제 Gemma decode-32 graph-derived ledger의 analytical replay입니다. 5,856개 coarse TileJob에서 S3는 S1보다 p95를 15.07%, p99를 14.61% 줄였습니다. 이 값은 실제 graph order와 tensor byte에서 파생됐지만 시간은 보드 측정이 아니라 분석 모델입니다.

아래 보조 수치인 18.12%는 별도의 1,000-job synthetic skew, 5개 seed 중앙값입니다. 실제 replay와 동일한 실험으로 합치는 값이 아니라, queue skew가 Tail 감소의 조건임을 확인하는 stress 결과입니다.

Mixed workload에서도 p95는 17.59% 줄고 p99도 감소합니다. 반면 balanced와 channel-hotspot에서는 stealing할 불균형이 없거나 memory channel 자체가 병목이어서 S3의 이득이 나타나지 않습니다.

따라서 결론은 S3가 항상 빠르다는 것이 아닙니다. 실제 graph-derived replay에서 방향성을 확인하고, 별도 stress로 정적 ownership의 queue skew가 효과 조건임을 분리해 확인했다는 것입니다.

그렇다면 줄어든 queue wait의 비용이 어디로 이동하는지 보겠습니다."""),
    SlideMeta("Tail 감소의 비용은 원격 이동이다", "TRADE-OFF", "S3는 S2보다 원격 가중치를 줄이지만 마지막 completion은 조금 늦어진다.", "0:55", """기억할 문장: 지역성 점수는 원격 이동을 줄이지만, 보수적인 선택 때문에 전체 완료시간은 소폭 늘 수 있습니다.

이 수치도 Gemma replay가 아니라 같은 synthetic skew, 5-seed 중앙값, full-overlap analytical model에서 나옵니다. 왼쪽은 S3와 S1의 비교로 p95가 18.12% 감소합니다. 실제 평균 queue wait 감소율은 17.86%이므로 p95와 같은 지표로 해석하면 안 됩니다. 가운데는 S3와 S2의 비교로 remote weight가 37.84% 줄어듭니다.

오른쪽도 S3와 S2의 비교입니다. 더 보수적인 선택 때문에 마지막 completion은 0.98% 길어집니다. 즉 세 수치는 비교 기준이 서로 다릅니다. S1은 정적 baseline이고, S2는 locality를 고려하지 않는 stealing baseline입니다.

그리고 service composition을 sequential 경계로 바꾸면 synthetic skew의 S2/S3 p95 순위가 뒤집힙니다. 따라서 S3의 S2 대비 Tail 우위는 서비스 중첩 가정에 조건부입니다.

이 결과는 queue wait가 줄어든 대신 link와 memory service가 더 중요해졌음을 보여줍니다. 다음 단계는 이 비용을 실제 보드 인터페이스에서 계측하는 것입니다."""),
    SlideMeta("물리 참조 설계로 다음 검증 범위를 고정했다", "BOUNDED PHYSICAL EVIDENCE", "Native KiCad coupon으로 링크와 제한 검증 범위를 고정했다.", "0:55", """기억할 문장: KiCad coupon의 역할은 제품 보드를 보여주는 것이 아니라 다음 물리 검증의 인터페이스와 범위를 고정하는 것입니다.

가운데 큰 이미지는 native KiCad source에서 렌더한 validation coupon입니다. 오른쪽 확대는 K26–Memory FPGA 연결 경계, 기준 클록 차동 경로, 대표 routed coupon 영역을 실제 render에서 잘라 보여줍니다.

현재 coupon에는 29 footprints, 20개의 routed GTH/refclock nets가 있고 선언한 제한 범위에서 ERC와 routed-subset DRC가 0입니다. 그러나 전체 보드에는 55 unrouted nets가 남아 있고, MIG pin placement, SI/PI, PDN, thermal과 timing closure는 수행하지 않았습니다.

따라서 이 이미지는 알고리즘을 어떤 인터페이스로 내려갈지 구체화한 reference coupon입니다. 제작 가능성을 주장하지 않고, 다음 검증의 대상과 범위를 명확히 한 것입니다."""),
    SlideMeta("기여와 한계: 조건부 효과를 규명했다", "CONCLUSION", "실제 workload·스케줄러·제한 물리 근거를 하나의 주장 사슬로 연결했다.", "0:50", """기억할 문장: 이 논문의 기여는 Work Stealing의 보편적 우월성이 아니라 효과가 생기는 조건과 아직 남은 검증을 함께 제시한 것입니다.

실제 Gemma graph에서 TileJob 변환을 정의했고, 별도의 제어된 스트레스에서 Tail과 이동 비용 조건을 분석했으며, 그 구조를 K26–Memory FPGA RTL과 KiCad 참조 설계까지 연결했습니다.

다음 단계는 DDR 응답부터 MatVec까지 폐루프를 완성하고 실제 대역폭·Tail·보드 전력을 측정하는 것입니다.

정적 큐의 유휴를 줄이면 Tail은 짧아지지만 링크·메모리 비용이 생깁니다. 이 연구는 그 재배치가 유효한 조건을 함께 제시합니다. 질문 받겠습니다."""),
)


def add_header(c: SlideCanvas, meta: SlideMeta) -> None:
    c.rect(86, 54, 54, 7, CYAN)
    c.text(meta.title, 86, 76, 1740, 80, size=34, bold=True)
    c.text(meta.conclusion, 88, 156, 1680, 58, size=22, color=MUTED)


def metric(c: SlideCanvas, x: int, y: int, label: str, value: str, color: str, w: int = 420) -> None:
    c.text(value, x, y, w, 82, size=36, color=color, bold=True, align="center")
    c.text(label, x, y + 86, w, 42, size=18, color=MUTED, align="center")


def slide_1(c: SlideCanvas) -> None:
    c.image_file(ASSETS / "cover_background.png", 0, 0, W, H, mode="cover")
    c.rect(0, 0, 1120, H, BG, alpha=70)
    c.rect(88, 245, 70, 8, CYAN)
    c.text("정적 큐의 Tail을 줄이는\nWork Stealing", 88, 282, 980, 230, size=43, bold=True)
    c.text("Gemma 3 1B · K26–Memory FPGA 논문", 92, 550, 1000, 58, size=22, color="#C8D5E3")
    c.text("연구 질문  |  Tail 감소는 언제 이동 비용보다 큰가?", 92, 630, 1000, 58, size=22, color=CYAN, bold=True)
    c.text("최윤혁 · 한국디지털미디어고등학교", 92, 940, 760, 42, size=18, color="#8298AE")


def slide_2(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[1])
    c.text("S1 정적 로컬 큐", 108, 258, 630, 45, size=22, color=CYAN, bold=True)
    lane_y = [360, 470, 580, 690]
    queues = [9, 2, 1, 1]
    for cluster, (y, count) in enumerate(zip(lane_y, queues)):
        c.text(f"클러스터 {cluster}", 100, y - 8, 205, 48, size=19, color=WHITE, bold=True)
        c.line(315, y + 18, 790, y + 18, GRID, width=2)
        for index in range(count):
            color = BLUE if cluster == 0 else GRAY
            c.rect(327 + index * 50, y - 3, 40, 42, color, radius=7)
        if cluster > 0:
            c.text("유휴", 650, y - 4, 120, 40, size=17, color=AMBER, bold=True, align="right")
    c.chevron(840, 488, 120, 86, CYAN)
    c.text("지역성 유지", 814, 420, 176, 40, size=18, color=MUTED, align="center")
    c.text("부하 균형 실패", 803, 590, 198, 40, size=18, color=AMBER, bold=True, align="center")
    c.text("느린 5%가 길어진다", 1090, 278, 640, 52, size=25, color=WHITE, bold=True)
    points = [(1100, 720), (1190, 710), (1280, 690), (1370, 650), (1460, 590), (1550, 490), (1640, 350), (1750, 270)]
    for p1, p2 in zip(points, points[1:]):
        c.line(*p1, *p2, BLUE, width=4)
    c.line(1590, 315, 1590, 760, CYAN, width=2)
    c.text("p95", 1524, 735, 132, 46, size=22, color=CYAN, bold=True, align="center")
    c.text("연구 공백  |  실제 LLM graph · p95/p99 · 원격 byte를 한 조건에서 연결", 1000, 832, 820, 92, size=22, color=AMBER, bold=True, align="center")


def slide_3(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[2])
    c.text("K26 연산 SoC", 110, 260, 610, 48, size=24, color=CYAN, bold=True)
    c.rect(100, 330, 690, 510, BG2, stroke="#284765", stroke_width=2, radius=24)
    c.circle(145, 385, 140, BG3, stroke=CYAN, stroke_width=2)
    c.text("타일\n스케줄러", 145, 420, 140, 60, size=15, bold=True, align="center", valign="middle")
    c.rect(330, 390, 370, 72, BG3, stroke=TEAL, stroke_width=2, radius=16)
    c.text("페이로드 저장소", 346, 405, 338, 42, size=20, color=TEAL, bold=True, align="center")
    for idx in range(4):
        x = 130 + (idx % 2) * 325
        y = 560 + (idx // 2) * 120
        c.rect(x, y, 285, 82, BG3, stroke=BLUE, stroke_width=2, radius=16)
        c.text(f"클러스터 {idx}\nMatVec", x + 12, y + 14, 261, 56, size=15, color=WHITE, bold=True, align="center", valign="middle")
    c.text("요청·제어 →", 780, 445, 250, 46, size=18, color=CYAN, bold=True, align="center")
    c.text("← 가중치 응답", 780, 535, 250, 46, size=18, color=TEAL, bold=True, align="center")
    c.text("4-bundle 링크", 780, 615, 250, 40, size=16, color=WHITE, bold=True, align="center")
    c.text("응답 폐루프 미구현", 780, 675, 250, 36, size=14, color=AMBER, align="center")
    c.text("Memory FPGA", 1070, 260, 650, 48, size=24, color=TEAL, bold=True)
    c.rect(1050, 330, 760, 510, BG2, stroke="#28594F", stroke_width=2, radius=24)
    c.rect(1100, 380, 660, 88, BG3, stroke=TEAL, stroke_width=2, radius=16)
    c.text("채널 친화도 + Bank 인식 큐", 1120, 402, 620, 44, size=20, color=TEAL, bold=True, align="center")
    for idx in range(4):
        x = 1110 + idx * 160
        c.rect(x, 550, 132, 185, BG3, stroke=BLUE if idx % 2 == 0 else TEAL, stroke_width=2, radius=14)
        c.text(f"DDR3L\nCH {idx}", x + 8, 598, 116, 76, size=17, bold=True, align="center", valign="middle")
    c.text("TileJob = 계산 소유권 + 데이터 친화도", 430, 902, 1060, 62, size=25, color=WHITE, bold=True, align="center")


def slide_4(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[3])
    c.image_file(ASSETS / "work_stealing_sequence_frame.png", 120, 235, 1680, 660, mode="contain")
    c.line(260, 938, 1660, 938, GRID, width=2)
    c.text("score = job age - remote weight - activation/reduction - bundle mismatch", 180, 960, 1560, 48, size=20, color=CYAN, bold=True, align="center")


def slide_5(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[4])
    items = [
        ("7,837", "ONNX 그래프 노드", BLUE),
        ("183", "프로젝션 / 토큰", TEAL),
        ("5,856", "TileJob · 디코드 32", CYAN),
    ]
    xs = [150, 720, 1290]
    for index, (value, label, color) in enumerate(items):
        c.circle(xs[index] + 80, 375, 260, BG2, stroke=color, stroke_width=3)
        c.text(value, xs[index] + 95, 438, 230, 80, size=40, color=color, bold=True, align="center")
        c.text(label, xs[index] + 40, 655, 340, 55, size=21, color=WHITE, bold=True, align="center")
        if index < 2:
            c.chevron(xs[index] + 420, 475, 120, 62, MUTED)
    c.text("실제 그래프", 225, 780, 270, 44, size=18, color=MUTED, align="center")
    c.text("프로젝션 필터", 795, 780, 270, 44, size=18, color=MUTED, align="center")
    c.text("Gemma replay\n5,856개 작업", 1345, 758, 360, 68, size=16, color=MUTED, align="center")
    c.line(300, 850, 1620, 850, GRID, width=2)
    c.text("Gemma replay  ≠  조건 분석용 합성 스트레스  ·  다음 결과는 별도 작업 집합", 260, 885, 1400, 52, size=19, color=CYAN, bold=True, align="center")


def load_events() -> list[dict[str, str]]:
    with (ASSETS / "s1_s3_timeline_events.csv").open(encoding="utf-8", newline="") as stream:
        return list(csv.DictReader(stream))


def timeline_panel(c: SlideCanvas, x: int, y: int, w: int, title: str, policy: str, events: list[dict[str, str]]) -> dict[int, tuple[float, float]]:
    start, end = 41_000, 43_000
    c.text(title, x, y, w, 44, size=22, color=CYAN if policy == "S3" else BLUE, bold=True)
    c.text("동일 2k cycle 구간", x, y + 58, w, 32, size=15, color=MUTED, align="right")
    lane_top = y + 105
    chosen: list[dict[str, str]] = []
    if policy == "S1":
        chosen = [row for row in events if row["scheduler"] == "S1" and start <= int(row["dispatch_cycle"]) < end][:5]
    else:
        for cluster in (1, 2, 3):
            matches = [row for row in events if row["scheduler"] == "S3" and row["stolen"] == "1" and int(row["dispatch_cluster"]) == cluster and start <= int(row["dispatch_cycle"]) < end]
            chosen.append(matches[0])
        local = [row for row in events if row["scheduler"] == "S3" and row["stolen"] == "0" and int(row["dispatch_cluster"]) == 0 and start <= int(row["dispatch_cycle"]) < end]
        if local:
            chosen.append(local[0])
    stolen_positions: dict[int, tuple[float, float]] = {}
    for cluster in range(4):
        ly = lane_top + cluster * 87
        c.text(f"C{cluster}", x, ly + 12, 54, 34, size=17, color=WHITE, bold=True)
        c.rect(x + 62, ly, w - 62, 58, BG3, radius=8)
        lane_events = [row for row in chosen if int(row["dispatch_cluster"]) == cluster]
        if not lane_events:
            c.rect(x + 62, ly, w - 62, 58, "#121B28", radius=8)
            c.text("유휴", x + w - 105, ly + 12, 84, 34, size=16, color=AMBER, bold=True, align="right")
        for row in lane_events:
            scale = (w - 62) / (end - start)
            dispatch = max(start, int(row["dispatch_cycle"]))
            compute_start = min(end, int(row["compute_start_cycle"]))
            finish = min(end, int(row["compute_end_cycle"]))
            bx = x + 62 + (dispatch - start) * scale
            compute_x = x + 62 + (compute_start - start) * scale
            finish_x = x + 62 + (finish - start) * scale
            c.line(x + 62, ly + 51, bx, ly + 51, GRAY, width=4)
            is_stolen = row["stolen"] == "1"
            color = TEAL if is_stolen else BLUE
            service_color = "#176A64" if is_stolen else "#234F96"
            service_w = max(8, min(compute_x, x + w) - bx)
            compute_w = max(10, min(finish_x, x + w) - compute_x)
            c.rect(bx, ly + 8, service_w, 42, service_color, stroke=color, stroke_width=2, radius=7)
            c.rect(min(compute_x, x + w - 10), ly + 8, min(compute_w, x + w - min(compute_x, x + w - 10)), 42, color, radius=5)
            if service_w > 95:
                tag = f"J{row['job_id']} · C{row['home_cluster']}→C{row['dispatch_cluster']}" if is_stolen else f"J{row['job_id']}"
                c.text(tag, bx + 8, ly + 13, service_w - 16, 30, size=12, color=WHITE, bold=True, align="center")
            if is_stolen:
                stolen_positions[cluster] = (bx + service_w / 2, ly + 29)
    return stolen_positions


def slide_6(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[5])
    events = load_events()
    c.text("동일 1,000-job stream · 동일 자원 · 정책만 변경", 420, 230, 1080, 42, size=19, color=AMBER, bold=True, align="center")
    c.text("회색 큐 대기 · 어두운 데이터 준비 · 밝은 연산 · 청록 훔친 작업", 410, 275, 1100, 34, size=14, color=MUTED, align="center")
    left_x, right_x, panel_y, panel_w = 90, 1010, 315, 820
    timeline_panel(c, left_x, panel_y, panel_w, "S1 정적 로컬 큐", "S1", events)
    positions = timeline_panel(c, right_x, panel_y, panel_w, "S3 지역성 인식 Work Stealing", "S3", events)
    c.chevron(910, 415, 78, 76, CYAN)
    c.text("재배치", 892, 505, 116, 38, size=16, color=CYAN, bold=True, align="center")
    for number, cluster in enumerate((1, 2, 3), 1):
        px, py = positions[cluster]
        c.circle(px - 18, py - 72, 36, CYAN)
        c.text(str(number), px - 18, py - 70, 36, 32, size=16, color=BG, bold=True, align="center", valign="middle")
        c.line(px, py - 34, px, py - 5, CYAN, width=3)
    c.line(90, 760, 1830, 760, GRID, width=2)
    metric(c, 145, 800, "유휴 클러스터 활용", "3 → 0", CYAN, 430)
    metric(c, 745, 800, "p95 지연 감소 · S3 vs S1", "-18.12%", BLUE, 430)
    metric(c, 1345, 800, "원격 가중치 이동 발생", "1.45 MiB", TEAL, 430)


def median_rows(workload: str, policies: tuple[str, ...]) -> dict[str, dict[str, float]]:
    with (ROOT / "results" / "experiments" / "scheduler_controlled.csv").open(encoding="utf-8", newline="") as stream:
        rows = list(csv.DictReader(stream))
    output: dict[str, dict[str, float]] = {}
    for policy in policies:
        selected = [row for row in rows if row["subset"] == "scheduler" and row["workload"] == workload and row["scheduler"] == policy and row["process_repetition"] == "0" and row["clusters"] == "4" and row["channels"] == "4" and row["bundles"] == "4" and row["link_width_bits"] == "128" and row["service_overlap_mode"] == "full"]
        if len(selected) != 5:
            raise RuntimeError(f"expected five rows for {workload}/{policy}")
        fields = ("p95_tile_latency_cycles", "p99_tile_latency_cycles", "queue_wait_mean_cycles", "remote_weight_bytes", "total_completion_cycles")
        output[policy] = {field: sorted(float(row[field]) for row in selected)[2] for field in fields}
    return output


def gemma_replay_rows() -> dict[str, dict[str, float]]:
    with (ROOT / "experiments" / "gemma3_1b" / "scheduler_replay.csv").open(encoding="utf-8", newline="") as stream:
        rows = list(csv.DictReader(stream))
    fields = ("p95_tile_latency_cycles", "p99_tile_latency_cycles")
    output: dict[str, dict[str, float]] = {}
    for policy in ("S1", "S3"):
        row = next(row for row in rows if row["decode_tokens"] == "32" and row["scheduler"] == policy)
        output[policy] = {field: float(row[field]) for field in fields}
    return output


def slide_7(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[6])
    data = gemma_replay_rows()
    groups = [("p95", "p95_tile_latency_cycles", 440), ("p99", "p99_tile_latency_cycles", 1150)]
    baseline_max = 1_500_000_000
    for label, field, cx in groups:
        c.text(label, cx - 220, 300, 440, 46, size=24, color=WHITE, bold=True, align="center")
        for index, policy in enumerate(("S1", "S3")):
            value = data[policy][field]
            bar_h = value / baseline_max * 470
            x = cx - 145 + index * 180
            y = 810 - bar_h
            color = GRAY if policy == "S1" else CYAN
            c.rect(x, y, 120, bar_h, color, radius=12)
            c.text(policy, x, 815, 120, 40, size=20, color=color, bold=True, align="center")
        reduction = (data["S1"][field] - data["S3"][field]) / data["S1"][field] * 100
        c.text(f"-{reduction:.2f}%", cx - 215, 870, 430, 64, size=36, color=CYAN, bold=True, align="center")
    c.line(950, 300, 950, 990, GRID, width=2)
    c.text("실제 Gemma projection ledger · 5,856 jobs · 분석 replay", 430, 235, 1060, 36, size=18, color=AMBER, bold=True, align="center")
    c.rect(570, 970, 780, 55, BG2, stroke=TEAL, stroke_width=2, radius=18)
    c.text("별도 skew stress: p95  -18.12%", 590, 974, 740, 44, size=18, color=TEAL, bold=True, align="center")


def slide_8(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[7])
    c.text("합성 skew · 1,000 jobs · 5-seed 중앙값 · 분석 모델", 340, 252, 1240, 44, size=20, color=AMBER, bold=True, align="center")
    stages = [
        ("1", "Tail 감소", "-18.12%", "p95 · S3 vs S1", BLUE),
        ("2", "원격 이동 최적화", "-37.84%", "가중치 · S3 vs S2", CYAN),
        ("3", "완료시간 절충", "+0.98%", "완료시간 · S3 vs S2", TEAL),
    ]
    xs = [150, 730, 1310]
    for index, (number, title, value, sub, color) in enumerate(stages):
        c.circle(xs[index], 345, 74, color)
        c.text(number, xs[index], 355, 74, 52, size=24, color=BG, bold=True, align="center", valign="middle")
        c.text(title, xs[index] - 20, 455, 390, 52, size=24, color=WHITE, bold=True)
        c.text(value, xs[index] - 20, 535, 390, 92, size=44, color=color, bold=True)
        c.text(sub, xs[index] - 20, 635, 390, 46, size=18, color=MUTED)
        if index < 2:
            c.chevron(xs[index] + 385, 520, 125, 64, MUTED)
    c.line(160, 760, 1750, 760, GRID, width=2)
    c.text("큐 감소", 200, 815, 310, 48, size=22, color=BLUE, bold=True, align="center")
    c.chevron(520, 820, 160, 48, BLUE)
    c.text("원격 가중치 이동", 690, 815, 410, 48, size=22, color=CYAN, bold=True, align="center")
    c.chevron(1110, 820, 160, 48, CYAN)
    c.text("링크·메모리 부담", 1280, 815, 420, 48, size=22, color=TEAL, bold=True, align="center")
    c.text("Sequential 경계에서는 skew의 S2/S3 p95 순위 역전", 260, 930, 1400, 46, size=20, color=AMBER, bold=True, align="center")


def prepare_kicad_crops() -> tuple[Path, Path, Path]:
    source = Image.open(ROOT / "paper" / "final" / "figures" / "paper_f07_kicad_coupon_render.png").convert("RGB")
    link = ASSETS / "kicad_link_crop.png"
    routed = ASSETS / "kicad_routed_crop.png"
    source.crop((420, 100, 1380, 610)).save(link)
    source.crop((260, 360, 1050, 850)).save(routed)

    top_render = ASSETS / "kicad_top_render.png"
    if not top_render.exists():
        kicad_cli = shutil.which("kicad-cli")
        if not kicad_cli:
            raise RuntimeError("kicad-cli is required to render the native refclk top view")
        subprocess.run(
            [
                kicad_cli,
                "pcb",
                "render",
                "--output",
                str(top_render),
                "--width",
                "1920",
                "--height",
                "1080",
                "--side",
                "top",
                "--background",
                "opaque",
                "--quality",
                "high",
                "--rotate",
                "0,0,0",
                str(ROOT / "hardware" / "kicad" / "k26_memory_coupon" / "k26_memory_coupon.kicad_pcb"),
            ],
            check=True,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
    top = Image.open(top_render).convert("RGB")
    refclk = ASSETS / "kicad_refclk_crop.png"
    top.crop((550, 220, 780, 340)).save(refclk)
    return link, refclk, routed


def slide_9(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[8])
    render = ROOT / "paper" / "final" / "figures" / "paper_f07_kicad_coupon_render.png"
    crops = prepare_kicad_crops()
    c.image_file(render, 80, 235, 1210, 680, mode="cover")
    c.rect(80, 235, 1210, 680, BG, stroke=TEAL, stroke_width=2, radius=10, alpha=0)
    labels = ["K26–Memory FPGA 링크", "기준 클록 차동쌍", "대표 라우팅 쿠폰"]
    for index, (path, label) in enumerate(zip(crops, labels)):
        y = 245 + index * 220
        c.image_file(path, 1370, y, 440, 150, mode="cover")
        c.rect(1370, y, 440, 150, BG, stroke=CYAN if index == 0 else TEAL, stroke_width=2, radius=8, alpha=0)
        if index == 1:
            c.line(1488, y + 66, 1680, y + 66, CYAN, width=3)
            c.line(1488, y + 78, 1680, y + 78, AMBER, width=3)
            c.text("N", 1460, y + 36, 30, 30, size=13, color=CYAN, bold=True, align="center")
            c.text("P", 1460, y + 86, 30, 30, size=13, color=AMBER, bold=True, align="center")
        c.text(label, 1370, y + 156, 440, 38, size=17, color=WHITE, bold=True, align="center")
        c.line(1290, y + 75, 1360, y + 75, CYAN, width=2)
    badges = [("풋프린트 29개", BLUE, 300), ("GTH/refclk 배선 20개", TEAL, 390), ("ERC 0 · 부분 DRC 0", CYAN, 420)]
    x = 120
    for label, color, width in badges:
        c.rect(x, 930, width, 54, BG2, stroke=color, stroke_width=2, radius=20)
        c.text(label, x + 12, 938, width - 24, 38, size=17, color=color, bold=True, align="center")
        x += width + 24
    c.text("NOT FOR FABRICATION · 참조 쿠폰", 1330, 942, 500, 36, size=13, color=RED, bold=True, align="center")


def slide_10(c: SlideCanvas) -> None:
    add_header(c, SLIDE_META[9])
    items = [
        ("실제 Gemma\n작업 부하", "그래프 → TileJob", BLUE),
        ("지역성 인식\n스케줄링", "Tail + 이동 비용", CYAN),
        ("물리 참조 설계", "RTL → KiCad", TEAL),
    ]
    xs = [120, 710, 1300]
    for index, (title, sub, color) in enumerate(items):
        c.circle(xs[index] + 90, 390, 190, BG2, stroke=color, stroke_width=3)
        c.text(str(index + 1), xs[index] + 140, 430, 90, 70, size=38, color=color, bold=True, align="center")
        c.text(title, xs[index], 610, 370, 88, size=21, color=WHITE, bold=True, align="center", valign="middle")
        c.text(sub, xs[index], 715, 370, 42, size=17, color=MUTED, align="center")
        if index < 2:
            c.chevron(xs[index] + 400, 455, 120, 60, color)
    c.line(170, 800, 1750, 800, GRID, width=2)
    c.text("다음 검증  |  DDR 응답 → Link 수신 → Weight FIFO → MatVec", 230, 840, 1460, 62, size=24, color=WHITE, bold=True, align="center")
    c.text("현재: analytical/RTL-bounded · 보드 성능·전력은 미측정", 450, 930, 1020, 46, size=17, color=MUTED, align="center")


SLIDE_BUILDERS = (slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7, slide_8, slide_9, slide_10)


def write_notes() -> None:
    lines = ["# VARP K26–Memory FPGA — 10분 발표 Speaker Notes", "", "총 목표 시간: 10:00", ""]
    for index, meta in enumerate(SLIDE_META, 1):
        lines += [f"## Slide {index}: {meta.title} ({meta.duration})", "", meta.notes, ""]
    (FINAL / "speaker_notes.md").write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def write_source_index() -> None:
    rows = [
        (1, "없음", "assets/cover_background.png", "image_gen + build_editable_deck.py", "두 연산 모듈 사이의 추상적 데이터 이동을 장식적으로 표현", "실제 칩·회로·구현 완료 상태"),
        (2, "results/experiments/scheduler_controlled.csv", "편집 가능한 queue/tail 도형", "build_editable_deck.py", "정적 ownership에서 skew와 idle이 생길 수 있음", "실제 보드 p95 또는 물리 queue timing"),
        (3, "없음", "paper_f01_evidence_path.svg; RTL sources", "build_editable_deck.py", "TileScheduler/payload store/compute cluster와 독립 memory/link plane 구조", "닫힌 DDR→link→MatVec payload loop"),
        (4, "없음", "work_stealing_sequence.mp4/.gif; work_stealing_storyboard.png; scheduler source", "manim_scenes.py; generate_animations.py; build_editable_deck.py", "S3의 victim search·locality score·exact-once 정책 개념", "분석 모델의 all-eligible search와 RTL victim-head 검사, ownership, scoring을 알고리즘적으로 동일시하거나 새 알고리즘 최적성·보드 성능으로 해석"),
        (5, "experiments/gemma3_1b/projection_trace.csv; experiments/gemma3_1b/scheduler_replay.csv", "편집 가능한 7,837→183→5,856 pipeline", "generate_trace.py; build_editable_deck.py", "실제 graph inventory와 graph-derived decode-32 replay; S1 대비 S3 p95 -15.07%, p99 -14.61%", "전체 모델 RTL 실행, ORT timing, 또는 다음 synthetic stress와 동일한 ledger"),
        (6, "results/experiments/scheduler_controlled.csv; assets/s1_s3_timeline_events.csv", "편집 가능한 41k–43k 확대 타임라인; scheduler_timeline.mp4/.gif", "generate_conference_figures.py; manim_scenes.py; build_editable_deck.py", "Gemma replay와 별개의 synthetic skew 1,000-job seed 23 full-overlap analytical event에서 S1 idle, S3 J172/J162/J143 이동, queue/data/compute 구간", "Gemma replay event, RTL cycle, 물리 link/DDR timing, 또는 compute가 dispatch 직후 시작했다는 해석"),
        (7, "experiments/gemma3_1b/scheduler_replay.csv; results/experiments/scheduler_controlled.csv", "편집 가능한 Gemma p95/p99 막대; tail_latency_results.mp4/.gif", "manim_scenes.py; build_editable_deck.py", "실제 graph-derived decode-32 replay의 S1 대비 S3 p95 -15.07%, p99 -14.61%; 별도 synthetic skew p95 -18.12%는 조건 일치 보조값", "두 data layer를 동일 실험으로 합치거나 측정 K26 latency·보편적 우월성으로 해석"),
        (8, "results/experiments/scheduler_controlled.csv; assets/bottleneck_shift_source.csv", "편집 가능한 3단계 trade-off; bottleneck_migration.mp4/.gif", "generate_conference_figures.py; manim_scenes.py; build_editable_deck.py", "같은 synthetic skew 5-seed full-overlap 중앙값에서 S1 대비 p95, S2 대비 remote/completion 절충", "p95 -18.12%를 queue-wait -17.86%와 동일시하거나 서로 다른 baseline 수치를 하나의 동일 비교로 해석"),
        (9, "없음", "paper_f07_kicad_coupon_render.png; assets/kicad_top_render.png; k26_scope_manifest.json; 실제 render crops", "kicad-cli pcb render; verify_k26_kicad.py; build_editable_deck.py", "Native KiCad reference coupon, board 좌표 기반 GTH_REFCLK0 P/N 두 segment 강조, coupon ERC 0과 routed-subset DRC 0 범위", "제작 가능 보드, 전체-board DRC 0, SI/PI/PDN/thermal closure"),
        (10, "없음", "docs/architecture.md; docs/evidence.md", "build_editable_deck.py", "기여·현재 한계·다음 integration gate 요약", "완성 accelerator, 보드 성능·전력, full 3B 실행"),
    ]
    lines = ["# Slide Source Index", "", "모든 경로는 repository root 기준이다. 화면에서 생략한 seed·모델·범위 조건은 이 인덱스와 speaker notes에서 유지한다.", "", "| Slide | 원본 CSV | 사용 Figure / 자료 | 생성·검증 스크립트 | 허용 해석 | 금지 해석 |", "|---:|---|---|---|---|---|"]
    for row in rows:
        lines.append("| " + " | ".join(str(value).replace("|", "\\|") for value in row) + " |")
    (FINAL / "slide_source_index.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    (FINAL / "qa_evidence_index.md").write_text("# Presentation QA Evidence Index\n\n슬라이드별 claim/source/허용·금지 해석은 [`slide_source_index.md`](slide_source_index.md)를 따른다. 객체 경계·가로/세로 overflow·텍스트 박스 교차 검사는 [`layout_audit.json`](layout_audit.json), 세 역할의 최종 검토는 [`independent_review.md`](independent_review.md), 발표 질의응답은 [`../../study/10_qna_bank.md`](../../study/10_qna_bank.md)를 사용한다.\n", encoding="utf-8")
    (FINAL / "demo_script.md").write_text("# 10-minute run of show\n\n" + "\n".join(f"- {meta.duration} · Slide {index}: {meta.title}" for index, meta in enumerate(SLIDE_META, 1)) + "\n\n전체 발표 대본은 [`speaker_notes.md`](speaker_notes.md)에 있다.\n", encoding="utf-8")


def build_contact_sheet(paths: list[Path]) -> Path:
    thumb_w, thumb_h = 576, 324
    sheet = Image.new("RGB", (thumb_w * 2 + 72, thumb_h * 5 + 144), rgb("#030812"))
    draw = ImageDraw.Draw(sheet)
    draw.text((36, 24), "VARP K26–Memory FPGA · 10-slide contact sheet", font=font(18, True), fill=rgb(WHITE))
    for index, path in enumerate(paths):
        image = Image.open(path).convert("RGB").resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        col, row = index % 2, index // 2
        x, y = 36 + col * thumb_w, 90 + row * thumb_h
        sheet.paste(image, (x, y))
        draw.rounded_rectangle((x + thumb_w - 58, y + 10, x + thumb_w - 10, y + 45), radius=8, fill=rgb("#030812"), outline=rgb(CYAN), width=2)
        draw.text((x + thumb_w - 34, y + 27), f"{index + 1:02d}", font=font(9, True), fill=rgb(CYAN), anchor="mm")
    output = FINAL / "slide_contact_sheet.png"
    sheet.save(output)
    return output


def build_pdf(paths: list[Path]) -> None:
    page_size = (960, 540)
    pdf = pdf_canvas.Canvas(str(PDF_PATH), pagesize=page_size, pageCompression=1)
    pdf.setTitle("VARP K26–Memory FPGA — 정적 큐의 Tail을 줄이는 Work Stealing")
    pdf.setAuthor("CHOI YUNHYUK")
    pdf.setCreator("presentation/tools/build_editable_deck.py")
    pdf.setDateFormatter(lambda *_: "D:20260731000000+09'00'")
    for path in paths:
        pdf.drawImage(str(path), 0, 0, width=page_size[0], height=page_size[1])
        pdf.showPage()
    pdf.save()


def normalize_pptx(path: Path) -> None:
    """Rewrite the OPC package with stable order and ZIP timestamps."""
    normalized = path.with_suffix(".normalized.pptx")
    fixed_time = (2026, 7, 31, 0, 0, 0)
    with zipfile.ZipFile(path, "r") as source, zipfile.ZipFile(
        normalized, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9
    ) as target:
        for name in sorted(source.namelist()):
            original = source.getinfo(name)
            info = zipfile.ZipInfo(name, fixed_time)
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = original.external_attr
            info.create_system = original.create_system
            target.writestr(info, source.read(name))
    normalized.replace(path)


def main() -> int:
    SLIDES.mkdir(parents=True, exist_ok=True)
    CACHE.mkdir(parents=True, exist_ok=True)
    for stale in SLIDES.glob("slide_*.png"):
        stale.unlink()
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "정적 큐의 Tail을 줄이는 Work Stealing"
    prs.core_properties.author = "CHOI YUNHYUK"
    prs.core_properties.creator = "presentation/tools/build_editable_deck.py"
    prs.core_properties.created = datetime(2026, 7, 31, 0, 0, 0)
    prs.core_properties.modified = datetime(2026, 7, 31, 0, 0, 0)
    paths: list[Path] = []
    for index, (meta, builder) in enumerate(zip(SLIDE_META, SLIDE_BUILDERS), 1):
        canvas = SlideCanvas(prs, meta.title, meta.notes)
        builder(canvas)
        paths.append(canvas.save(index))
    prs.save(PPTX_PATH)
    normalize_pptx(PPTX_PATH)
    write_notes()
    write_source_index()
    build_pdf(paths)
    contact = build_contact_sheet(paths)
    print(f"pptx={PPTX_PATH} slides={len(prs.slides)}")
    print(f"pdf={PDF_PATH}")
    print(f"contact_sheet={contact}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
