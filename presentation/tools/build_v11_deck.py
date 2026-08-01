#!/usr/bin/env python3
"""Build the 17-slide v11 Korean final-round conference deck.

The presentation keeps text, diagrams, bars, arrows, and annotations editable.
Raster media is limited to a decorative cover, Manim fallback frames, the
native KiCad render/crops, and the repository QR code.
"""

from __future__ import annotations

import csv
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
import shutil
import subprocess

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from reportlab import rl_config
from reportlab.graphics import renderSVG
from reportlab.graphics.barcode.qr import QrCodeWidget
from reportlab.graphics.shapes import Drawing
from reportlab.pdfgen import canvas as pdf_canvas

from build_editable_deck import (
    AMBER,
    BG,
    BG2,
    BG3,
    BLUE,
    CYAN,
    GRID,
    H,
    MUTED,
    RED,
    SLIDES,
    TEAL,
    W,
    WHITE,
    SlideCanvas,
    font,
    inch,
    normalize_pptx,
    ppt_rgb,
    rgb,
)


ROOT = Path(__file__).resolve().parents[2]
FINAL = ROOT / "presentation/final"
ASSETS = FINAL / "assets"
PPTX = FINAL / "presentation.pptx"
PDF = FINAL / "presentation.pdf"
REPO_URL = "https://github.com/snowman0919/varp-k26-memory-fpga"
SECTIONS = ("문제", "구조", "평가", "결론")


@dataclass(frozen=True)
class SlideMeta:
    title: str
    section: int
    conclusion: str
    duration: int
    memory: str
    speech: str
    transition: str
    supplement: str


META = (
    SlideMeta("온디바이스 sLLM용 K26–Memory FPGA 가속기 설계", 0, "실제 Gemma 작업에서 계산·메모리 분리의 채택 조건을 분석했다.", 10, "외부 메모리는 무조건 추가하는 부품이 아니라 채택 조건을 검증할 확장 수단입니다.", "저는 외부 메모리를 무조건 추가하는 것이 아니라, 언제 외부 메모리가 실제로 이득이 되는지를 확인하려고 했습니다.", "오늘 답할 네 가지 질문부터 보겠습니다.", "표지의 배경은 개념 이미지이며 실제 구현 보드가 아닙니다."),
    SlideMeta("오늘 답할 네 가지 질문", 0, "문제 재정의에서 구조·검증·채택 조건까지 직선적으로 답한다.", 20, "발표는 문제, 구조, 실제 모델 검증, 외부 메모리 채택 조건의 순서입니다.", "첫째 어떤 문제를 발견했는지, 둘째 어떤 구조로 풀었는지, 셋째 실제 모델에서 무엇이 달라졌는지, 넷째 외부 메모리는 언제 필요한지를 답하겠습니다.", "먼저 출발점이었던 용량 가설의 반례를 보겠습니다.", "목차는 결과가 아니라 발표의 논리 순서입니다."),
    SlideMeta("1B 모델은 4GB에 들어갔다—그래서 질문을 바꿨다", 0, "외부 메모리는 용량보다 공급 분리의 이득이 이동 비용을 넘을 때 의미가 있다.", 45, "Gemma 3 1B는 용량 모델상 K26 4GB에 들어가므로 외부 8GB는 필수가 아닙니다.", "처음에는 외부 8GB가 단순히 모델을 담기 위해 필요하다고 생각했습니다. 그런데 실제 용량을 계산해 보니 Gemma 3 1B는 K26의 4GB 안에도 들어갔습니다. 그래서 질문을 바꿨습니다. 외부 메모리는 더 큰 창고이기만 해서는 의미가 없고, 가중치 공급을 분리해 얻는 이득이 운반 비용보다 커야 합니다. 외부 8GB는 3B 이하 확장과 공급 경로 분리를 평가하기 위한 구조입니다.", "이 질문을 두 장치의 역할 분리로 구체화했습니다.", "2.43 GiB는 INT8·문맥 32K·실행 여유의 용량 모델이며 보드 실측이 아닙니다."),
    SlideMeta("계산과 가중치 공급을 두 장치로 분리했다", 1, "K26은 연산·제어를, Memory FPGA는 확장 메모리와 가중치 공급을 맡는다.", 50, "K26의 계산 위치와 Memory FPGA의 데이터 위치를 연산 타일에서 함께 추적합니다.", "K26은 운영체제와 제어, 네 연산 클러스터, 로컬 대기열, MatVec를 맡습니다. Memory FPGA는 가중치 저장과 네 메모리 채널, 채널 선택, 링크 전송을 맡습니다. 연산 타일, 즉 TileJob은 계산할 클러스터와 가중치가 가까운 채널을 함께 기록합니다. 실선은 현재 RTL이나 모델에 존재하는 경로이고, 점선은 아직 남은 실제 직렬 링크와 메모리 제어 경로입니다.", "이 구조가 기존 접근과 어떻게 다른지 비교하겠습니다.", "논리 경로 일부는 RTL에 있으나 실제 GTH·CDC·MIG 물리 경로는 닫히지 않았습니다."),
    SlideMeta("메모리 용량만 늘리지 않고 작업 배치와 데이터 위치를 함께 설계했다", 1, "새 스케줄러보다 외부 메모리 가속기가 유효해지는 시스템 조건을 찾는 것이 핵심이다.", 35, "독창성은 Work Stealing 자체가 아니라 실제 모델·초기 배치·이동 비용·물리 인터페이스를 한 설계 흐름에서 연결한 점입니다.", "기존 연구는 연산량이나 대역폭, 정적 배정, 합성 행렬, 또는 개별 RTL 구현에 집중하는 경우가 많았습니다. 이 연구는 연산·메모리·링크의 균형을 함께 보고, 이동 비용을 고려해 작업을 선택적으로 재분배하며, 실제 Gemma 그래프에서 작업을 추출해 모델에서 RTL과 PCB 참조 설계까지 연결했습니다. Work Stealing은 기존 기술이며, 여기서는 구조 채택 조건을 찾는 도구입니다.", "이제 직접 검증한 MatVec 범위를 보겠습니다.", "FTRANS·DFX·FlightLLM·고전적 Work Stealing·LAWS와 논문에 적힌 비교 항목만 사용합니다."),
    SlideMeta("실제 Gemma 가중치를 MatVec RTL에서 검증했다", 1, "세 대표 가중치 타일이 소프트웨어 INT32 기준값과 정확히 일치했다.", 35, "직접 증명한 것은 실제 Gemma 가중치 세 타일의 16×4 INT8 MatVec 산술 일치입니다.", "실제 Gemma 가중치를 추출해 시험 입력을 만들고, 16×4 INT8 MatVec RTL에 넣은 뒤 소프트웨어 INT32 기준값과 비교했습니다. gate projection, lm_head, output projection 세 타일이 모두 정확히 일치했습니다. 이 시험은 MatVec 산술을 확인한 범위이며 실제 DDR 읽기나 FPGA 간 링크 성능은 포함하지 않습니다.", "정확한 산술만으로는 정적 배정의 유휴 문제를 해결할 수 없습니다.", "응답 데이터는 시험 경계에서 주입했으며 전체 Gemma 추론·실제 DDR·물리 링크는 검증하지 않았습니다."),
    SlideMeta("정적 배정은 데이터를 가까이 두지만 작업을 고르게 나누지 못한다", 1, "한 큐에 작업이 몰리면 다른 연산기가 비어 있어도 마지막 작업은 기다린다.", 30, "정적 로컬 큐는 가중치 이동을 줄이지만 작업 쏠림을 스스로 고치지 못합니다.", "정적 로컬 큐, S1은 각 작업대 앞에 고정된 대기 줄을 둔 것과 같습니다. 가중치를 가까이 둘 수 있지만 C0에 일이 몰리면 C2와 C3가 비어 있어도 돕지 못합니다. 그래서 평균보다 마지막 작업과 p95·p99 꼬리 지연이 길어질 수 있습니다.", "빈 작업대가 이득이 있는 일만 가져오도록 바꿨습니다.", "화면의 큐는 개념도이며 성능 수치는 뒤의 분석 모델에서 제시합니다."),
    SlideMeta("빈 연산기가 이득이 있는 작업만 가져온다", 1, "기다림 감소 이득과 가중치·활성값 이동 비용을 함께 계산한다.", 40, "지역성 인식 작업 재분배는 대기 감소가 운반 비용보다 클 때만 동작합니다.", "지역성 인식 작업 재분배, S3는 빈 작업대가 옆 줄의 일을 가져오는 방식입니다. 먼저 유휴 연산기를 찾고, 가져올 수 있는 작업의 기다림 감소와 가중치·활성값 이동 비용을 비교합니다. 이득이 더 클 때만 소유권을 옮깁니다. 따라서 Work Stealing은 연구의 목적이 아니라 제한된 자원을 쓰기 위한 선택 수단입니다.", "이 정책을 실제 Gemma 작업 단위로 바꾸겠습니다.", "정책 영상은 상태 전이를 설명하며 최적성 증명이나 실제 보드 타이밍이 아닙니다."),
    SlideMeta("Gemma 그래프를 802개 연산 타일로 나누고 단계 의존 규칙을 부여했다", 2, "큰 행렬을 출력 방향으로 분할하고 보수적 단계 의존 규칙을 부여했다.", 35, "실제 모델의 7,837개 노드에서 183개 투영을 찾고 802개 연산 타일로 분할했습니다.", "실제 Gemma ONNX 그래프의 7,837개 노드에서 토큰당 183개 행렬 투영 연산을 찾았습니다. 큰 행렬 하나를 여러 출력 타일로 나누면 802개 연산 타일이 됩니다. 질의·키·값 다음 출력 투영, 게이트·상향 다음 하향, 그다음 계층이라는 보수적 단계 의존 규칙을 부여했습니다. 32토큰에서는 25,664개입니다.", "이제 같은 작업에서 정책만 바꾼 대표 사례를 보겠습니다.", "802와 25,664는 고정한 그래프 형상과 명시한 타일 분할·보수적 단계 규칙에서 나온 모델링 수치입니다."),
    SlideMeta("같은 작업에서 정적 배정과 작업 재분배만 비교했다", 2, "동일 작업·초기값·하드웨어 가정에서 정책만 바꿔 대기 감소와 이동 비용을 분리했다.", 45, "대표 연산 타일 하나에서 긴 큐 대기가 원격 준비를 포함한 짧은 경로로 바뀝니다.", "합성 불균형 부하에서 작업, 초기값, 자원 가정을 고정하고 정책만 바꿨습니다. 대표 연산 타일은 정적 배정에서 C0 대기열에 오래 남지만, 재분배에서는 유휴 C3로 옮겨집니다. 큐 대기는 약 6만 7천 사이클 줄고, 대신 원격 가중치 준비 구간이 추가됩니다. 모든 정책에서 작업 ID의 중복과 누락이 없는지도 확인했습니다.", "이 변화가 전체 분포에서는 언제 효과적인지 보겠습니다.", "원본 CSV의 분석 사건이며 cycle은 실제 DDR·GTH 물리 사이클이 아닙니다."),
    SlideMeta("작업이 치우친 경우에만 재분배가 효과적이었다", 2, "초기 배치가 이미 균형이면 이동 비용만 추가될 수 있다.", 45, "재분배 효과는 부하 불균형과 초기 데이터 배치에 조건부입니다.", "합성 불균형 부하에서는 정적 배정보다 p95가 19.13%, p99가 18.71% 줄었습니다. 하지만 Gemma 형상에서는 기존 초기 배치가 p95를 0.28% 늘렸고, 메모리 채널을 고려한 초기 배치에서는 19.79% 줄였습니다. 즉 스케줄링만 따로 떼어 항상 좋다고 말할 수 없고, 처음 데이터를 어디에 놓았는지와 함께 봐야 합니다.", "꼬리 지연을 줄인 대가가 무엇인지 이어서 보겠습니다.", "p95·p99는 의존성 해제 뒤 연산 타일 완료 분포이며 사용자 요청 지연이 아닙니다."),
    SlideMeta("유휴 연산은 줄었지만 데이터 이동이 새로운 비용이 됐다", 2, "작업 재분배의 이득은 링크·메모리 비용과 함께 판단해야 한다.", 40, "지역성 점수는 무제약 재분배보다 원격 가중치를 줄이지만 완료시간은 조금 늘 수 있습니다.", "정적 배정과 비교하면 지역성 인식 재분배의 p95는 19.13% 줄었습니다. 반면 무제약 재분배와 비교하면 원격 가중치는 35.49% 줄었지만 마지막 완료시간은 0.01% 늘었습니다. 앞의 수치는 정적 배정 대비 효과이고 뒤의 두 수치는 재분배 방식끼리의 비교입니다. 서로 다른 기준이므로 하나의 인과효과처럼 연결하지 않습니다.", "이 비용까지 포함해 외부 메모리 채택 조건을 보겠습니다.", "분석 모델의 비교 기준이 다르며 서비스 중첩 가정에 따라 정책 순위가 바뀔 수 있습니다."),
    SlideMeta("현재 모델 가정에서는 K26 로컬 경로가 유리했다", 2, "외부 메모리는 용량·대역폭 격리의 이득이 링크 비용을 넘을 때만 채택한다.", 50, "시험한 민감도 범위에서는 외부 메모리 곡선이 K26 로컬 기준선과 교차하지 않았습니다.", "민감도 CSV의 유효 공급 대역폭과 예상 완료 사이클을 같은 축에 놓았습니다. K26 로컬은 4.8에서 14.4 GB/s, 외부 4채널은 3.2에서 12.8 GB/s를 시험했습니다. 이 범위에서는 외부 곡선이 로컬 기준선과 교차하지 않았고 중앙 조건도 24.6M 대 120.6M사이클이었습니다. 더 큰 모델, 긴 문맥, 실제 로컬 메모리 경합이 생기면 같은 방식으로 다시 검증해야 합니다.", "다음은 이 구조를 물리 인터페이스로 어디까지 내렸는지입니다.", "모든 GB/s와 cycle은 분석 민감도이며 실제 하드웨어 성능 측정값이 아닙니다."),
    SlideMeta("핵심 인터페이스를 실제 PCB 객체로 구체화했다", 2, "K26–Memory FPGA 링크와 DDR3L x16 대표 슬라이스를 참조 라우팅으로 검증했다.", 30, "KiCad 결과는 제작 보드가 아니라 다음 물리 검증 범위를 고정한 참조 쿠폰입니다.", "실제 KiCad 객체로 범용 경계 링크, 기준 클록 차동쌍, DDR3L x16 대표 슬라이스를 배치하고 일부를 라우팅했습니다. 풋프린트 29개, 고속·클록 배선 20개이며 선언한 제한 범위의 ERC와 DRC는 0입니다. 그러나 전체 116개 중 55개가 미배선이고 전원·신호 무결성도 남아 있어 제작용이 아닙니다.", "남은 과제를 다음 물리 검증 순서로 정리하겠습니다.", "실제 K26 SOM 핀·FPGA BGA ball·네 MIG 채널·최종 메모리 구성·전원 트리·SI/PI는 포함하지 않습니다."),
    SlideMeta("다음 단계는 분석 모델을 실제 하드웨어 경로로 닫는 것이다", 3, "K26 기준선부터 물리 경로와 전력·신호 무결성까지 순서대로 검증한다.", 30, "현재의 제한 모델을 K26 실측, 물리 경로 통합, SI/PI·전력·열 검증으로 닫습니다.", "현재 결과는 제한된 RTL과 분석 모델입니다. 먼저 같은 Gemma 작업으로 K26 로컬 기준선을 측정하고, 다음으로 DDR에서 링크를 거쳐 MatVec까지 물리 경로를 통합합니다. 마지막으로 SI/PI, 전력, 열을 검증합니다. 정확성, 대역폭, 꼬리 지연, 전력과 신호 무결성의 순서로 닫겠습니다.", "이제 연구의 결론과 기여를 세 문장으로 정리하겠습니다.", "GTH·CDC·크레딧·패킷·MIG와 보드 계측은 후속 구현 범위입니다."),
    SlideMeta("결론과 기여", 3, "구조 제안, 조건 분석, 채택 기준 도출을 하나의 설계 흐름으로 연결했다.", 35, "연산 배치와 메모리 위치를 함께 고려해 외부 메모리의 채택 조건을 구체화했습니다.", "K26 계산부와 외부 다채널 공급부를 분리한 구조를 설계했습니다. Gemma 그래프 기반 분석에서는 불균형 이득이 이동 비용보다 클 때만 재분배가 유효했습니다. 이번 민감도에서는 1B가 K26 로컬 우선이었습니다. 따라서 외부 메모리는 더 큰 모델, 긴 문맥, 실제 경합에서 재평가해야 합니다.", "이제 질문과 토론으로 넘어가겠습니다.", "구조와 분석 결과이며 완성 제품·보드 실측·보편적 최적 정책을 주장하지 않습니다."),
    SlideMeta("질문과 토론", 3, "감사합니다.", 0, "질문에는 결론→근거→경계→다음 검증 순으로 답합니다.", "질문 받겠습니다.", "감사합니다.", "Q&A는 발표 시간 9분 35초에 포함하지 않으며 상세 답변은 발표자 참고 자료에만 둡니다."),
)

_note_extensions = ["" for _ in META]
_note_extensions[4] = "관련 연구 비교 기준은 FTRANS의 대형 모델 전송, DFX의 분산 FPGA 실행, FlightLLM의 메모리 중심 가속, 고전적 Work Stealing의 부하 균형, LAWS의 지역성 고려입니다. 화면은 문헌에 명시된 범위만 요약합니다."
NOTE_EXTENSIONS = tuple(_note_extensions)

VISUAL_ORDERS = (
    "제목 → 채택 조건 부제 → 개념 이미지 표시",
    "문제 → 구조 → 실제 모델 → 채택 조건",
    "처음 가설 → 2.43 GiB → 의사결정식 → 새 질문",
    "K26 역할 → 요청 경로 → Memory FPGA 역할 → 응답 경로 → 선 경계",
    "기존 접근 → 본 연구 → 하단 연구 핵심",
    "실제 가중치 → 시험 입력 → MatVec RTL → 기준값 비교 → 3/3",
    "C0 집중 → C2·C3 유휴 → 장점과 한계",
    "유휴 발견 → 이득·비용 비교 → 조건부 이전 → 부등식",
    "7,837 → 183 → 802 → 분할 이유 → 의존 순서",
    "동일 조건 → 정적 배정 → 지역성 인식 재분배 → 대표 cycle",
    "합성 불균형 → Gemma 기존 배치 → 채널 고려 배치 → 조건 문장",
    "정적 배정 대비 → 무제약 재분배 대비 → 기준 차이",
    "모델 입력 유효 메모리 대역폭 축 → K26 로컬 곡선 → 외부 곡선 → 미교차 구간",
    "전체 KiCad 렌더 → 경계 링크 → 기준 클록 → DDR3L 채널 → 배지",
    "현재 → 다음 1 → 다음 2 → 다음 3 → 검증 순서",
    "구조 → 효과 조건 → 현재 채택 판단 → 기여 → QR",
    "중앙 질문과 토론 → 감사 인사",
)

EXPECTED_QUESTIONS = (
    "왜 K26을 계산 장치로 선택했는가?",
    "왜 이 네 질문 순서로 발표하는가?",
    "1B가 4GB에 들어가는데 왜 외부 8GB인가?",
    "왜 K26과 Memory FPGA의 역할을 분리했는가?",
    "Work Stealing은 기존 기술인데 무엇이 새로운가?",
    "실제 FPGA 성능을 측정했는가?",
    "왜 중앙 큐가 아니라 로컬 큐인가?",
    "어떤 조건이면 작업을 옮기지 않는가?",
    "183개 투영과 802개 연산 타일은 무엇이 다른가?",
    "동일 조건을 어떻게 보장했는가?",
    "왜 어떤 조건에서는 재분배가 더 느린가?",
    "세 수치의 비교 기준이 왜 다른가?",
    "GB/s와 cycle은 실측인가?",
    "PCB는 제작 가능한가?",
    "가장 먼저 수행할 실제 보드 실험은 무엇인가?",
    "외부 Memory FPGA 결과가 나쁜데 연구가 실패한 것 아닌가?",
    "가장 강한 공격 질문에는 어떤 순서로 답할 것인가?",
)

ANSWER_CORES = (
    "K26은 온디바이스 제어와 FPGA 연산을 함께 제공하며 이번 연구의 로컬 4GB 기준선이기 때문입니다.",
    "문제에서 시작해 구조와 직접 검증을 거친 뒤 채택 판단으로 끝내야 설계 의사결정이 보입니다.",
    "1B 필수 용량이 아니라 공급 분리와 더 큰 모델·긴 문맥 확장 조건을 평가하기 위해서입니다.",
    "K26은 계산 소유권을, Memory FPGA는 데이터 위치와 다채널 공급을 담당하게 해 두 비용을 분리합니다.",
    "알고리즘 신규성이 아니라 실제 Gemma 형상·초기 배치·원격 비용·RTL·PCB를 하나의 채택 분석으로 연결한 점입니다.",
    "아닙니다. 실제 가중치 세 타일의 MatVec RTL 일치만 직접 확인했고 성능은 분석 모델입니다.",
    "로컬 큐는 가중치 지역성과 분산 제어를 유지하면서 정적 배정의 불균형을 관찰하기 위한 기준선입니다.",
    "대기 감소가 가중치·활성값 이동 비용보다 작거나 작업이 옮길 수 없으면 재분배하지 않습니다.",
    "183은 그래프 연산 수이고 802는 큰 행렬을 출력 방향으로 나눈 하드웨어 작업 수입니다.",
    "같은 작업 목록·초기값·자원 가정·완료 ID를 고정하고 정책만 바꿨습니다.",
    "초기 배치가 이미 균형이거나 메모리 채널이 병목이면 이동 비용만 추가될 수 있습니다.",
    "정적 배정 효과와 무제약 재분배 대비 지역성 효과를 분리하기 위한 서로 다른 기준입니다.",
    "아닙니다. 유효 대역폭과 예상 cycle을 바꾼 분석 민감도입니다.",
    "아닙니다. 실제 객체와 일부 배선만 있는 참조 쿠폰이며 미배선·SI/PI·전원 과제가 남았습니다.",
    "동일 Gemma 작업의 K26 로컬 메모리 기준선을 보드에서 먼저 측정합니다.",
    "실패가 아니라 현재 1B에서는 채택하지 말아야 한다는 조건을 도출한 결과입니다.",
    "결론 한 문장, 근거 유형과 수치, 해석 경계, 다음 검증 순으로 20초 안에 답합니다.",
)

CAUTIONS = (
    "개념 이미지를 실제 보드로 소개하지 않는다.",
    "목차를 결과처럼 설명하지 않는다.",
    "외부 8GB가 Gemma 1B의 필수 용량이라고 말하지 않는다.",
    "실선을 전체 물리 경로 구현으로 확대하지 않는다.",
    "Work Stealing 자체를 알고리즘 독창성으로 주장하지 않는다.",
    "3/3 일치를 전체 모델 실행·성능으로 확대하지 않는다.",
    "정적 로컬 큐가 모든 조건에서 나쁘다고 일반화하지 않는다.",
    "정책 점수를 최적성 증명이나 보드 타이밍으로 말하지 않는다.",
    "그래프 기반 작업 수를 실제 FPGA 실행 횟수로 부르지 않는다.",
    "분석 cycle을 물리 cycle로 부르지 않는다.",
    "연산 타일 p95를 사용자 요청 지연으로 부르지 않는다.",
    "서로 다른 기준 수치를 하나의 인과식으로 묶지 않는다.",
    "민감도 곡선을 보드 실측이나 미래 모든 조건으로 일반화하지 않는다.",
    "제한 ERC/DRC 0을 제작 준비 완료로 말하지 않는다.",
    "한계를 변명처럼 나열하지 말고 다음 검증으로 연결한다.",
    "현재 1B 채택 판단을 외부 메모리의 보편적 실패로 말하지 않는다.",
    "준비된 질문과 답변을 화면에 추가하지 않는다.",
)

EVIDENCE_BADGES = (
    None, None, "분석 모델", "RTL 검증", None, "RTL 검증", None, None,
    "실제 모델", "분석 모델", "분석 모델", "분석 모델", "분석 모델",
    "PCB 참조 설계", None, None, None,
)

SOURCE_FOOTERS = (
    "개념 이미지 · 실제 구현 보드 아님",
    "발표 구성 · 본선 평가 기준 대응",
    "results/model_level/k26_local_external_sensitivity.csv",
    "docs/architecture.md · ClosedLoopVirtualPrototypeTop.scala",
    "논문 II장 · 문헌 [6–8, 15–16]",
    "evidence/model/gemma3_1b_rtl_tile_parity.csv",
    "정적 로컬 큐 개념 · 수치는 Slide 11",
    "src/varp/k26_scheduler_model.py",
    "projection_trace.csv · gemma3_1b_dependency_manifest.json",
    "presentation/final/assets/s1_s3_timeline_events.csv",
    "paired_policy_effects.csv · gemma3_1b_policy_effects.csv",
    "results/experiments/paired_policy_effects.csv",
    "results/model_level/k26_local_external_sensitivity.csv",
    "hardware/kicad/k26_reports/k26_scope_manifest.json",
    "docs/architecture.md · hardware/kicad/controlled_review.md",
    "paper/final/submission_manuscript.md · research/v11_research_freeze.md",
    "",
)


def add_progress(c: SlideCanvas, section: int) -> None:
    x0, y, width, gap = 1160, 42, 140, 14
    for index, label in enumerate(SECTIONS):
        color = CYAN if index == section else "#294057"
        c.rect(x0 + index * (width + gap), y, width, 7, color)
        c.text(label, x0 + index * (width + gap), 8, width, 28, size=12, color=color if index == section else MUTED, bold=index == section, align="center")


def header(c: SlideCanvas, meta: SlideMeta) -> None:
    c.rect(82, 56, 58, 7, CYAN)
    title_size = 24 if len(meta.title) >= 29 else 27 if len(meta.title) >= 25 else 30
    c.text(meta.title, 82, 78, 1680, 82, size=title_size, bold=True)
    c.text(meta.conclusion, 84, 160, 1690, 62, size=20, color=MUTED)
    add_progress(c, meta.section)


def footer(c: SlideCanvas, index: int) -> None:
    source = SOURCE_FOOTERS[index]
    badge = EVIDENCE_BADGES[index]
    if source:
        c.text(f"출처 · {source}", 84, 1028, 1320, 28, size=10, color="#71869A")
    if badge:
        colors = {"실제 모델": BLUE, "RTL 검증": TEAL, "분석 모델": AMBER, "PCB 참조 설계": CYAN}
        color = colors[badge]
        c.rect(1600, 1018, 235, 34, BG2, stroke=color, stroke_width=1.5, radius=12)
        c.text(badge, 1610, 1023, 215, 23, size=11, color=color, bold=True, align="center")


def dashed_line(c: SlideCanvas, x1: float, y1: float, x2: float, y2: float, color: str, *, width: float = 2.0, parts: int = 12) -> None:
    for part in range(parts):
        if part % 2:
            continue
        start = part / parts
        end = min(1.0, (part + 1) / parts)
        c.line(x1 + (x2 - x1) * start, y1 + (y2 - y1) * start, x1 + (x2 - x1) * end, y1 + (y2 - y1) * end, color, width=width)


def arrow_between(c: SlideCanvas, x: int, y: int, w: int = 90, color: str = CYAN) -> None:
    c.chevron(x, y, w, 50, color)


def embedded_movie(c: SlideCanvas, stem: str, x: int, y: int, w: int, h: int) -> None:
    frame = ASSETS / f"{stem}_frame.png"
    movie = ASSETS / f"{stem}.mp4"
    c.image_file(frame, x, y, w, h, mode="contain")
    c.slide.shapes.add_movie(str(movie), inch(x), inch(y), inch(w), inch(h), poster_frame_image=str(frame), mime_type="video/mp4")
    c.rect(x, y, w, h, BG, stroke=CYAN, stroke_width=2, radius=10, alpha=0)
    c.circle(x + w - 90, y + h - 90, 58, CYAN)
    c.text("▶", x + w - 78, y + h - 83, 34, 34, size=20, color=BG, bold=True, align="center", valign="middle")


def metric(c: SlideCanvas, x: int, y: int, value: str, label: str, color: str, width: int = 430) -> None:
    c.text(value, x, y, width, 76, size=37, color=color, bold=True, align="center")
    c.text(label, x, y + 78, width, 44, size=17, color=MUTED, align="center")


def slide_1(c: SlideCanvas) -> None:
    c.image_file(ASSETS / "cover_background.png", 0, 0, W, H, mode="cover")
    c.rect(0, 0, 1250, H, BG, alpha=66)
    c.rect(92, 220, 78, 8, CYAN)
    c.text("온디바이스 sLLM용\nK26–Memory FPGA 가속기 설계", 92, 262, 1280, 220, size=40, bold=True)
    c.text("실제 Gemma 작업에서 계산·메모리 분리의 채택 조건 분석", 96, 540, 1320, 62, size=22, color=CYAN, bold=True)
    c.text("최윤혁 · 한국디지털미디어고등학교", 96, 935, 760, 38, size=17, color=MUTED)
    c.text("개념 이미지", 1640, 945, 190, 32, size=12, color=AMBER, bold=True, align="right")


def slide_2(c: SlideCanvas) -> None:
    header(c, META[1])
    labels = [("1", "어떤 문제를\n발견했나"), ("2", "어떤 구조로\n풀었나"), ("3", "실제 모델에서\n무엇이 달라졌나"), ("4", "외부 메모리는\n언제 필요한가")]
    xs = [150, 580, 1010, 1440]
    for index, (number, title) in enumerate(labels):
        c.circle(xs[index], 350, 104, CYAN if index == 0 else BG3, stroke=CYAN, stroke_width=3)
        c.text(number, xs[index], 370, 104, 62, size=32, color=BG if index == 0 else CYAN, bold=True, align="center", valign="middle")
        c.text(title, xs[index] - 115, 520, 340, 96, size=22, bold=True, align="center", valign="middle")
        if index < 3:
            arrow_between(c, xs[index] + 155, 380, 120, "#294057")
    c.text("문제 재정의  →  구조  →  검증  →  채택 조건", 360, 835, 1200, 64, size=27, color=CYAN, bold=True, align="center")


def slide_3(c: SlideCanvas) -> None:
    header(c, META[2])
    c.text("처음 가설", 125, 270, 480, 44, size=21, color=AMBER, bold=True, align="center")
    c.text("외부 8GB가\n1B 모델의 필수 용량", 120, 345, 490, 120, size=27, color=WHITE, bold=True, align="center", valign="middle")
    arrow_between(c, 700, 375, 145, AMBER)
    c.text("용량 계산", 980, 270, 500, 44, size=21, color=CYAN, bold=True, align="center")
    metric(c, 1010, 340, "2.43 GiB", "INT8 + 문맥 32K + 실행 여유", CYAN, 450)
    c.line(100, 535, 1820, 535, GRID, width=2)
    c.text("외부 메모리의 이득", 120, 605, 400, 48, size=23, color=AMBER, bold=True)
    c.text("=", 505, 608, 80, 46, size=29, color=WHITE, bold=True, align="center")
    c.text("가중치 공급 분리", 590, 590, 370, 76, size=23, color=CYAN, bold=True, align="center", valign="middle")
    c.text("+", 970, 606, 70, 50, size=29, color=WHITE, bold=True, align="center")
    c.text("유휴 연산 감소", 1045, 590, 340, 76, size=23, color=TEAL, bold=True, align="center", valign="middle")
    c.text("-", 1390, 606, 70, 50, size=29, color=WHITE, bold=True, align="center")
    c.text("원격 데이터 이동 비용", 1455, 590, 360, 76, size=22, color=AMBER, bold=True, align="center", valign="middle")
    c.rect(320, 765, 1280, 130, BG2, stroke=CYAN, stroke_width=2, radius=18)
    c.text("외부 메모리의 이득이 이동 비용을 넘는 조건은 무엇인가?", 370, 793, 1180, 72, size=26, color=CYAN, bold=True, align="center", valign="middle")


def slide_4(c: SlideCanvas) -> None:
    header(c, META[3])
    c.rect(80, 275, 575, 610, BG2, stroke=BLUE, stroke_width=3, radius=22)
    c.text("Kria K26", 120, 305, 495, 50, size=28, color=BLUE, bold=True, align="center")
    for index, label in enumerate(("운영체제·제어", "네 연산 클러스터", "로컬 대기열", "MatVec")):
        y = 405 + index * 100
        c.rect(150, y, 435, 62, BG3, stroke=BLUE if index != 2 else CYAN, stroke_width=2, radius=10)
        c.text(label, 170, y + 12, 395, 38, size=20, color=WHITE, bold=True, align="center")
    c.rect(1265, 275, 575, 610, BG2, stroke=TEAL, stroke_width=3, radius=22)
    c.text("Memory FPGA", 1305, 305, 495, 50, size=28, color=TEAL, bold=True, align="center")
    for index, label in enumerate(("가중치 저장", "네 메모리 채널", "채널 선택", "링크 전송")):
        y = 405 + index * 100
        c.rect(1335, y, 435, 62, BG3, stroke=TEAL if index != 3 else CYAN, stroke_width=2, radius=10)
        c.text(label, 1355, y + 12, 395, 38, size=20, color=WHITE, bold=True, align="center")
    c.text("연산 요청 · 가중치 요청  →", 720, 395, 475, 48, size=20, color=CYAN, bold=True, align="center")
    c.line(680, 460, 1240, 460, CYAN, width=6)
    c.text("←  가중치 응답 · 연산 완료", 720, 610, 475, 48, size=20, color=TEAL, bold=True, align="center")
    c.line(680, 585, 1240, 585, TEAL, width=6)
    c.rect(735, 690, 450, 82, BG3, stroke=CYAN, stroke_width=2, radius=14)
    c.text("연산 타일(TileJob)\n계산 위치 + 데이터 위치", 755, 700, 410, 60, size=19, color=WHITE, bold=True, align="center", valign="middle")
    dashed_line(c, 700, 840, 1220, 840, RED, width=3)
    c.text("후속 물리 경로 · GTH · CDC · MIG", 730, 865, 460, 38, size=17, color=RED, bold=True, align="center")


def slide_5(c: SlideCanvas) -> None:
    header(c, META[4])
    c.text("기존 접근", 145, 270, 650, 52, size=27, color=MUTED, bold=True, align="center")
    c.text("본 연구", 1125, 270, 650, 52, size=27, color=CYAN, bold=True, align="center")
    c.line(960, 260, 960, 825, GRID, width=2)
    left = ("연산량 또는 대역폭 중심", "정적 작업 배정", "합성 행렬 중심", "RTL과 시스템 제안이 분리")
    right = ("연산·메모리·링크 균형", "이동 비용을 고려한 선택적 재분배", "실제 Gemma 그래프에서 작업 추출", "모델 → RTL → PCB 흐름 연결")
    for index, (before, after) in enumerate(zip(left, right)):
        y = 385 + index * 110
        c.text(before, 150, y, 640, 62, size=22, color=WHITE, align="center", valign="middle")
        arrow_between(c, 855, y + 6, 115, CYAN)
        c.text(after, 1080, y, 740, 62, size=22, color=WHITE, bold=True, align="center", valign="middle")
    c.rect(245, 855, 1430, 82, BG2, stroke=AMBER, stroke_width=2, radius=18)
    c.text("연구 핵심 · 외부 메모리가 유효해지는 시스템 조건 도출", 285, 873, 1350, 44, size=21, color=AMBER, bold=True, align="center", valign="middle")


def slide_6(c: SlideCanvas) -> None:
    header(c, META[5])
    embedded_movie(c, "tile_dataflow", 70, 275, 1120, 620)
    steps = (("1", "실제 가중치 추출"), ("2", "시험 입력 생성"), ("3", "16×4 INT8 MatVec RTL"), ("4", "소프트웨어 기준값 비교"))
    for index, (number, label) in enumerate(steps):
        y = 300 + index * 122
        c.circle(1285, y, 58, BG3, stroke=CYAN, stroke_width=2)
        c.text(number, 1285, y + 10, 58, 34, size=17, color=CYAN, bold=True, align="center", valign="middle")
        c.text(label, 1370, y + 4, 420, 52, size=20, color=WHITE, bold=True, valign="middle")
        if index < 3:
            c.line(1314, y + 62, 1314, y + 112, GRID, width=2)
    c.rect(1280, 805, 505, 86, BG2, stroke=TEAL, stroke_width=3, radius=20)
    c.text("3/3 정확히 일치", 1310, 822, 445, 50, size=27, color=TEAL, bold=True, align="center", valign="middle")
    c.text("실제 DDR 및 FPGA 간 링크는 이 시험에 포함하지 않음", 1180, 930, 650, 38, size=14, color=AMBER, bold=True, align="right")


def slide_7(c: SlideCanvas) -> None:
    header(c, META[6])
    c.text("정적 로컬 큐(S1)", 150, 260, 620, 48, size=24, color=BLUE, bold=True, align="center")
    lanes_y = [390, 510, 630, 750]
    counts = [8, 3, 0, 0]
    for cluster, (y, count) in enumerate(zip(lanes_y, counts)):
        c.text(f"C{cluster}", 125, y + 6, 80, 36, size=20, bold=True, align="center")
        c.rect(230, y, 600, 56, BG2, stroke=GRID, stroke_width=2, radius=9)
        for j in range(count):
            c.rect(245 + j * 66, y + 10, 50, 36, BLUE if cluster == 0 else "#2B667B", radius=6)
        if count == 0:
            c.text("유휴", 460, y + 8, 140, 34, size=20, color=AMBER, bold=True, align="center")
    arrow_between(c, 900, 530, 130, AMBER)
    c.text("마지막 작업 대기", 860, 600, 210, 38, size=19, color=AMBER, bold=True, align="center")
    c.line(1130, 360, 1130, 820, GRID, width=2)
    c.text("장점", 1325, 390, 270, 44, size=22, color=TEAL, bold=True, align="center")
    c.text("가중치 이동이 적다", 1190, 465, 540, 54, size=25, color=WHITE, bold=True, align="center")
    c.line(1190, 575, 1730, 575, GRID, width=2)
    c.text("한계", 1325, 620, 270, 44, size=22, color=AMBER, bold=True, align="center")
    c.text("작업이 한곳에 몰릴 수 있다", 1160, 700, 600, 58, size=25, color=WHITE, bold=True, align="center")


def slide_8(c: SlideCanvas) -> None:
    header(c, META[7])
    labels = ("1  유휴 연산기 발견", "2  이동 이득과 비용 비교", "3  이득이 클 때만 작업 이전")
    for index, label in enumerate(labels):
        x = 130 + index * 590
        c.text(label, x, 240, 520, 48, size=20, color=CYAN if index == 2 else WHITE, bold=True, align="center")
    embedded_movie(c, "work_stealing_sequence", 160, 310, 1600, 540)
    c.rect(470, 885, 980, 60, BG2, stroke=CYAN, stroke_width=2, radius=20)
    c.text("대기 감소 이득 > 데이터 이동 비용", 510, 897, 900, 40, size=23, color=CYAN, bold=True, align="center")
    c.text("지역성 인식 작업 재분배(S3)", 650, 958, 620, 34, size=16, color=MUTED, bold=True, align="center")


def slide_9(c: SlideCanvas) -> None:
    header(c, META[8])
    stages = [("전체 ONNX 노드", "7,837", BLUE), ("토큰당 행렬 투영", "183", CYAN), ("분할된 연산 타일", "802", TEAL)]
    xs = [140, 770, 1400]
    for index, (label, value, color) in enumerate(stages):
        c.circle(xs[index], 330, 240, BG2, stroke=color, stroke_width=3)
        c.text(value, xs[index] + 10, 385, 220, 70, size=38, color=color, bold=True, align="center")
        c.text(label, xs[index] - 30, 600, 300, 44, size=22, color=WHITE, bold=True, align="center")
        if index < 2:
            arrow_between(c, xs[index] + 315, 425, 130, color)
    c.text("큰 행렬을 출력 방향의 여러 타일로 분할", 590, 705, 740, 42, size=19, color=AMBER, bold=True, align="center")
    c.line(160, 780, 1760, 780, GRID, width=2)
    c.text("보수적 단계 규칙 · 질의·키·값  →  출력 투영  →  게이트·상향  →  하향  →  다음 계층", 180, 815, 1560, 54, size=21, color=CYAN, bold=True, align="center")
    c.text("32토큰 = 25,664개 연산 타일", 500, 905, 920, 46, size=22, color=WHITE, bold=True, align="center")


def slide_10(c: SlideCanvas) -> None:
    header(c, META[9])
    c.text("동일 작업 부하 · 동일 초기값 · 동일 하드웨어 가정", 360, 230, 1200, 38, size=19, color=AMBER, bold=True, align="center")
    embedded_movie(c, "scheduler_timeline", 140, 285, 1640, 620)
    c.text("대기 감소가 이동 비용보다 큰 대표 사례", 520, 930, 880, 44, size=22, color=CYAN, bold=True, align="center")


def slide_11(c: SlideCanvas) -> None:
    header(c, META[10])
    embedded_movie(c, "tail_latency_results", 80, 270, 1000, 590)
    c.text("합성 불균형 부하", 1150, 285, 600, 42, size=22, color=CYAN, bold=True, align="center")
    metric(c, 1140, 350, "-19.13%", "p95", CYAN, 310)
    metric(c, 1460, 350, "-18.71%", "p99", TEAL, 310)
    c.line(1140, 535, 1770, 535, GRID, width=2)
    c.text("Gemma 초기 배치", 1150, 575, 600, 42, size=22, color=AMBER, bold=True, align="center")
    c.text("기존 초기 배치", 1160, 665, 270, 38, size=18, color=MUTED, align="center")
    c.text("+0.28%", 1160, 710, 270, 58, size=28, color=AMBER, bold=True, align="center")
    c.text("메모리 채널 고려", 1470, 665, 290, 38, size=18, color=MUTED, align="center")
    c.text("-19.79%", 1470, 710, 290, 58, size=28, color=TEAL, bold=True, align="center")
    c.text("스케줄링 효과는 초기 데이터 배치와 분리할 수 없다", 1110, 855, 700, 50, size=21, color=WHITE, bold=True, align="center")


def slide_12(c: SlideCanvas) -> None:
    header(c, META[11])
    embedded_movie(c, "bottleneck_migration", 170, 265, 1580, 620)
    c.text("서로 다른 기준 비교이며 하나의 인과효과로 합치지 않는다", 350, 920, 1220, 46, size=20, color=AMBER, bold=True, align="center")


def slide_13(c: SlideCanvas) -> None:
    header(c, META[12])
    rows = read_csv(ROOT / "results/model_level/k26_local_external_sensitivity.csv")
    local = sorted((float(r["modeled_aggregate_memory_gbps"]), float(r["total_completion_cycles"]) / 1e6) for r in rows if r["architecture"] == "k26_local_shared_ddr4" and r["placement"] == "size_aware")
    external = sorted((float(r["modeled_aggregate_memory_gbps"]), float(r["total_completion_cycles"]) / 1e6) for r in rows if r["architecture"] == "external_4ch_ddr3l_candidate" and r["placement"] == "size_aware" and r["scheduler"] == "S1" and r["link_width_bits_per_bundle"] == "128")
    x0, y0, x1, y1 = 210, 850, 1640, 300
    c.line(x0, y0, x1, y0, MUTED, width=3)
    c.line(x0, y0, x0, y1, MUTED, width=3)
    for tick in (3, 6, 9, 12, 15):
        x = x0 + (tick - 3) / 12 * (x1 - x0)
        c.line(x, y0, x, y0 + 14, MUTED, width=2)
        c.text(str(tick), x - 35, y0 + 20, 70, 30, size=13, color=MUTED, align="center")
    for tick in (0, 50, 100, 150, 200, 250):
        y = y0 - tick / 250 * (y0 - y1)
        c.line(x0 - 12, y, x0, y, MUTED, width=2)
        c.line(x0, y, x1, y, GRID, width=1)
        c.text(str(tick), 125, y - 14, 70, 28, size=12, color=MUTED, align="right")
    c.text("모델 입력 유효 메모리 대역폭 (GB/s)", 550, 915, 760, 38, size=18, color=MUTED, align="center")
    c.text("완료 사이클\n(백만)", 30, 500, 150, 76, size=16, color=MUTED, align="center", valign="middle")
    def draw_series(points, color):
        coords = []
        for bandwidth, cycles in points:
            x = x0 + (bandwidth - 3) / 12 * (x1 - x0)
            y = y0 - cycles / 250 * (y0 - y1)
            coords.append((x, y))
        for start, end in zip(coords, coords[1:]):
            c.line(start[0], start[1], end[0], end[1], color, width=5)
        for x, y in coords:
            c.circle(x - 10, y - 10, 20, color, stroke=WHITE, stroke_width=1)
    draw_series(local, BLUE)
    draw_series(external, TEAL)
    c.text("K26 로컬", 1480, 690, 230, 36, size=18, color=BLUE, bold=True)
    c.text("외부 4채널", 1480, 385, 250, 36, size=18, color=TEAL, bold=True)
    c.rect(650, 235, 620, 70, BG2, stroke=AMBER, stroke_width=2, radius=16)
    c.text("시험 범위에서 교차 없음", 680, 249, 560, 44, size=24, color=AMBER, bold=True, align="center")
    c.text("9.6 GB/s · 24.6M", 980, 765, 320, 34, size=15, color=BLUE, bold=True)
    c.text("6.4 GB/s · 120.6M", 650, 555, 340, 34, size=15, color=TEAL, bold=True)
    c.text("더 큰 모델 · 긴 문맥 · 실제 로컬 경합에서 재검증", 520, 965, 920, 34, size=17, color=WHITE, bold=True, align="center")


def prepare_kicad() -> tuple[Path, Path, Path]:
    source_path = ROOT / "paper/final/figures/paper_f07_kicad_coupon_render.png"
    source = Image.open(source_path).convert("RGB")
    crops = (
        (ASSETS / "kicad_link_crop.png", (420, 100, 1380, 610)),
        (ASSETS / "kicad_refclk_crop.png", (570, 210, 1180, 620)),
        (ASSETS / "kicad_routed_crop.png", (260, 360, 1050, 850)),
    )
    for path, bounds in crops:
        source.crop(bounds).save(path)
    return tuple(path for path, _ in crops)


def slide_14(c: SlideCanvas) -> None:
    header(c, META[13])
    render = ROOT / "paper/final/figures/paper_f07_kicad_coupon_render.png"
    crops = prepare_kicad()
    c.image_file(render, 70, 250, 1210, 650, mode="cover")
    c.rect(70, 250, 1210, 650, BG, stroke=TEAL, stroke_width=2, radius=12, alpha=0)
    c.rect(400, 535, 540, 58, BG, stroke=RED, stroke_width=2, radius=12, alpha=30)
    c.text("참조 라우팅 쿠폰", 420, 545, 500, 36, size=21, color=RED, bold=True, align="center")
    labels = ("범용 경계 링크", "기준 클록 차동쌍", "DDR3L x16 대표 슬라이스")
    for index, (crop, label) in enumerate(zip(crops, labels)):
        y = 255 + index * 205
        c.image_file(crop, 1360, y, 450, 145, mode="cover")
        c.rect(1360, y, 450, 145, BG, stroke=CYAN if index == 0 else TEAL, stroke_width=2, radius=8, alpha=0)
        c.text(label, 1340, y + 150, 490, 38, size=19, color=WHITE, bold=True, align="center")
        c.line(1285, y + 72, 1345, y + 72, CYAN, width=2)
    badges = (("풋프린트 29개", BLUE, 270, 18), ("GTH·기준 클록 20개", TEAL, 380, 16), ("제한 범위 ERC/DRC 0", CYAN, 440, 16))
    x = 100
    for label, color, width, label_size in badges:
        c.rect(x, 930, width, 54, BG2, stroke=color, stroke_width=2, radius=18)
        c.text(label, x + 10, 938, width - 20, 36, size=label_size, color=color, bold=True, align="center")
        x += width + 24
    c.text("전체 116개 중 55개 미배선", 1330, 900, 500, 34, size=17, color=AMBER, bold=True, align="center")
    c.text("참조 쿠폰 · 제작용 아님", 1330, 960, 500, 34, size=14, color=RED, bold=True, align="center")


def make_qr() -> Path:
    import cairosvg

    output = ASSETS / "repository_qr.png"
    widget = QrCodeWidget(REPO_URL)
    x1, y1, x2, y2 = widget.getBounds()
    size = 360
    drawing = Drawing(size, size, transform=[size / (x2 - x1), 0, 0, size / (y2 - y1), 0, 0])
    drawing.add(widget)
    svg_path = ASSETS / "repository_qr.svg"
    renderSVG.drawToFile(drawing, str(svg_path))
    cairosvg.svg2png(url=str(svg_path), write_to=str(output), output_width=720, output_height=720, background_color="white")
    return output


def slide_15(c: SlideCanvas) -> None:
    header(c, META[14])
    stages = (("현재", "제한된 RTL과\n분석 모델", AMBER), ("다음 1", "K26 로컬\n기준선 실측", BLUE), ("다음 2", "DDR → 링크 → MatVec\n물리 경로 통합", CYAN), ("다음 3", "SI/PI · 전력 · 열\n검증", TEAL))
    xs = [85, 535, 985, 1435]
    detail_layout = [(0, 340), (400, 360), (790, 470), (1320, 400)]
    for index, (label, detail, color) in enumerate(stages):
        c.circle(xs[index], 300, 80, BG3, stroke=color, stroke_width=3)
        c.text(label, xs[index], 320, 80, 40, size=16, color=color, bold=True, align="center", valign="middle")
        detail_x, detail_w = detail_layout[index]
        c.text(detail, detail_x, 455, detail_w, 125, size=20, color=WHITE, bold=True, align="center", valign="middle")
        if index < 3:
            arrow_between(c, xs[index] + 260, 330, 115, color)
    c.line(120, 730, 1800, 730, GRID, width=2)
    c.text("검증 순서", 140, 795, 180, 40, size=18, color=CYAN, bold=True)
    c.text("정확성  →  대역폭  →  꼬리 지연  →  전력·신호 무결성", 340, 775, 1330, 72, size=25, color=WHITE, bold=True, align="center", valign="middle")
    c.text("분석을 측정으로 바꾸는 순서", 610, 900, 700, 42, size=20, color=MUTED, align="center")


def slide_16(c: SlideCanvas) -> None:
    header(c, META[15])
    lines = [
        ("1", "K26 계산부와 외부 다채널 가중치 공급부를 분리한 확장형 구조를 설계했다.", BLUE),
        ("2", "Gemma 그래프 기반 분석에서 불균형 이득이 이동 비용보다 클 때만 재분배가 유효했다.", CYAN),
        ("3", "이번 민감도에서 1B는 K26 로컬 우선—더 큰 모델·긴 문맥·실제 경합에서 재평가한다.", TEAL),
    ]
    for index, (number, sentence, color) in enumerate(lines):
        y = 270 + index * 190
        c.text(number, 110, y, 130, 74, size=31, color=color, bold=True, align="center", valign="middle")
        c.line(265, y + 36, 335, y + 36, color, width=4)
        c.text(sentence, 370, y - 2, 1160, 95, size=22, color=WHITE, bold=True, valign="middle")
    qr = make_qr()
    c.rect(1595, 330, 220, 220, WHITE, radius=8)
    c.image_file(qr, 1610, 345, 190, 190, mode="contain")
    c.text("코드 · 데이터 · 논문", 1525, 580, 360, 38, size=17, color=CYAN, bold=True, align="center")
    c.text("github.com/snowman0919/\nvarp-k26-memory-fpga", 1490, 635, 430, 70, size=13, color=MUTED, align="center")
    c.rect(250, 855, 1280, 72, BG2, stroke=AMBER, stroke_width=2, radius=18)
    c.text("기여 · 연산 배치와 메모리 위치를 함께 고려한 설계 흐름", 290, 870, 1200, 42, size=20, color=AMBER, bold=True, align="center", valign="middle")


def slide_17(c: SlideCanvas) -> None:
    c.rect(0, 0, W, H, BG)
    c.circle(660, 120, 600, BG2, stroke="#123B51", stroke_width=3)
    c.circle(760, 220, 400, BG, stroke=CYAN, stroke_width=2)
    c.rect(770, 610, 380, 8, CYAN)
    c.text("질문과 토론", 410, 370, 1100, 190, size=62, color=WHITE, bold=True, align="center", valign="middle")
    c.text("감사합니다", 760, 930, 400, 46, size=18, color=MUTED, align="center")


BUILDERS = (slide_1, slide_2, slide_3, slide_4, slide_5, slide_6, slide_7, slide_8, slide_9, slide_10, slide_11, slide_12, slide_13, slide_14, slide_15, slide_16, slide_17)


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8", newline="") as handle:
        return list(csv.DictReader(handle))


def write_notes() -> None:
    lines = [
        "# VARP v11 — 본선 9분 35초 발표자 노트",
        "",
        "총 발표 분량: 9:35 · Slide 17 Q&A는 발표 시간에서 제외",
        "",
        "## 본선 심사 기준 대응",
        "",
        "- 주제 이해도 30점: 모든 답을 필요성→구조→검증→채택 조건으로 연결한다.",
        "- 전달성 10점: 한 슬라이드 한 문장, 그림은 아래 순서대로 가리키며 읽지 말고 설명한다.",
        "- 질의응답 20점: 결론 한 문장→근거 유형·수치→해석 경계→다음 검증 순서로 20초 안에 답한다.",
        "- 참여도 5점: 전체 세션의 80% 이상 접속하고 다른 발표에 근거 있는 질문을 준비한다.",
        "- 시간 점검: Slide 4 끝 2:05 · Slide 8 끝 4:25 · Slide 12 끝 7:10 · Slide 14 끝 8:30 · Slide 16 끝 9:35.",
        "",
    ]
    for index, meta in enumerate(META, 1):
        minutes, seconds = divmod(meta.duration, 60)
        speech = " ".join(part for part in (meta.speech, NOTE_EXTENSIONS[index - 1]) if part)
        lines += [
            f"## Slide {index}: {meta.title} ({minutes}:{seconds:02d})",
            "",
            f"목표 시간: {minutes}:{seconds:02d}{' · 발표 시간 제외' if index == 17 else ''}",
            "",
            f"기억할 문장: {meta.memory}",
            "",
            f"발화문: {speech}",
            "",
            f"그림 설명 순서: {VISUAL_ORDERS[index - 1]}",
            "",
            f"전환: {meta.transition or '이상으로 발표를 마치겠습니다. 질문 받겠습니다.'}",
            "",
            f"예상 질문: {EXPECTED_QUESTIONS[index - 1]}",
            "",
            f"20초 답변: {ANSWER_CORES[index - 1]}",
            "",
            f"근거 파일: {SOURCE_FOOTERS[index - 1] or 'presentation/final/qna_bank.md'}",
            "",
            f"해석 경계: {meta.supplement}",
            "",
            f"주의할 표현: {CAUTIONS[index - 1]}",
            "",
        ]
    (FINAL / "speaker_notes.md").write_text("\n".join(lines).rstrip() + "\n", encoding="utf-8")


def write_outline() -> None:
    lines = ["# 온디바이스 sLLM용 K26–Memory FPGA 가속기 설계", "", "형식: 16:9 · 17장 · 발표 목표 9:35 · Slide 17 Q&A 시간 제외 · 한국어 기술 컨퍼런스 발표", "", "중심 문장: K26의 계산부와 외부 Memory FPGA의 가중치 공급부를 분리한 확장형 가속기 구조를 설계하고, 실제 Gemma 3 1B 작업에서 작업 불균형을 줄이는 이득이 원격 데이터 이동 비용보다 커지는 조건을 분석하였다.", "", "## 구성", ""]
    for index, meta in enumerate(META, 1):
        minutes, seconds = divmod(meta.duration, 60)
        lines.append(f"{index}. {meta.title} — {minutes}:{seconds:02d} · {meta.conclusion}")
    (FINAL / "outline.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def write_source_index() -> None:
    rows = [
        (1, "없음", "cover_background.png", "image_gen; build_v11_deck.py", "연구 주제의 개념 이미지", "실제 회로나 구현 보드"),
        (2, "없음", "편집 가능한 목차", "build_v11_deck.py", "발표 질문과 진행 순서", "연구 결과"),
        (3, "k26_local_external_sensitivity.csv", "편집 가능한 의사결정식", "run_k26_local_baseline.py; build_v11_deck.py", "2.43 GiB 용량 모델과 연구 질문 전환", "K26 보드 측정"),
        (4, "없음", "편집 가능한 K26↔Memory FPGA 구조", "build_v11_deck.py; ClosedLoopVirtualPrototypeTop.scala", "역할 분리 구조와 논리/물리 경계", "GTH·CDC·MIG 물리 완료"),
        (5, "없음", "편집 가능한 관련 연구 비교", "build_v11_deck.py; submission_manuscript.md", "문헌에 적힌 비교 가능한 연구 공백", "Work Stealing 알고리즘 신규성"),
        (6, "gemma3_1b_rtl_tile_parity.csv", "tile_dataflow.mp4/frame; 편집 가능한 검증 흐름", "manim_scenes.py; GemmaWeightTileRtlParitySpec.scala", "실제 가중치 3개 타일의 MatVec RTL 산술 일치", "전체 Gemma·DDR·링크·성능"),
        (7, "없음", "편집 가능한 정적 큐 개념", "build_v11_deck.py", "정적 소유권의 지역성과 불균형", "수치 성능"),
        (8, "없음", "work_stealing_sequence.mp4/frame", "manim_scenes.py; k26_scheduler_model.py", "이득과 이동 비용을 비교하는 정책", "최적성·보드 타이밍"),
        (9, "projection_trace.csv; gemma3_1b_dependency_manifest.json", "편집 가능한 7,837→183→802 흐름", "gemma_dependency_model.py; build_v11_deck.py", "ONNX 형상과 모델링한 보수적 단계 의존·타일 분할", "전체 ONNX 의존성 재현·실제 FPGA 실행 횟수"),
        (10, "s1_s3_timeline_events.csv", "scheduler_timeline.mp4/frame", "generate_conference_figures.py; manim_scenes.py", "동일 합성 작업·초기값에서 정책만 바꾼 대표 사건", "Gemma 사건·물리 cycle"),
        (11, "paired_policy_effects.csv; gemma3_1b_policy_effects.csv", "tail_latency_results.mp4/frame; 편집 가능한 결과", "build_v11_research_summary.py; manim_scenes.py", "합성 짝비교와 Gemma 배치 민감도", "사용자 요청 latency·보편적 우월성"),
        (12, "paired_policy_effects.csv", "bottleneck_migration.mp4/frame", "manim_scenes.py", "정적·무제약 기준을 분리한 비용 비교", "세 수치가 같은 비교라는 해석"),
        (13, "k26_local_external_sensitivity.csv", "편집 가능한 민감도 곡선", "run_k26_local_baseline.py; build_v11_deck.py", "시험 범위에서 곡선 미교차와 K26 로컬 우선", "보드 실측·미래 모든 조건의 우월성"),
        (14, "없음", "실제 KiCad 렌더와 확대 이미지", "verify_k26_kicad.py; kicad-cli; build_v11_deck.py", "참조 쿠폰 객체·DDR3L x16 대표 슬라이스·부분 배선·제한 검사", "완성 메모리 채널·제작 가능 보드·전체 DRC 0"),
        (15, "없음", "편집 가능한 물리 검증 순서", "build_v11_deck.py; architecture.md", "현재 한계와 후속 실측 순서", "실측 완료"),
        (16, "없음", "편집 가능한 3문장 결론·QR", "build_v11_deck.py", "구조·조건·현재 채택 판단", "외부 메모리 보편 실패·실측 완료"),
        (17, "없음", "편집 가능한 질문과 토론 화면", "build_v11_deck.py; qna_bank.md", "구두 질의응답 시작", "추가 기술 결과"),
    ]
    lines = ["# Slide Source Index", "", "모든 경로는 저장소 root 기준이다. 화면에서 생략한 시드·분석 모델 조건은 speaker notes와 연구 CSV에 남긴다.", "", "| Slide | 원본 CSV | 사용 Figure / 자료 | 생성·검증 스크립트 | 허용 해석 | 금지 해석 |", "|---:|---|---|---|---|---|"]
    for row in rows:
        lines.append("| " + " | ".join(str(value).replace("|", "\\|") for value in row) + " |")
    (FINAL / "slide_source_index.md").write_text("\n".join(lines) + "\n", encoding="utf-8")


def build_pdf(paths: list[Path]) -> None:
    rl_config.invariant = 1
    page_size = (960, 540)
    pdf = pdf_canvas.Canvas(str(PDF), pagesize=page_size, pageCompression=1)
    pdf.setTitle("온디바이스 sLLM용 K26–Memory FPGA 가속기 설계")
    pdf.setAuthor("CHOI YUNHYUK")
    for path in paths:
        pdf.drawImage(str(path), 0, 0, width=960, height=540)
        pdf.showPage()
    pdf.save()


def contact_sheet(paths: list[Path]) -> Path:
    tw, th = 480, 270
    rows = (len(paths) + 1) // 2
    sheet = Image.new("RGB", (tw * 2 + 84, th * rows + 135), rgb("#030812"))
    draw = ImageDraw.Draw(sheet)
    draw.text((42, 26), f"VARP v11 · {len(paths)}-slide contact sheet", font=font(18, True), fill=rgb(WHITE))
    for index, path in enumerate(paths):
        image = Image.open(path).convert("RGB").resize((tw, th), Image.Resampling.LANCZOS)
        col, row = index % 2, index // 2
        x, y = 42 + col * tw, 90 + row * th
        sheet.paste(image, (x, y))
        draw.text((x + tw - 30, y + 24), f"{index + 1:02d}", font=font(9, True), fill=rgb(CYAN), anchor="mm")
    output = FINAL / "slide_contact_sheet.png"
    sheet.save(output)
    return output


def main() -> int:
    SLIDES.mkdir(parents=True, exist_ok=True)
    for stale in SLIDES.glob("slide_*.png"):
        stale.unlink()
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = META[0].title
    prs.core_properties.author = "CHOI YUNHYUK"
    prs.core_properties.created = datetime(2026, 8, 1, 0, 0, 0)
    prs.core_properties.modified = datetime(2026, 8, 1, 0, 0, 0)
    paths: list[Path] = []
    for index, (meta, builder) in enumerate(zip(META, BUILDERS), 1):
        speech = f"{meta.speech} {NOTE_EXTENSIONS[index - 1]}"
        canvas = SlideCanvas(
            prs,
            meta.title,
            f"기억할 문장: {meta.memory}\n\n"
            f"발화문: {speech}\n\n"
            f"그림 설명 순서: {VISUAL_ORDERS[index - 1]}\n\n"
            f"전환: {meta.transition}\n\n"
            f"예상 질문: {EXPECTED_QUESTIONS[index - 1]}\n\n"
            f"20초 답변: {ANSWER_CORES[index - 1]}\n\n"
            f"근거 파일: {SOURCE_FOOTERS[index - 1] or 'presentation/final/qna_bank.md'}\n\n"
            f"해석 경계: {meta.supplement}\n\n"
            f"주의할 표현: {CAUTIONS[index - 1]}",
        )
        builder(canvas)
        if index < 17:
            footer(canvas, index - 1)
        paths.append(canvas.save(index))
    prs.save(PPTX)
    normalize_pptx(PPTX)
    write_notes()
    write_outline()
    write_source_index()
    build_pdf(paths)
    sheet = contact_sheet(paths)
    print(f"pptx={PPTX} slides={len(prs.slides)}")
    print(f"pdf={PDF}")
    print(f"contact_sheet={sheet}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
