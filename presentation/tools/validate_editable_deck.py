#!/usr/bin/env python3
"""Fail-closed validation for the 16-slide final-round conference deck."""

from __future__ import annotations

import json
import hashlib
from pathlib import Path
import re
import subprocess
import sys
import zipfile

from PIL import Image, ImageFont
from fontTools.ttLib import TTFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader


ROOT = Path(__file__).resolve().parents[2]
FINAL = ROOT / "presentation/final"
PPTX = FINAL / "presentation.pptx"
PDF = FINAL / "presentation.pdf"
TITLES = [
    "온디바이스 sLLM용 K26–Memory FPGA 가속기 구조",
    "오늘 답할 네 가지 질문",
    "Gemma 3 1B는 4GB 안에 들어갔다—그래서 질문을 바꿨다",
    "K26은 계산하고, Memory FPGA는 가중치를 공급한다",
    "메모리 용량만 늘리지 않고, 작업 배치와 데이터 위치를 함께 설계했다",
    "실제 Gemma 가중치를 MatVec RTL에서 검증했다",
    "고정 배정은 데이터를 가까이 두지만 작업을 고르게 나누지 못한다",
    "빈 연산기가 이득이 있는 작업만 가져온다",
    "Gemma 그래프를 802개 연산 타일로 나눴다",
    "같은 작업에서 고정 배정과 작업 재배분만 바꿨다",
    "작업이 치우친 경우에만 재배분이 효과적이었다",
    "유휴 연산은 줄었지만 데이터 이동이 새로운 비용이 됐다",
    "현재 가정에서는 K26의 로컬 메모리 경로가 더 유리했다",
    "다음 단계는 분석 모델을 실제 하드웨어 경로로 닫는 것이다",
    "결론과 기여",
    "질문과 토론",
]
MOVIES = ("tile_dataflow", "work_stealing_sequence", "scheduler_timeline", "tail_latency_results", "bottleneck_migration")
FONT_REGULAR = "/usr/share/fonts/truetype/nanum/NanumSquareR.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/nanum/NanumSquareB.ttf"


def fail(message: str) -> None:
    raise AssertionError(message)


def main() -> int:
    required = [
        PPTX, PDF, FINAL / "speaker_notes.md", FINAL / "outline.md",
        FINAL / "full_script.md", FINAL / "rehearsal_script.txt",
        FINAL / "slide_source_index.md", FINAL / "qna_bank.md",
        FINAL / "peer_question_bank.md", FINAL / "review_report.md",
        FINAL / "slide_contact_sheet.png", FINAL / "presentation.mp4",
        FINAL / "presentation.gif", FINAL / "animation_manifest.md",
        ROOT / "research/v11_research_freeze.json",
    ]
    for stem in MOVIES:
        required += [FINAL / "assets" / f"{stem}.mp4", FINAL / "assets" / f"{stem}.gif", FINAL / "assets" / f"{stem}_frame.png"]
    for path in required:
        if not path.is_file() or path.stat().st_size == 0:
            fail(f"missing artifact: {path.relative_to(ROOT)}")

    with zipfile.ZipFile(PPTX) as package:
        if package.testzip():
            fail("PPTX ZIP contains a corrupt entry")
        movies = [name for name in package.namelist() if name.startswith("ppt/media/") and name.endswith(".mp4")]
        if len(movies) != len(MOVIES):
            fail(f"expected {len(MOVIES)} embedded Manim movies, got {len(movies)}")
        packaged_hashes = {hashlib.sha256(package.read(name)).hexdigest() for name in movies}
        asset_hashes = {
            hashlib.sha256((FINAL / "assets" / f"{stem}.mp4").read_bytes()).hexdigest()
            for stem in MOVIES
        }
        if packaged_hashes != asset_hashes:
            fail("embedded movies do not match the final MP4 assets")
        for name in package.namelist():
            if not name.startswith("ppt/slides/_rels/") or not name.endswith(".rels"):
                continue
            rels = package.read(name).decode("utf-8")
            if "TargetMode=\"External\"" in rels and ("video" in rels or ".mp4" in rels):
                fail(f"external video relationship remains: {name}")
        for slide_number in (6, 8, 10, 11, 12):
            rel_name = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
            rels = package.read(rel_name).decode("utf-8")
            if rels.count("/relationships/video") != 1 or rels.count("/relationships/media") != 1:
                fail(f"slide {slide_number}: broken embedded-video relationships")
            if "TargetMode=\"External\"" in rels or "../media/" not in rels:
                fail(f"slide {slide_number}: movie is not internally embedded")
            if "/relationships/image" not in rels:
                fail(f"slide {slide_number}: poster-frame relationship missing")

    for stem in MOVIES:
        movie = FINAL / "assets" / f"{stem}.mp4"
        probe = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration:stream=codec_name,width,height,codec_type", "-of", "json", str(movie)],
            check=True, capture_output=True, text=True,
        )
        payload = json.loads(probe.stdout)
        duration = float(payload["format"]["duration"])
        video_streams = [stream for stream in payload["streams"] if stream.get("codec_type") == "video"]
        audio_streams = [stream for stream in payload["streams"] if stream.get("codec_type") == "audio"]
        if not 8.0 <= duration <= 15.0:
            fail(f"{stem}: duration {duration:.2f}s is outside 8–15s")
        if len(video_streams) != 1 or audio_streams:
            fail(f"{stem}: expected one silent video stream")
        stream = video_streams[0]
        if stream.get("codec_name") != "h264" or stream.get("width", 0) < 1280 or stream.get("height", 0) < 720:
            fail(f"{stem}: expected H.264 at 1280×720 or larger")

    prs = Presentation(PPTX)
    if len(prs.slides) != 16:
        fail(f"expected 16 slides, got {len(prs.slides)}")
    if abs(prs.slide_width / prs.slide_height - 16 / 9) > 0.002:
        fail("deck is not 16:9")

    min_font = 999.0
    checked_runs = pictures = media = 0
    slide_w, slide_h = prs.slide_width, prs.slide_height
    tolerance = 2 * 914400 / 144
    cmaps = {False: TTFont(FONT_REGULAR).getBestCmap(), True: TTFont(FONT_BOLD).getBestCmap()}
    forbidden_terms = (
        "후보 구조 평가", "이번 민감도", "기존 산출", "채널 친화", "논리 폐루프",
        "원본 사건 CSV", "reference coupon", "not for fabrication", "TileJob",
        "payload", "contract", "bounded", "full-overlap", "home ownership",
    )
    for index, (slide, title) in enumerate(zip(prs.slides, TITLES), 1):
        texts = [shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        screen = " ".join(value.replace("\n", " ") for value in texts)
        if title not in screen:
            fail(f"slide {index}: title missing")
        if index == 16:
            allowed = {"질문과 토론", "감사합니다"}
            visible = {value.strip() for value in texts if value.strip()}
            if not visible.issubset(allowed):
                fail(f"slide 16: Q&A detail must be presenter-only: {sorted(visible - allowed)!r}")
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if not notes.startswith("기억할 문장:"):
            fail(f"slide {index}: notes do not start with memory sentence")
        for term in forbidden_terms:
            if term.lower() in screen.lower():
                fail(f"slide {index}: forbidden screen wording: {term!r}")
        if any(term in screen.lower() for term in ("kicad", "pcb", "쿠폰", "제작용")):
            fail(f"slide {index}: removed routing-coupon content remains visible")
        if any(token in screen for token in ("-15.07%", "-14.61%", "-18.12%", "-37.84%", "+0.98%", "5,856")):
            fail(f"slide {index}: obsolete result is visible")
        if "□" in screen or "−" in screen:
            fail(f"slide {index}: unsupported glyph is visible")
        if re.search(r"\bS1\b", screen) and index != 7:
            fail(f"slide {index}: S1 may appear only at its first-definition slide")
        if re.search(r"\bS3\b", screen) and index != 8:
            fail(f"slide {index}: S3 may appear only at its first-definition slide")
        pictures += sum(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        media += sum(shape.shape_type == MSO_SHAPE_TYPE.MEDIA for shape in slide.shapes)
        for shape in slide.shapes:
            if shape.left < -tolerance or shape.top < -tolerance or shape.left + shape.width > slide_w + tolerance or shape.top + shape.height > slide_h + tolerance:
                fail(f"slide {index}: shape outside canvas: {shape.name}")
            if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
                continue
            explicit_lines = max(1, shape.text.count("\n") + 1)
            run_sizes = [run.font.size.pt for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text.strip() and run.font.size]
            if run_sizes:
                estimated_height_px = explicit_lines * max(run_sizes) * 2 * 1.05
                shape_height_px = shape.height / 914400 * 144
                if estimated_height_px > shape_height_px * 1.06:
                    fail(f"slide {index}: vertical text overflow risk: {shape.text!r}")
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not run.text.strip() or run.font.size is None:
                        continue
                    size = run.font.size.pt
                    min_font = min(min_font, size)
                    checked_runs += 1
                    top_px = shape.top / 914400 * 144
                    is_navigation = top_px < 65
                    is_footer = top_px >= 1000
                    is_cover_qualifier = index == 1 and top_px >= 900
                    floor = 9.5 if is_navigation or is_footer or is_cover_qualifier else 17.5
                    if size < floor:
                        fail(f"slide {index}: text below {floor:.0f}pt: {run.text!r}")
                    bold = bool(run.font.bold)
                    face = ImageFont.truetype(FONT_BOLD if bold else FONT_REGULAR, max(10, round(size * 2)))
                    missing = sorted({char for char in run.text if not char.isspace() and ord(char) not in cmaps[bold]})
                    if missing:
                        fail(f"slide {index}: font lacks glyphs {missing!r} in {run.text!r}")
                    width_px = shape.width / 914400 * 144
                    for line in run.text.splitlines() or [run.text]:
                        if line != "▶" and face.getlength(line) > width_px * 1.03:
                            fail(f"slide {index}: horizontal text overflow risk: {line!r}")

    if media != 5:
        fail(f"expected five movie shapes, got {media}")
    for slide_number in (8, 10, 12):
        count = sum(shape.shape_type == MSO_SHAPE_TYPE.MEDIA for shape in prs.slides[slide_number - 1].shapes)
        if count != 1:
            fail(f"slide {slide_number}: required embedded movie missing")
    if pictures < 8:
        fail(f"expected poster/cover/QR raster pictures, got {pictures}")
    if len(PdfReader(str(PDF)).pages) != 16:
        fail("PDF does not have 16 pages")
    pngs = sorted((FINAL / "slides").glob("slide_*.png"))
    if len(pngs) != 16:
        fail(f"expected 16 slide PNGs, got {len(pngs)}")
    for path in pngs:
        with Image.open(path) as image:
            if image.size != (1920, 1080):
                fail(f"wrong slide PNG size: {path.name} {image.size}")

    notes_md = (FINAL / "speaker_notes.md").read_text(encoding="utf-8")
    if notes_md.count("기억할 문장:") != 16:
        fail("speaker notes must contain 16 memory sentences")
    for field in ("목표 시간:", "그림 설명 순서:", "예상 질문:", "20초 답변:", "근거 파일:", "해석 경계:", "주의할 표현:"):
        if notes_md.count(field) != 16:
            fail(f"speaker notes must contain 16 {field} entries")
    for rubric in ("주제 이해도 30점", "전달성 10점", "질의응답 20점", "참여도 5점"):
        if rubric not in notes_md:
            fail(f"speaker notes missing final-round rubric: {rubric}")
    durations = re.findall(r"^## Slide \d+:.+\((\d+):(\d{2})\)$", notes_md, flags=re.MULTILINE)
    total_seconds = sum(int(minutes) * 60 + int(seconds) for minutes, seconds in durations)
    if total_seconds != 570:
        fail(f"speaker-note target is {total_seconds}s, expected 570s")
    spoken_area = notes_md.split("## Slide 16:", 1)[0]
    spoken_sections = re.findall(r"^발화문: (.+)$|^전환: (.+)$", spoken_area, flags=re.MULTILINE)
    spoken = "".join((a or b) for a, b in spoken_sections).replace(" ", "")
    sentence_pauses = sum(spoken.count(mark) for mark in ".?!") * 0.35
    transition_pauses = 15 * 1.5
    # 290 non-space Korean characters/minute is a conservative technical-talk
    # pace; explicit sentence and slide-transition pauses are added separately.
    opening_pause = 3.0
    estimated_seconds = len(spoken) / 290 * 60 + sentence_pauses + transition_pauses + opening_pause
    if not 560 <= estimated_seconds <= 585:
        fail(f"script reading estimate {estimated_seconds:.1f}s outside 9:20–9:45")

    for file_name in ("speaker_notes.md", "full_script.md", "rehearsal_script.txt"):
        script_text = (FINAL / file_name).read_text(encoding="utf-8")
        for term in ("이번 민감도", "후보 구조 평가", "논리 폐루프", "원본 사건 CSV", "full-overlap"):
            if term in script_text:
                fail(f"{file_name}: forbidden spoken wording remains: {term!r}")

    source_index = (FINAL / "slide_source_index.md").read_text(encoding="utf-8")
    for index in range(1, 17):
        if f"| {index} |" not in source_index:
            fail(f"source index missing slide {index}")
    for header in ("원본 CSV", "생성·검증 스크립트", "허용 해석", "금지 해석"):
        if header not in source_index:
            fail(f"source index missing {header}")

    qna = (FINAL / "qna_bank.md").read_text(encoding="utf-8")
    if qna.count("**20초 핵심 답변:**") < 40 or qna.count("**60초 확장 답변:**") < 40:
        fail("presenter Q&A bank must contain at least 40 complete questions")
    for attack in ("1B가 4GB에 들어가는데", "실제 FPGA 성능", "Work Stealing은 기존 기술", "왜 중앙 큐", "더 느린가", "GB/s와 cycle은 실측", "PCB는 제작", "전체 FPGA 경로", "연구가 실패"):
        if attack not in qna:
            fail(f"Q&A bank missing attack question: {attack}")
    peer = (FINAL / "peer_question_bank.md").read_text(encoding="utf-8")
    if len(re.findall(r"^## \d+\.", peer, flags=re.MULTILINE)) < 10:
        fail("peer question bank must contain at least ten questions")
    for field in ("AI 모델", "하드웨어", "데이터 분석", "웹·서비스", "보안", "로보틱스"):
        if field not in peer:
            fail(f"peer question bank missing field: {field}")

    audit = {
        "status": "PASS", "slides": 16, "aspect_ratio": "16:9",
        "embedded_manim_movies": media, "scoped_raster_pictures": pictures,
        "checked_text_runs": checked_runs, "minimum_font_pt": min_font,
        "out_of_bounds_shapes": 0, "obsolete_screen_values": 0,
        "speaker_note_memory_sentences": 16, "target_duration_seconds": total_seconds,
        "script_estimate_seconds_at_290_chars_per_minute_plus_pauses": round(estimated_seconds, 1),
        "qna_questions": qna.count("**20초 핵심 답변:**"),
        "peer_questions": len(re.findall(r"^## \d+\.", peer, flags=re.MULTILINE)),
        "rendered_png_size": [1920, 1080],
    }
    (FINAL / "layout_audit.json").write_text(json.dumps(audit, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"validated: 16 slides, 5 movies, 16 notes, {total_seconds}s target, {estimated_seconds:.1f}s reading estimate, minimum {min_font:.1f}pt")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except AssertionError as error:
        print(f"validation failed: {error}", file=sys.stderr)
        raise SystemExit(1)
