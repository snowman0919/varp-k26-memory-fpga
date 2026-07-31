#!/usr/bin/env python3
"""Build deterministic visual-evidence and presentation packages.

The generator reads only committed CSV evidence and native KiCad source/report
metadata. Generated publication files are written below ``build/`` and are
never tracked as source.
It preserves hybrid/estimated/cost-denominator qualifiers and renders only
remaining physical evidence gaps as BLOCKED panels.
"""

from __future__ import annotations

import csv
import hashlib
import json
import os
import shutil
import subprocess
import textwrap
import zipfile
from io import BytesIO
from dataclasses import dataclass
from pathlib import Path
from typing import Callable, Iterable

from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.pdfgen import canvas as pdfcanvas
from reportlab.lib.utils import ImageReader


ROOT = Path(__file__).resolve().parents[1]
PUB = ROOT / "build" / "publication_assets"
PRES = ROOT / "build" / "presentation"
SRC = ROOT / "data" / "publication"
RAW = ROOT / "results" / "experiments" / "scheduler_controlled.csv"
W, H = 3200, 1800
PAPER_W, PAPER_H = 3000, 1800
FONT_REGULAR = "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"
FONT_BOLD = "/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc"
FONT_MONO = "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf"

# Wong colorblind-safe palette, plus neutral status colors.
NAVY = "#16324F"
BLUE = "#0072B2"
SKY = "#56B4E9"
ORANGE = "#E69F00"
VERMILION = "#D55E00"
GREEN = "#009E73"
PALE_GREEN = "#BCE5D7"
PURPLE = "#CC79A7"
INK = "#17212B"
MID = "#5C6770"
LIGHT = "#EDF2F5"
GRID = "#CDD7DE"
WHITE = "#FFFFFF"
BLOCKED = "#6B7280"
RED = "#B3261E"


@dataclass(frozen=True)
class FigureSpec:
    id: str
    title_ko: str
    title_en: str
    question: str
    evidence: str
    claim: str
    blocked: str
    sources: tuple[str, ...]
    renderer: Callable[[Image.Image, ImageDraw.ImageDraw, bool], None]


def font(size: int, bold: bool = False, mono: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_MONO if mono else (FONT_BOLD if bold else FONT_REGULAR)
    return ImageFont.truetype(path, size)


def read_csv(path: Path) -> list[dict[str, str]]:
    with path.open(encoding="utf-8", newline="") as stream:
        return list(csv.DictReader(stream))


def write_csv(path: Path, rows: list[dict[str, object]], fields: list[str] | None = None) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    if not fields:
        fields = []
        for row in rows:
            for key in row:
                if key not in fields:
                    fields.append(key)
        fields = fields or ["status"]
    with path.open("w", encoding="utf-8", newline="") as stream:
        writer = csv.DictWriter(stream, fieldnames=fields, lineterminator="\n")
        writer.writeheader()
        writer.writerows(rows or [{"status": "empty"}])


def rounded(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], fill: str,
            outline: str = GRID, width: int = 4, radius: int = 28) -> None:
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def arrow(draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int],
          color: str = BLUE, width: int = 12, dash: bool = False) -> None:
    if dash:
        x1, y1 = start
        x2, y2 = end
        steps = 18
        for i in range(0, steps, 2):
            a, b = i / steps, min((i + 1) / steps, 1)
            draw.line((x1 + (x2-x1)*a, y1 + (y2-y1)*a,
                       x1 + (x2-x1)*b, y1 + (y2-y1)*b), fill=color, width=width)
    else:
        draw.line((*start, *end), fill=color, width=width)
    import math
    angle = math.atan2(end[1] - start[1], end[0] - start[0])
    for delta in (2.55, -2.55):
        draw.line((end[0], end[1],
                   end[0] + 34 * math.cos(angle + delta),
                   end[1] + 34 * math.sin(angle + delta)), fill=color, width=width)


def label(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, size: int = 34,
          fill: str = INK, bold: bool = False, anchor: str = "la", mono: bool = False) -> None:
    draw.text(xy, text, font=font(size, bold, mono), fill=fill, anchor=anchor)


def multiline(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], text: str,
              size: int = 30, fill: str = INK, bold: bool = False,
              spacing: int = 10, align: str = "left", min_size: int = 20) -> None:
    """Fit wrapped text inside a declared box or fail instead of overflowing."""
    x1, y1, x2, y2 = box
    if x2 <= x1 or y2 <= y1:
        raise ValueError(f"invalid text box: {box}")
    for candidate in range(size, min_size - 1, -1):
        face = font(candidate, bold)
        avg = max(6, int((x2 - x1) / (candidate * 0.54)))
        wrapped = "\n".join(
            textwrap.wrap(
                text,
                width=avg,
                break_long_words=True,
                break_on_hyphens=True,
            )
        )
        bounds = draw.multiline_textbbox(
            (0, 0), wrapped, font=face, spacing=spacing, align=align
        )
        text_width = bounds[2] - bounds[0]
        text_height = bounds[3] - bounds[1]
        if text_width <= x2 - x1 and text_height <= y2 - y1:
            draw_x = x1 if align == "left" else x1 + (x2 - x1 - text_width) / 2
            draw_y = y1 + (y2 - y1 - text_height) / 2 - bounds[1]
            draw.multiline_text(
                (draw_x, draw_y),
                wrapped,
                font=face,
                fill=fill,
                spacing=spacing,
                align=align,
            )
            return
    raise ValueError(f"text does not fit {box}: {text!r}")


def badge(draw: ImageDraw.ImageDraw, x: int, y: int, text: str,
          fill: str, width: int | None = None) -> None:
    width = width or max(220, int(len(text) * 21 + 60))
    rounded(draw, (x, y, x + width, y + 58), fill, outline=fill, radius=24)
    multiline(
        draw,
        (x + 14, y + 5, x + width - 14, y + 53),
        text,
        25,
        WHITE,
        True,
        2,
        "center",
        18,
    )


def header(draw: ImageDraw.ImageDraw, title: str, subtitle: str = "") -> None:
    draw.rectangle((0, 0, W, 170), fill=NAVY)
    label(draw, (110, 78), title, 52, WHITE, True, "lm")
    if subtitle:
        label(draw, (3090, 82), subtitle, 27, "#DDEAF2", False, "rm")


def evidence_legend(draw: ImageDraw.ImageDraw, y: int = 1690) -> None:
    items = [
        ("[R] RTL-generated", BLUE),
        ("[K] KiCad-native-checked", GREEN),
        ("[M] analytical-model", ORANGE),
        ("[X] BLOCKED / not-run", BLOCKED),
    ]
    x = 110
    for text, color in items:
        draw.rectangle((x, y, x + 32, y + 32), fill=color)
        label(draw, (x + 45, y + 16), text, 24, MID, False, "lm")
        x += 660


def blocked_panel(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int],
                  title: str, reason: str) -> None:
    rounded(draw, box, "#F3F4F6", outline=BLOCKED, width=7)
    x1, y1, x2, y2 = box
    badge(draw, x1 + 40, y1 + 35, "BLOCKED — EXCLUDED FROM CLAIMS", BLOCKED, 570)
    label(draw, ((x1+x2)//2, y1 + 175), title, 42, BLOCKED, True, "ma")
    multiline(draw, (x1 + 70, y1 + 235, x2 - 70, y2 - 50), reason,
              31, BLOCKED, False, 14, "center")


def arch_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Figure 1 · Current RTL planes and the missing system loop",
           "functional structure; not timing closure")
    rows = [
        ("COMPUTE PLANE", 310, BLUE,
         ["TileJob + payload", "TileScheduler", "payload store", "DecodeMatVecInt8", "result"]),
        ("MEMORY COMMAND PLANE", 760, ORANGE,
         ["external memoryRequest", "channel ingress", "bank-aware queue", "memoryCommands"]),
        ("LINK ROUTING PLANE", 1210, GREEN,
         ["external linkInput", "BundleRouter", "linkBundles"]),
    ]
    for row_name, y, color, names in rows:
        label(draw, (110, y + 115), row_name, 28, color, True, "lm")
        x0, span = 640, 2320
        width = 350
        gap = (span - width * len(names)) // max(1, len(names) - 1)
        for i, name in enumerate(names):
            x = x0 + i * (width + gap)
            rounded(draw, (x, y, x + width, y + 230), WHITE, color, 7)
            multiline(draw, (x + 25, y + 70, x + width - 25, y + 185),
                      name, 28, INK, True, 10, "center")
            if i < len(names) - 1:
                arrow(draw, (x + width + 6, y + 115),
                      (x + width + gap - 6, y + 115), color, 8)
    label(draw, (1650, 670),
          "NO response/data connection into the compute payload store",
          31, RED, True, "ma")
    draw.line((1650, 705, 1650, 760), fill=RED, width=8)
    draw.line((1650, 1090, 1650, 1210), fill=RED, width=8)
    badge(draw, 870, 1555, "END-TO-END MEMORY/LINK INTEGRATION: BLOCKED", RED, 1260)
    multiline(draw, (2200, 1548, 3060, 1648),
              "DDR PHY, GT wrappers, responses, and DMA sequencing are "
              "external/unimplemented.", 24, MID, False, 8)
    evidence_legend(draw)


def steal_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Figure 2 · Locality-aware work-stealing sequence",
           "RTL decision/dispatch; modeled migration cost")
    lanes = ["Idle cluster", "TileScheduler", "Victim queue", "MatVec"]
    lx = [320, 1050, 1800, 2680]
    for x, name in zip(lx, lanes):
        label(draw, (x, 260), name, 34, NAVY, True, "ma")
        draw.line((x, 320, x, 1510), fill=GRID, width=5)
    events = [
        (0, 1, 420, "local queue empty", BLUE),
        (1, 2, 560, "scan oldest eligible head", BLUE),
        (2, 1, 700, "job identity + owner/channel/bundle metadata", BLUE),
        (1, 1, 850, "age/locality score\nand steal decision", BLUE),
        (1, 3, 1245, "dispatch stolen TileJob to bounded MatVec harness", BLUE),
        (3, 1, 1430, "exact-once completion", PURPLE),
    ]
    for a, b, y, text, color in events:
        if a == b:
            box = (lx[a]-260, y-58, lx[a]+260, y+58)
            rounded(draw, box, LIGHT, color, 5, 18)
            multiline(draw, (box[0]+20, box[1]+10, box[2]-20, box[3]-10),
                      text, 27, color, True, 6, "center")
        else:
            arrow(draw, (lx[a], y), (lx[b], y), color, 8)
            left, right = sorted((lx[a], lx[b]))
            multiline(draw, (left+35, y-72, right-35, y-12),
                      text, 24, color, True, 4, "center")
    model_box = (610, 980, 1490, 1135)
    rounded(draw, model_box, "#FFF3D6", ORANGE, 6, 20)
    multiline(
        draw,
        (model_box[0]+30, model_box[1]+18, model_box[2]-30, model_box[3]-18),
        "Analytical only: remote bytes and migration latency.\n"
        "No integrated DMA, link, DDR request/response, or weight return.",
        26,
        ORANGE,
        True,
        7,
        "center",
    )
    badge(draw, 90, 1585, "RTL: scan → decision → MatVec completion", BLUE, 730)
    badge(draw, 860, 1585, "modeled migration cost", ORANGE, 430)
    label(draw, (1340, 1614),
          "No integrated remote-weight transaction or universal-optimum claim.",
          27, MID, False, "lm")
    evidence_legend(draw, 1718)


def gantt_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Figure 3 · Gemma 3 1B graph-derived replay and hybrid timing",
           "coarse analytical replay; not cycle trace")
    rows = read_csv(ROOT / "experiments/gemma3_1b/scheduler_replay.csv")
    rows = [r for r in rows if r["decode_tokens"] == "1" and r["scheduler"] in ("S1", "S3")]
    x0, x1 = 260, 1890
    max_p95 = max(float(r["p95_tile_latency_cycles"]) for r in rows)
    for idx, row in enumerate(rows):
        y = 410 + idx * 350
        label(draw, (130, y + 70), row["scheduler"], 50, NAVY, True, "mm")
        p50 = float(row["p50_tile_latency_cycles"])
        p95 = float(row["p95_tile_latency_cycles"])
        p99 = float(row["p99_tile_latency_cycles"])
        scale = (x1-x0) / (max_p95*1.18)
        draw.rectangle((x0, y, x0+p50*scale, y+58), fill=SKY)
        draw.rectangle((x0, y+78, x0+p95*scale, y+136), fill=ORANGE)
        draw.rectangle((x0, y+156, x0+p99*scale, y+214), fill=PURPLE)
        for yy, text, val in ((y+29, "p50", p50), (y+107, "p95", p95), (y+185, "p99", p99)):
            label(draw, (x0+12, yy), f"{text}  {val:,.0f} cycles", 25, INK, True, "lm")
        label(draw, (x0, y+260),
              f"modeled MAC duty={100*float(row['cluster_utilization_mean']):.2f}% · "
              f"unreserved idle={float(row['cluster_idle_cycles']):,.0f} cycles",
              27, MID)
    hybrid = read_csv(ROOT / "results/model_level/gemma3_1b_hybrid.csv")
    hybrid = [r for r in hybrid if r["scenario"] == "batch1_decode1"
              and r["scheduler"] in ("S0-physical", "S1", "S3", "Oracle")]
    box = (2050, 300, 3090, 1505)
    rounded(draw, box, WHITE, ORANGE, 7)
    label(draw, (box[0]+45, box[1]+70), "Hybrid decode-1 composition", 38, NAVY, True)
    max_total = max(float(r["hybrid_total_ms"]) for r in hybrid)
    for i, r in enumerate(hybrid):
        y = box[1]+190+i*230
        proj = float(r["projection_ms"])
        host = float(r["host_non_projection_fallback_ms"])
        scale = 780/max_total
        label(draw, (box[0]+55, y+38), r["scheduler"], 27, INK, True, "lm")
        x = box[0]+250
        draw.rectangle((x, y, x+proj*scale, y+75), fill=ORANGE)
        draw.rectangle((x+proj*scale, y, x+(proj+host)*scale, y+75), fill=SKY)
        label(draw, (x, y+105), f"{proj:.1f} ms modeled projection + "
              f"{host:.1f} ms measured host fallback", 23, MID)
    badge(draw, box[0]+55, box[3]-150, "hybrid-modeled", ORANGE, 330)
    label(draw, (box[0]+420, box[3]-120),
          "Oracle = offline clairvoyant bound; not implementable.", 23, RED)
    badge(draw, 190, 1430, "graph-derived + analytical replay", ORANGE, 600)
    evidence_legend(draw)


def results_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Figure 4 · Scheduler outcomes: latency, occupancy, traffic, energy, cost",
           "5 seeds × 3 deterministic repetitions")
    rows = read_csv(SRC / "k26_scheduler_summary.csv")
    scheduler_order = ("S0", "S1", "S2", "S3")
    rows = [
        r for r in rows
        if r["workload"] == "mixed" and r["scheduler"] in scheduler_order
    ]
    rows.sort(key=lambda row: scheduler_order.index(row["scheduler"]))
    panels = [(90, 240, 1570, 890), (1630, 240, 3110, 890),
              (90, 950, 1570, 1600), (1630, 950, 3110, 1600)]
    # (a) grouped latency bars.
    box = panels[0]
    rounded(draw, box, WHITE)
    label(draw, (box[0]+40, box[1]+55), "(a) Tail latency · mixed", 35, NAVY, True)
    maxv = max(float(r["p99_tile_latency_cycles_median"]) for r in rows)
    colors = [BLUE, ORANGE, PURPLE]
    metrics = ["p50_tile_latency_cycles_median", "p95_tile_latency_cycles_median",
               "p99_tile_latency_cycles_median"]
    for i, r in enumerate(rows):
        x = box[0]+120+i*335
        label(draw, (x+90, box[3]-45), r["scheduler"], 29, INK, True, "ma")
        for j, m in enumerate(metrics):
            v = float(r[m])
            h = 450*v/maxv
            draw.rectangle((x+j*62, box[3]-100-h, x+j*62+48, box[3]-100), fill=colors[j])
    for j, txt in enumerate(("p50", "p95", "p99")):
        badge(draw, box[0]+780+j*190, box[1]+35, txt, colors[j], 160)
    # (b) compute duty, reservation occupancy, and unreserved idle.
    box = panels[1]
    rounded(draw, box, WHITE)
    label(draw, (box[0]+40, box[1]+55),
          "(b) MAC duty / reservation occupancy · mixed", 34, NAVY, True)
    for i, r in enumerate(rows):
        y = box[1]+145+i*112
        compute = float(r["cluster_compute_utilization_mean_median"])
        occupancy = float(r["cluster_reservation_occupancy_mean_median"])
        label(draw, (box[0]+55, y+28), r["scheduler"], 28, INK, True, "lm")
        draw.rectangle((box[0]+180, y, box[0]+1130, y+56), fill=LIGHT)
        draw.rectangle(
            (box[0]+180, y, box[0]+180+950*occupancy, y+56),
            fill=PALE_GREEN,
        )
        draw.rectangle(
            (box[0]+180, y, box[0]+180+950*compute, y+56),
            fill=GREEN,
        )
        label(draw, (box[0]+1160, y+28),
              f"{compute*100:.2f}% / {occupancy*100:.2f}% / "
              f"{float(r['cluster_idle_cycles_median']):,.0f} idle",
              22, MID, False, "lm")
    # (c) traffic and steals.
    box = panels[2]
    rounded(draw, box, WHITE)
    label(draw, (box[0]+40, box[1]+55), "(c) Remote weights and successful steals", 35, NAVY, True)
    max_remote = max(float(r["remote_weight_bytes_median"]) for r in rows) or 1
    for i, r in enumerate(rows):
        y = box[1]+150+i*110
        remote = float(r["remote_weight_bytes_median"])
        steals = float(r["successful_steals_median"])
        label(draw, (box[0]+55, y+26), r["scheduler"], 28, INK, True, "lm")
        draw.rectangle((box[0]+180, y, box[0]+1050, y+52), fill=LIGHT)
        draw.rectangle((box[0]+180, y, box[0]+180+870*remote/max_remote, y+52), fill=VERMILION)
        label(draw, (box[0]+1080, y+26),
              f"{remote/1e6:.2f} MB · {steals:.0f} steals", 25, MID, False, "lm")
    # (d) estimated energy and DRAM-die-only cost arithmetic.
    box = panels[3]
    rounded(draw, box, WHITE, ORANGE, 6)
    label(draw, (box[0]+40, box[1]+55),
          "(d) Gemma decode-32 · S3 estimated energy/cost", 34, NAVY, True)
    energy = read_csv(ROOT/"results/power_cost/gemma3_1b_energy_join.csv")
    energy = [r for r in energy if r["scenario"] == "batch1_decode32"
              and r["scheduler"] == "S3"]
    energy.sort(key=lambda r: ("low", "central", "high").index(r["energy_case"]))
    colors_e = [GREEN, ORANGE, VERMILION]
    max_e = max(float(r["estimated_total_dynamic_j_per_token"]) for r in energy)
    for i, (r, color) in enumerate(zip(energy, colors_e)):
        y = box[1]+150+i*115
        v = float(r["estimated_total_dynamic_j_per_token"])
        label(draw, (box[0]+55, y+26), r["energy_case"], 27, INK, True, "lm")
        draw.rectangle((box[0]+220, y, box[0]+220+680*v/max_e, y+52), fill=color)
        label(draw, (box[0]+930, y+26), f"{v:.4f} estimated dynamic J/token",
              25, MID, False, "lm")
    cost = read_csv(ROOT/"cost/gemma3_1b_cost_normalized.csv")
    c = next(r for r in cost if r["scenario"] == "batch1_decode32"
             and r["scheduler"] == "S3"
             and r["price_case"] == "midpoint_sensitivity")
    multiline(draw, (box[0]+55, box[1]+480, box[2]-55, box[1]+535),
              f"{float(c['estimated_tokens_per_s']):.4f} hybrid tokens/s ÷ "
              f"${float(c['dram_die_cost_usd']):.2f} DRAM dies = "
              f"{float(c['estimated_tokens_per_s_per_dram_die_dollar']):.6f} tokens/s/$",
              27, NAVY, True, 5, "center")
    multiline(draw, (box[0]+55, box[1]+545, box[2]-55, box[3]-30),
              "Estimated dynamic energy excludes refresh, idle, controller, PHY, "
              "and board power. Cost denominator is DRAM dies only; FPGA, PCB, "
              "power delivery, cooling, assembly, and software are excluded.",
              22, RED, False, 6, "center", 17)
    badge(draw, 110, 1630, "analytical-model", ORANGE, 310)
    evidence_legend(draw, 1718)


def kicad_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Figure 5 · KiCad-native validation coupon",
           "fixed KiCad 3D render · NOT FOR FABRICATION")
    # KiCad's OpenGL renderer is not pixel-deterministic between invocations.
    # Reuse the reviewed native render tracked with the manuscript; native
    # ERC/DRC and PCB/source parity remain enforced by verify_k26_kicad.py.
    render = ROOT / "paper/final/figures/paper_f07_kicad_coupon_render.png"
    board = Image.open(render).convert("RGB")
    board = ImageOps.contain(board, (1960, 1260), Image.Resampling.LANCZOS)
    panel = (90, 240, 2110, 1530)
    rounded(draw, panel, "#E9EDF3", GREEN, 8)
    board_x = panel[0] + (panel[2] - panel[0] - board.width) // 2
    board_y = panel[1] + (panel[3] - panel[1] - board.height) // 2
    img.paste(board, (board_x, board_y))

    status = (2200, 245, 3100, 1530)
    rounded(draw, status, WHITE, GREEN, 8)
    badge(draw, 2260, 305, "NATIVE SOURCE RENDER", GREEN, 760)
    multiline(draw, (2260, 430, 3040, 610),
              "29 footprints · 116 nets · 65 tracks · 50 vias · 4 copper layers",
              34, NAVY, True, 14, "center")
    badge(draw, 2260, 670, "BOUNDED PASS", GREEN, 760)
    multiline(draw, (2260, 785, 3040, 960),
              "coupon ERC 0 · hierarchy ERC 0 · bounded routed-subset DRC 0",
              31, GREEN, True, 12, "center")
    badge(draw, 2260, 1020, "FABRICATION BLOCKERS", RED, 760)
    multiline(draw, (2260, 1130, 3040, 1435),
              "55 unrouted nets · 9 return-path crossings · "
              "no ground stitching · 0/116 test points · "
              "MPN/datasheet incomplete · no SI/PI closure",
              28, RED, True, 12, "center")
    evidence_legend(draw, 1718)


def onnx_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Figure 6 · ONNX graph-to-scheduler evidence flow",
           "direct graph inspection + separate ONNX Runtime reference")
    stages = [
        ("Authorized local artifact", "model.onnx + external data\nSHA-256 bound; weights excluded", PURPLE),
        ("ONNX protobuf inspection", "checker + graph load\nexternal tensors not materialized", GREEN),
        ("Graph inventory", "7,837 ordered nodes\nname · op · shape · initializer", GREEN),
        ("Projection filter", "183 dense nodes/token\n104 attention · 78 MLP · 1 head", GREEN),
        ("Deterministic ledger", "decode-1 / decode-32\nTileJob identity + byte geometry", ORANGE),
        ("Evidence branches", "analytical scheduler replay\n+ three bounded RTL weight tiles", BLUE),
    ]
    margin, gap = 90, 42
    width = (W - 2 * margin - gap * (len(stages) - 1)) // len(stages)
    for index, (name, detail, color) in enumerate(stages):
        x = margin + index * (width + gap)
        box = (x, 310, x + width, 1030)
        rounded(draw, box, WHITE, color, 7)
        badge(draw, x + 28, 355, f"STEP {index + 1}", color, width - 56)
        multiline(draw, (x + 28, 490, x + width - 28, 680),
                  name, 35, NAVY, True, 10, "center")
        multiline(draw, (x + 28, 735, x + width - 28, 955),
                  detail, 27, MID, False, 10, "center")
        if index < len(stages) - 1:
            arrow(draw, (x + width + 5, 670),
                  (x + width + gap - 5, 670), BLUE, 8)

    ort_box = (220, 1160, 2080, 1580)
    rounded(draw, ort_box, "#EEF5FB", BLUE, 7)
    badge(draw, 280, 1215, "ONNX RUNTIME · ANDROID CPU EP", BLUE, 740)
    multiline(draw, (280, 1320, 2020, 1515),
              "Functional reference only: batch 1 · sequence 1 · artificial past 1 · "
              "three pre-existing runs · 408.445 ms mean wall clock.",
              31, NAVY, True, 12, "center")

    boundary = (2160, 1160, 2980, 1580)
    rounded(draw, boundary, "#F3F4F6", BLOCKED, 7)
    badge(draw, 2220, 1215, "BOUNDARY", BLOCKED, 700)
    multiline(draw, (2220, 1320, 2920, 1515),
              "The ORT run did not replay decode-32 token_trace.jsonl and is not "
              "RTL, DRAMsim3, or accelerator timing.",
              29, BLOCKED, True, 12, "center")
    evidence_legend(draw, 1718)


def bottleneck_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Technical T01 · Exclusive wait-counter attribution",
           "analytical-model; medians")
    rows = read_csv(SRC / "k26_bottleneck_counters.csv")
    scheduler_order = ("S0", "S1", "S2", "S3")
    rows = [r for r in rows if r["scheduler"] in scheduler_order]
    rows.sort(key=lambda row: scheduler_order.index(row["scheduler"]))
    keys = ["request_serializer_wait_median", "bundle_contention_wait_median",
            "memory_command_wait_median"]
    names = ["serializer", "bundle contention", "memory command"]
    colors = [BLUE, ORANGE, GREEN]
    maxv = max(float(r[k]) for r in rows for k in keys)
    for i, r in enumerate(rows):
        y = 340+i*275
        label(draw, (135, y+72), r["scheduler"], 42, NAVY, True, "mm")
        x = 310
        for k, name, c in zip(keys, names, colors):
            v = float(r[k])
            width = int(2340*v/maxv)
            draw.rectangle((x, y, x+width, y+70), fill=c)
            label(draw, (x+width+20, y+35), f"{name}: {v:,.0f}", 25, INK, False, "lm")
            y += 78
        y -= 234
    badge(draw, 110, 1600, "exclusive counters; do not sum overlapping waits", ORANGE, 690)
    evidence_legend(draw, 1718)


def heatmap_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Technical T02 · Channel × link-width evidence coverage",
           "modeled cells and not-run cells")
    rows = read_csv(SRC / "k26_channel_link_matrix.csv")
    channels = [1, 2, 4]
    widths = [32, 64, 128, 256]
    lookup = {(int(r["channels"]), int(r["link_width_bits"])): r for r in rows}
    x0, y0, cw, ch = 600, 390, 560, 300
    for j, width in enumerate(widths):
        label(draw, (x0+j*cw+cw//2, 300), f"{width} bit", 36, NAVY, True, "ma")
    for i, channel in enumerate(channels):
        label(draw, (450, y0+i*ch+ch//2), f"{channel} channel", 36, NAVY, True, "rm")
        for j, width in enumerate(widths):
            r = lookup[(channel, width)]
            box = (x0+j*cw, y0+i*ch, x0+(j+1)*cw-20, y0+(i+1)*ch-20)
            modeled = r["status"] == "modeled"
            rounded(draw, box, "#FFF3D6" if modeled else "#F3F4F6",
                    ORANGE if modeled else BLOCKED, 6)
            if modeled:
                label(draw, ((box[0]+box[2])//2, box[1]+85),
                      f"{float(r['p95_cycles_median']):,.0f} p95 cycles",
                      28, INK, True, "ma")
                label(draw, ((box[0]+box[2])//2, box[1]+150),
                      r["source_subset"], 24, MID, False, "ma")
                badge(draw, box[0]+70, box[1]+205, "analytical-model", ORANGE, 360)
            else:
                label(draw, ((box[0]+box[2])//2, box[1]+120), "NOT RUN", 35, BLOCKED, True, "mm")
    label(draw, (600, 1435),
          "Blank matrix cells are preserved as not-run; interpolation is forbidden.",
          32, RED, True)
    evidence_legend(draw)


def evidence_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Technical T03 · Claim-to-evidence chain and physical gates",
           "what is known vs still required")
    stages = [
        ("K26 interface", "AMD datasheet", "verified", GREEN),
        ("SpinalHDL RTL", "generated + Verilator", "functional", BLUE),
        ("Scheduler outcomes", "780 rows", "analytical", ORANGE),
        ("KiCad scope", "ERC/limited DRC 0", "bounded", GREEN),
        ("Vivado/MIG", "tool + placement absent", "BLOCKED", BLOCKED),
        ("Gemma/energy/cost", "graph + hybrid + estimates", "modeled", ORANGE),
    ]
    for i, (name, source, status, color) in enumerate(stages):
        x = 100+i*510
        rounded(draw, (x, 440, x+430, 1170), WHITE, color, 8)
        badge(draw, x+35, 480, status, color, 360)
        multiline(draw, (x+35, 620, x+395, 850), name, 37, NAVY, True, 15, "center")
        multiline(draw, (x+35, 900, x+395, 1080), source, 28, MID, False, 12, "center")
        if i < len(stages)-1:
            arrow(draw, (x+432, 805), (x+505, 805), color if i < 3 else BLOCKED, 8, i >= 3)
    multiline(draw, (150, 1290, 3050, 1550),
              "Publication rule: only green/blue/orange nodes may support claims, "
              "and each claim retains its evidence-type qualifier. Gray nodes are "
              "future gates, never silently inferred from adjacent evidence.",
              36, INK, True, 18, "center")
    evidence_legend(draw)


def scaling_renderer(img: Image.Image, draw: ImageDraw.ImageDraw, bw: bool) -> None:
    header(draw, "Technical T04 · One-factor scaling without optimum extrapolation",
           "analytical-model")
    rows = read_csv(SRC / "k26_one_factor_summary.csv")
    rows = [r for r in rows if r["workload"] == "mixed"]
    factors = ["clusters", "channels", "bundles", "link_width_bits"]
    titles = ["Clusters", "Channels", "Bundles", "Link width"]
    panels = [(100, 250, 1550, 850), (1650, 250, 3100, 850),
              (100, 930, 1550, 1530), (1650, 930, 3100, 1530)]
    for factor_name, title, box in zip(factors, titles, panels):
        rounded(draw, box, WHITE)
        label(draw, (box[0]+45, box[1]+55), title, 34, NAVY, True)
        selected = [r for r in rows if r["factor"] == factor_name]
        selected.sort(key=lambda r: float(r["factor_value"]))
        if not selected:
            blocked_panel(draw, box, title, "No modeled rows in current one-factor source.")
            continue
        maxv = max(float(r["p95_tile_latency_cycles_median"]) for r in selected)
        for i, r in enumerate(selected):
            x = box[0]+160+i*300
            v = float(r["p95_tile_latency_cycles_median"])
            h = 390*v/maxv
            draw.rectangle((x, box[3]-100-h, x+150, box[3]-100), fill=ORANGE)
            label(draw, (x+75, box[3]-65), f"{float(r['factor_value']):g}", 27, INK, True, "ma")
            label(draw, (x+75, box[3]-120-h), f"{v/1000:.0f}k", 23, MID, False, "ma")
    badge(draw, 100, 1610, "No physical optimum or linear scaling claim", RED, 650)
    evidence_legend(draw, 1718)


def flow_renderer(flow_id: str, title: str, stages: list[tuple[str, str, str]]) -> Image.Image:
    img = Image.new("RGB", (W, H), WHITE)
    draw = ImageDraw.Draw(img)
    header(draw, f"{flow_id} · {title}", "16:9 flow asset")
    count = len(stages)
    margin, gap = 110, 70
    bw = (W-2*margin-(count-1)*gap)//count
    for i, (name, detail, evidence) in enumerate(stages):
        x = margin+i*(bw+gap)
        color = {"RTL": BLUE, "model": ORANGE, "checked": GREEN,
                 "blocked": BLOCKED, "input": PURPLE}.get(evidence, MID)
        rounded(draw, (x, 430, x+bw, 1250), WHITE, color, 8)
        badge(draw, x+30, 475, evidence, color, bw-60)
        multiline(draw, (x+35, 640, x+bw-35, 850), name, 37, NAVY, True, 16, "center")
        multiline(draw, (x+35, 910, x+bw-35, 1170), detail, 27, MID, False, 13, "center")
        if i < count-1:
            arrow(draw, (x+bw+5, 840), (x+bw+gap-5, 840),
                  BLOCKED if evidence == "blocked" else BLUE, 9,
                  evidence == "blocked")
    evidence_legend(draw)
    return img


def write_svg_wrapper(path: Path, png: Path, width: int, height: int) -> None:
    import base64
    data = base64.b64encode(png.read_bytes()).decode("ascii")
    path.write_text(
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}" '
        f'viewBox="0 0 {width} {height}"><image width="{width}" height="{height}" '
        f'href="data:image/png;base64,{data}"/></svg>\n', encoding="utf-8")


def save_pdf_images(images: list[Image.Image], path: Path, dpi: int = 150) -> None:
    """Write image pages with fixed PDF identifiers/timestamps."""
    first = images[0]
    page_size = (first.width * 72 / dpi, first.height * 72 / dpi)
    c = pdfcanvas.Canvas(str(path), pagesize=page_size, pageCompression=1, invariant=1)
    c.setTitle(path.stem)
    c.setAuthor("CHOI YUNHYUK")
    for im in images:
        buffer = BytesIO()
        im.convert("RGB").save(buffer, format="PNG", compress_level=9)
        buffer.seek(0)
        c.drawImage(ImageReader(buffer), 0, 0, width=page_size[0], height=page_size[1],
                    preserveAspectRatio=False, mask=None)
        c.showPage()
    c.save()


def normalize_zip(path: Path) -> None:
    """Rewrite an OOXML ZIP with stable entry order and timestamps."""
    tmp = path.with_suffix(path.suffix + ".stable")
    with zipfile.ZipFile(path, "r") as src, zipfile.ZipFile(
            tmp, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=9) as dst:
        for name in sorted(src.namelist()):
            old = src.getinfo(name)
            info = zipfile.ZipInfo(name, date_time=(2000, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.external_attr = old.external_attr
            info.create_system = 3
            dst.writestr(info, src.read(name), compress_type=zipfile.ZIP_DEFLATED,
                         compresslevel=9)
    tmp.replace(path)


def save_variants(img: Image.Image, out: Path) -> None:
    out.mkdir(parents=True, exist_ok=True)
    png = out / "figure.png"
    # 150 dpi is the declared placement density: 23 px = 11.04 pt.
    img.save(png, dpi=(150, 150), compress_level=9)
    save_pdf_images([img], out / "figure.pdf", 150)
    write_svg_wrapper(out / "figure.svg", png, W, H)
    # BW variant preserves differences through luminance plus hatch-like borders/text.
    gray = ImageOps.grayscale(img).convert("RGB")
    gray.save(out / "figure_bw.png", dpi=(150, 150), compress_level=9)
    img.save(out / "figure_16x9.png", dpi=(150, 150), compress_level=9)


def source_rows_for(spec: FigureSpec) -> list[dict[str, object]]:
    rows: list[dict[str, object]] = []
    for rel in spec.sources:
        path = ROOT / rel
        if path.suffix == ".csv":
            for row in read_csv(path):
                rows.append({"source_file": rel, **row})
        else:
            rows.append({"source_file": rel, "status": "binary evidence input"})
    return rows


def figure_metadata(spec: FigureSpec) -> str:
    return f"""# {spec.id} — {spec.title_en}

- ID: `{spec.id}`
- Korean title: {spec.title_ko}
- Paper location: core/technical publication asset
- Design question: {spec.question}
- Evidence type: {spec.evidence}
- Source data: {", ".join(f"`{s}`" for s in spec.sources)}
- Allowed claim: {spec.claim}
- Blocked/excluded: {spec.blocked}
- Generator: `python3 publication_tools/generate_publication_and_presentation.py`
- Variants: color SVG/PDF/PNG, grayscale PNG, 16:9 PNG
- Minimum designed text: 23 px at 3200×1800; validation checks 8 pt equivalent.

The SVG is a deterministic archival wrapper around the rendered PNG. It is not
claimed to be a fully editable vector reconstruction.
"""


def make_figures() -> list[FigureSpec]:
    specs = [
        FigureSpec("F01", "실제 RTL 전체 데이터 경로", "Actual RTL data path",
                   "Which blocks are implemented and which physical links remain conditional?",
                   "RTL-generated + datasheet candidate",
                   "Functional RTL integration and explicit clock/evidence boundaries.",
                   "Timing closure, MIG placement, measured bandwidth.",
                   ("data/publication/k26_system_architecture.csv",
                    "data/publication/k26_scheduler_policies.csv"), arch_renderer),
        FigureSpec("F02", "Work Stealing 실행 순서", "Work-stealing sequence",
                   "How does an idle cluster steal while preserving identity?",
                   "RTL decision/dispatch + analytical migration cost",
                   "Victim scan, locality scoring, dispatch, and exact-once MatVec completion.",
                   "Integrated DMA/link/DDR transfer, universal S3 optimality, or measured latency.",
                   ("data/publication/k26_scheduler_policies.csv",), steal_renderer),
        FigureSpec("F03", "Gemma graph-derived replay와 hybrid timing",
                   "Gemma graph-derived replay and hybrid timing",
                   "How do schedulers affect graph-derived projection replay and hybrid timing?",
                   "graph-derived + analytical replay + host-measured fallback",
                   "Decode-1 S1/S3 percentiles and hybrid S0-physical/S1/S3/Oracle composition.",
                   "Cycle-accurate event Gantt, hardware latency, implementable Oracle.",
                   ("experiments/gemma3_1b/scheduler_replay.csv",
                    "results/model_level/gemma3_1b_hybrid.csv"), gantt_renderer),
        FigureSpec("F04", "지연·점유율·traffic·추정 energy/cost",
                   "Latency, occupancy, traffic, and estimated energy/cost",
                   "Which modeled scheduler effects are supported?",
                   "analytical-model + graph-derived hybrid/energy/cost estimates",
                   "Mixed MAC-duty/reservation metrics plus decode-32 S3 energy range and DRAM-die-only cost arithmetic.",
                   "Measured board energy, whole-system price, product performance.",
                   ("data/publication/k26_scheduler_summary.csv",
                    "results/power_cost/gemma3_1b_energy_join.csv",
                    "cost/gemma3_1b_cost_normalized.csv"), results_renderer),
        FigureSpec("F05", "KiCad 참조 설계와 검증 쿠폰",
                   "KiCad reference design and validation coupon",
                   "What native design scope is checked?",
                   "KiCad-native-checked",
                   "Declared schematic/coupon ERC and bounded DRC status.",
                   "Fabrication readiness, SI/PI, full-board routing, hardware behavior.",
                   ("data/publication/k26_kicad_status.csv",
                    "hardware/kicad/k26_memory_reference/k26_memory_reference.kicad_sch",
                    "hardware/kicad/k26_memory_coupon/k26_memory_coupon.kicad_sch",
                    "hardware/kicad/k26_memory_coupon/k26_memory_coupon.kicad_pcb",
                    "paper/final/figures/paper_f07_kicad_coupon_render.png"),
                   kicad_renderer),
        FigureSpec("F06", "ONNX graph와 Runtime 경계",
                   "ONNX graph and Runtime evidence boundary",
                   "How does the on-device graph become a scheduler ledger, and what did ONNX Runtime measure?",
                   "graph-derived + pre-existing ONNX Runtime host measurement",
                   "Ordered graph inventory, dense-projection filter, deterministic ledger, and the bounded Android CPU EP reference.",
                   "Decode-32 ORT replay, accelerator timing, full-model RTL inference, or hardware latency.",
                   ("experiments/gemma3_1b/projection_trace.csv",
                    "experiments/gemma3_1b/trace_manifest.json"),
                   onnx_renderer),
        FigureSpec("T01", "배타적 wait counter", "Exclusive wait-counter attribution",
                   "Which modeled waits dominate by scheduler?",
                   "analytical-model", "Median exclusive counter values.",
                   "Physical cycle attribution or overlapping-wait summation.",
                   ("data/publication/k26_bottleneck_counters.csv",), bottleneck_renderer),
        FigureSpec("T02", "Channel×link evidence coverage", "Channel × link evidence coverage",
                   "Which combinations were actually modeled?",
                   "analytical-model + not-run",
                   "Only populated cells and explicit not-run coverage.",
                   "Interpolation into blank cells.",
                   ("data/publication/k26_channel_link_matrix.csv",), heatmap_renderer),
        FigureSpec("T03", "주장-증거 chain", "Claim-to-evidence chain",
                   "Where do claims stop?",
                   "mixed evidence",
                   "Evidence-layer boundaries and future gates.",
                   "Vivado/P&R, implemented MIG/link timing, board power, SI/PI, or fabrication closure.",
                   ("data/publication/k26_research_flow.csv",
                    "data/publication/k26_design_decisions.csv"), evidence_renderer),
        FigureSpec("T04", "One-factor scaling", "One-factor scaling",
                   "How do modeled factors change p95?",
                   "analytical-model", "One-factor p95 medians.",
                   "Physical optimum, linear extrapolation, unmodeled combinations.",
                   ("data/publication/k26_one_factor_summary.csv",), scaling_renderer),
    ]
    for spec in specs:
        out = PUB / "figures" / spec.id
        img = Image.new("RGB", (W, H), WHITE)
        draw = ImageDraw.Draw(img)
        spec.renderer(img, draw, False)
        save_variants(img, out)
        rows = source_rows_for(spec)
        write_csv(out / "source.csv", rows)
        (out / "metadata.md").write_text(figure_metadata(spec), encoding="utf-8")
        (out / "generate.py").write_text(
            "#!/usr/bin/env python3\n"
            "from pathlib import Path\nimport subprocess, sys\n"
            "root = Path(__file__).resolve().parents[3]\n"
            "raise SystemExit(subprocess.call([sys.executable, "
            "str(root/'publication_tools/generate_publication_and_presentation.py')], cwd=root))\n",
            encoding="utf-8")
    return specs


def make_flows() -> list[Path]:
    definitions = [
        ("F01_research_decision_flow", "Research decision flow",
         [("Gemma workload", "183 graph-derived projections", "input"),
          ("K26 target", "four exposed GTH lanes", "checked"),
          ("RTL prototype", "real MatVec + scheduler", "RTL"),
          ("Analytical sweep", "780 controlled rows", "model"),
          ("Physical gates", "Vivado/MIG/board", "blocked")]),
        ("F02_token_dataflow", "Token / model data flow",
         [("ONNX graph", "node order/shape/weights", "input"),
          ("TileJob", "identity + shape", "RTL"),
          ("Multi-queue", "home assignment", "RTL"),
          ("Weight fetch", "bundle/channel model", "model"),
          ("MatVec result", "exact-once RTL", "RTL")]),
        ("F03_work_stealing_sequence", "Work-stealing sequence",
         [("Imbalance", "idle local queue", "RTL"),
          ("Victim scan", "oldest eligible", "RTL"),
          ("Locality score", "owner/channel/bundle", "RTL"),
          ("Migration cost", "bytes/latency; DMA unintegrated", "model"),
          ("Completion", "exact-once identity", "RTL")]),
        ("F04_cycle_timeline", "Cycle timeline",
         [("Queue wait", "modeled aggregate", "model"),
          ("Link wait", "exclusive counters", "model"),
          ("Memory wait", "command counter", "model"),
          ("Compute", "real MatVec function", "RTL"),
          ("Gemma hybrid", "modeled + host fallback", "model")]),
        ("F05_physical_data_path", "Physical data path",
         [("K26", "compute SOM", "checked"),
          ("4× GTH", "datasheet candidate", "checked"),
          ("XC7K160T", "conditional fallback", "model"),
          ("4× DDR3L", "8 GB arithmetic", "model"),
          ("Board sign-off", "SI/PI + bring-up", "blocked")]),
    ]
    paths = []
    editable = PUB / "flow" / "editable_sources"
    editable.mkdir(parents=True, exist_ok=True)
    for stem, title, stages in definitions:
        img = flow_renderer(stem.split("_")[0], title, stages)
        png = PUB / "flow" / f"{stem}.png"
        pdf = PUB / "flow" / f"{stem}.pdf"
        svg = PUB / "flow" / f"{stem}.svg"
        img.save(png, dpi=(150, 150), compress_level=9)
        save_pdf_images([img], pdf, 150)
        write_svg_wrapper(svg, png, W, H)
        (editable / f"{stem}.json").write_text(
            json.dumps({"id": stem, "title": title, "stages": stages},
                       ensure_ascii=False, indent=2, sort_keys=True) + "\n",
            encoding="utf-8")
        paths.append(png)
    return paths


def make_storyboard_and_animation() -> None:
    frames_dir = PUB / "flow" / "animation_frames"
    frames_dir.mkdir(parents=True, exist_ok=True)
    labels = [
        ("1", "Queue imbalance", "Cluster 3 becomes idle while another local queue remains deep."),
        ("2", "Victim search", "The scheduler scans oldest eligible jobs."),
        ("3", "Locality score", "Age is balanced against owner/channel/bundle locality."),
        ("4", "Steal decision", "A single job identity migrates to the idle cluster."),
        ("5", "Modeled migration cost",
         "Remote bytes and latency are charged analytically; DMA/link/DDR is not integrated."),
        ("6", "MatVec execute", "The real DecodeMatVecInt8 datapath consumes the tile."),
        ("7", "Exact-once completion", "The live scoreboard retires one matching identity."),
        ("8", "Modeled tail outcome", "Aggregate p95 may improve versus static local FCFS."),
    ]
    frames = []
    for i, (n, title, body) in enumerate(labels):
        img = Image.new("RGB", (1920, 1080), NAVY)
        d = ImageDraw.Draw(img)
        label(d, (120, 100), "Locality-aware Work Stealing", 54, WHITE, True)
        label(d, (1800, 100), f"{n}/8", 38, "#DDEAF2", True, "rm")
        rounded(d, (160, 250, 1760, 840), WHITE, WHITE, 2, 36)
        step_color = ORANGE if i in (4, 7) else (PURPLE if i == 6 else BLUE)
        badge(d, 230, 320, f"STEP {n}", step_color, 250)
        label(d, (960, 500), title, 62, NAVY, True, "mm")
        multiline(d, (300, 610, 1620, 760), body, 34, MID, False, 16, "center")
        label(d, (960, 980),
              "Illustrative sequence · quantitative claims require linked CSV evidence",
              28, "#DDEAF2", False, "ma")
        frame = frames_dir / f"frame_{i:03d}.png"
        img.save(frame, compress_level=9)
        frames.append(frame)
    frames[0].parent.mkdir(exist_ok=True)
    imgs = [Image.open(p).convert("RGB") for p in frames]
    save_pdf_images(imgs, PUB/"flow/storyboard_16x9.pdf", 150)
    # Two seconds per frame, deterministic metadata-free streams.
    if shutil.which("ffmpeg"):
        subprocess.run([
            "ffmpeg", "-y", "-loglevel", "error", "-framerate", "1/2",
            "-i", str(frames_dir/"frame_%03d.png"), "-vf", "format=yuv420p",
            "-c:v", "libx264", "-preset", "slow", "-crf", "20",
            "-map_metadata", "-1", "-movflags", "+faststart",
            str(PUB/"flow/work_stealing_animation.mp4")], check=True)
        subprocess.run([
            "ffmpeg", "-y", "-loglevel", "error", "-framerate", "1/2",
            "-i", str(frames_dir/"frame_%03d.png"), "-vf",
            "fps=8,scale=960:-1:flags=lanczos,split[s0][s1];"
            "[s0]palettegen=max_colors=128[p];[s1][p]paletteuse=dither=bayer",
            "-loop", "0", "-map_metadata", "-1",
            str(PUB/"flow/work_stealing_animation.gif")], check=True)
    (PUB/"flow/README.md").write_text(
        "# Flow and animation assets\n\n"
        "F01–F06 are 3200×1800 PNG/PDF/SVG assets. The storyboard contains eight "
        "16:9 pages. MP4/GIF are illustrative policy sequences, not waveform or "
        "Gemma trace evidence. Animation frames are retained for deterministic rebuilds.\n",
        encoding="utf-8")


def captions(specs: list[FigureSpec]) -> None:
    ko = ["# 그림 캡션\n"]
    en = ["# Figure captions\n"]
    for s in specs:
        ko.append(
            f"**{s.id}. {s.title_ko}.** 증거 유형: {s.evidence}. "
            f"허용 주장: {s.claim} 제외: {s.blocked}\n")
        en.append(
            f"**{s.id}. {s.title_en}.** Evidence type: {s.evidence}. "
            f"Allowed: {s.claim} Excluded: {s.blocked}\n")
    (PUB/"captions_ko.md").write_text("\n".join(ko), encoding="utf-8")
    (PUB/"captions_en.md").write_text("\n".join(en), encoding="utf-8")


def make_evidence_map(specs: list[FigureSpec]) -> None:
    rows = []
    for s in specs:
        rows.append({
            "asset_id": s.id,
            "title": s.title_en,
            "evidence_type": s.evidence,
            "source_paths": ";".join(s.sources),
            "allowed_claim": s.claim,
            "excluded_or_blocked": s.blocked,
            "generator": "publication_tools/generate_publication_and_presentation.py",
        })
    write_csv(PUB/"evidence_map.csv", rows)


def make_tables_and_images() -> None:
    """Create bounded, source-linked table/image inventories."""
    scheduler = read_csv(SRC/"k26_scheduler_summary.csv")
    selected = [
        r for r in scheduler
        if r["workload"] in ("mixed", "skew") and r["scheduler"] in ("S0", "S1", "S2", "S3")
    ]
    fields = [
        "workload", "scheduler", "evidence_type", "seed_count", "process_repetitions",
        "p50_tile_latency_cycles_median", "p95_tile_latency_cycles_median",
        "p99_tile_latency_cycles_median",
        "cluster_compute_utilization_mean_median",
        "cluster_reservation_occupancy_mean_median",
        "cluster_idle_cycles_median", "successful_steals_median",
        "remote_weight_bytes_median",
    ]
    write_csv(PUB/"tables/scheduler_core_metrics.csv",
              [{k: r[k] for k in fields} for r in selected], fields)
    gaps = [
        {"required_metric": "post-route performance/power", "status": "BLOCKED",
         "reason": "Vivado synthesis/P&R/timing/power absent", "claim_use": "excluded"},
        {"required_metric": "measured whole-board energy", "status": "BLOCKED",
         "reason": "no fabricated board or calibrated power measurement", "claim_use": "excluded"},
    ]
    write_csv(PUB/"tables/blocked_evidence.csv", gaps)
    energy = read_csv(ROOT/"results/power_cost/gemma3_1b_energy_join.csv")
    energy = [r for r in energy if r["scenario"] == "batch1_decode32"
              and r["scheduler"] in ("S0-physical", "S3", "Oracle")]
    cost = read_csv(ROOT/"cost/gemma3_1b_cost_normalized.csv")
    cost = [r for r in cost if r["scenario"] == "batch1_decode32"
            and r["scheduler"] in ("S0-physical", "S3", "Oracle")
            and r["price_case"] == "midpoint_sensitivity"]
    write_csv(PUB/"tables/gemma_energy_sensitivity.csv",
              [{"source_file": "results/power_cost/gemma3_1b_energy_join.csv", **r}
               for r in energy])
    write_csv(PUB/"tables/gemma_cost_midpoint.csv",
              [{"source_file": "cost/gemma3_1b_cost_normalized.csv", **r}
               for r in cost])
    (PUB/"tables/README.md").write_text(
        "# Tables\n\n`scheduler_core_metrics.csv` is a bounded projection of the committed "
        "`k26_scheduler_summary.csv`. Gemma energy and cost tables preserve the "
        "analytical/hybrid and DRAM-die-only qualifiers. `blocked_evidence.csv` "
        "records only remaining physical gates.\n", encoding="utf-8")
    (PUB/"images/README.md").parent.mkdir(parents=True, exist_ok=True)
    (PUB/"images/README.md").write_text(
        "# Images\n\nKiCad pages are rasterized during figure generation from the committed "
        "NOT FOR FABRICATION PDF exports. No external or generated photorealistic board "
        "image is used.\n", encoding="utf-8")


def make_contact_sheets(specs: list[FigureSpec], flow_paths: list[Path]) -> None:
    def sheet(paths: list[Path], target: Path, title: str) -> None:
        thumb_w, thumb_h = 1200, 675
        cols = 2
        rows = (len(paths)+1)//2
        canvas = Image.new("RGB", (2600, 160 + rows*790), WHITE)
        d = ImageDraw.Draw(canvas)
        label(d, (1300, 75), title, 48, NAVY, True, "ma")
        for i, p in enumerate(paths):
            im = Image.open(p).convert("RGB")
            im = ImageOps.contain(im, (thumb_w, thumb_h))
            x = 80+(i%cols)*1260
            y = 150+(i//cols)*790
            canvas.paste(im, (x+(thumb_w-im.width)//2, y))
            label(d, (x+thumb_w//2, y+710), p.parent.name if p.name == "figure.png" else p.stem,
                  28, INK, True, "ma")
        target.parent.mkdir(parents=True, exist_ok=True)
        canvas.save(target, dpi=(150, 150))
        save_pdf_images([canvas], target.with_suffix(".pdf"), 150)
    fig_paths = [PUB/"figures"/s.id/"figure.png" for s in specs]
    sheet(fig_paths, PUB/"contact_sheets/figures_contact_sheet.png", "Publication figures")
    sheet(flow_paths, PUB/"contact_sheets/flows_contact_sheet.png", "Flow assets")


def slide_image(slide, path: Path) -> None:
    slide.shapes.add_picture(str(path), 0, 0, width=Inches(13.333), height=Inches(7.5))


def set_slide_notes(slide, note: str) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.text = note


def title_slide_image(title: str, subtitle: str, badge_text: str = "") -> Image.Image:
    img = Image.new("RGB", (1920, 1080), NAVY)
    d = ImageDraw.Draw(img)
    draw = d
    label(draw, (130, 120), "VARP · K26–Memory FPGA", 34, "#A9D6EC", True)
    multiline(draw, (130, 270, 1790, 650), title, 67, WHITE, True, 24)
    multiline(draw, (135, 700, 1700, 850), subtitle, 34, "#DDEAF2", False, 16)
    if badge_text:
        badge(draw, 135, 900, badge_text, ORANGE, 650)
    return img


def make_presentation(specs: list[FigureSpec]) -> None:
    assets = PRES/"assets"
    assets.mkdir(parents=True, exist_ok=True)
    slides: list[tuple[str, str, Image.Image | Path, str]] = []
    slides.append(("문제와 연구 질문", "Can a K26 + memory-FPGA architecture reduce "
                   "imbalance without hiding link/memory cost?",
                   title_slide_image("K26–Memory FPGA 다채널 가속기와\nlocality-aware Work Stealing",
                                     "독립적인 시각 증거 패키지 · 분석 모델과 물리 검증 경계를 분리",
                                     "PASS WITH EXPLICIT LIMITATIONS"),
                   "연구 질문을 구조, 스케줄링, 물리 증거의 세 층으로 소개한다. "
                   "분석 모델 결과를 실측 성능처럼 말하지 않는다."))
    slides.append(("Gemma 3 1B workload",
                   "ONNX graph inspection and the separate Android CPU EP reference are explicitly bounded.",
                   PUB/"figures/F06/figure_16x9.png",
                   "ONNX protobuf inspection이 7,837-node graph와 183 projection ledger를 만든다. "
                   "별도의 ONNX Runtime Android CPU EP 408.445 ms 참조는 decode-32 ledger를 "
                   "실행하지 않았으며 accelerator timing으로 해석하지 않는다."))
    mapping = [
        ("K26–Memory FPGA 전체 구조", "Implemented RTL is separated from conditional physical candidates.", "F01"),
        ("단일/정적 queue의 문제", "Graph-derived projection replay exposes scheduler tradeoffs.", "F03"),
        ("Multi-Queue FCFS", "S1 keeps home-local FCFS queues; it avoids migration but can strand work.", "F02"),
        ("Work Stealing 순서", "S3 scans eligibility and scores age/locality before exact-once completion.", "F02"),
        ("S1/S3 timing comparison", "Gemma decode-1 projection replay and hybrid composition are modeled.", "F03"),
        ("시간·이용률 결과", "S3 improves S1 under mixed/skew, but is not universally best.", "F04"),
        ("traffic·energy·cost 결과", "Dynamic energy is estimated; cost divides hybrid throughput by DRAM-die dollars only.", "F04"),
        ("KiCad 참조 설계", "Native ERC/limited DRC pass only within the declared reference/coupon scope.", "F05"),
        ("핵심 기여와 한계", "Functional integration and bounded evidence are complete; physical and model gates remain.", "T03"),
        ("다음 물리 검증", "Vivado/MIG, board SI/PI, post-route power, and measured whole-board energy.", "T03"),
    ]
    notes = {
        "F01": "실제 Scala module 이름과 valid/ready 경로를 따라간다. DDR3L과 GTH는 조건부 후보다.",
        "F02": "idle → victim scan → locality score → fetch → MatVec → exact-once 순서다.",
        "F03": "graph-derived projection replay와 measured host fallback의 hybrid 조합이다. Oracle은 구현 불가능한 offline bound다.",
        "F04": "energy는 dynamic estimate, cost는 DRAM-die-only denominator다. 보드 전력·전체 시스템 가격이 아니다.",
        "F05": "ERC/DRC 0을 fabrication ready로 확대 해석하지 않는다.",
        "F06": "ONNX graph inspection과 ONNX Runtime functional reference를 같은 실행 경로로 합치지 않는다.",
        "T03": "지원된 주장과 다음 gate를 색으로 구분한다.",
    }
    for title, message, fid in mapping:
        slides.append((title, message, PUB/"figures"/fid/"figure_16x9.png", notes[fid]))
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    page_imgs = []
    for idx, (title, message, art, note) in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank)
        if isinstance(art, Image.Image):
            path = assets/f"slide_{idx:02d}.png"
            art.save(path, compress_level=9)
        else:
            base = Image.open(art).convert("RGB")
            path = assets/f"slide_{idx:02d}.png"
            base.resize((1920, 1080), Image.Resampling.LANCZOS).save(path, compress_level=9)
        # Add a subtle message strip for figure slides.
        im = Image.open(path).convert("RGB")
        d = ImageDraw.Draw(im)
        if not isinstance(art, Image.Image):
            focus = {
                "단일/정적 queue의 문제": "FOCUS · imbalance premise",
                "Multi-Queue FCFS": "FOCUS · S1 static locality",
                "Work Stealing 순서": "FOCUS · S3 decision boundary",
                "S1/S3 timing comparison": "FOCUS · modeled timing",
                "시간·이용률 결과": "FOCUS · latency + real MAC duty",
                "traffic·energy·cost 결과": "FOCUS · denominator + exclusions",
                "핵심 기여와 한계": "FOCUS · supported claims",
                "다음 물리 검증": "FOCUS · blocked gates",
            }.get(title, f"FOCUS · {title}")
            rounded(d, (1340, 180, 1860, 245), WHITE, ORANGE, 3, 18)
            multiline(d, (1360, 190, 1840, 235), focus, 23,
                      ORANGE, True, 2, "center", 18)
        d.rectangle((0, 1000, 1920, 1080), fill=NAVY)
        multiline(d, (60, 1010, 1740, 1070), message, 25,
                  WHITE, False, 3, "left", 19)
        label(d, (1870, 1040), f"{idx:02d}/12", 23, "#DDEAF2", True, "rm")
        im.save(path, compress_level=9)
        slide_image(slide, path)
        set_slide_notes(slide, f"{title}\n\n{note}\n\nEvidence discipline: {message}")
        page_imgs.append(Image.open(path).convert("RGB"))
    # Remove the automatically created first blank slide if present.
    if len(prs.slides) > 12:
        slide_id = prs.slides._sldIdLst[0]
        prs.part.drop_rel(slide_id.rId)
        del prs.slides._sldIdLst[0]
    core = prs.core_properties
    core.title = "VARP K26–Memory FPGA Paper 10 Presentation"
    core.subject = "Evidence-bounded 12-slide technical presentation"
    core.author = "VARP research workflow"
    core.keywords = "K26, FPGA, work stealing, evidence"
    core.created = core.modified = __import__("datetime").datetime(2000, 1, 1)
    prs.save(PRES/"presentation.pptx")
    normalize_zip(PRES/"presentation.pptx")
    save_pdf_images(page_imgs, PRES/"presentation.pdf", 150)
    speaker = ["# Speaker notes\n"]
    for i, (title, message, _, note) in enumerate(slides, 1):
        speaker.append(f"## {i}. {title}\n\n{note}\n\nEvidence boundary: {message}\n")
    (PRES/"speaker_notes.md").write_text("\n".join(speaker), encoding="utf-8")
    (PRES/"qa_evidence_index.md").write_text("""# 예상 Q&A와 증거 위치

1. **S3가 항상 가장 빠른가?** 아니다. `publication_assets/figures/F04/source.csv`에서 S0–S3를 함께 확인한다.
2. **Gemma 3 1B를 실제로 실행했는가?** graph order/shape/initializer는 ONNX-derived이고 projection은 analytical replay, non-projection은 measured Y700 fallback이다. 실물 end-to-end 측정은 아니다.
3. **에너지와 비용 이득은 얼마인가?** S3 decode-32 dynamic energy는 0.2919–1.1850 estimated J/token이고 중앙값은 0.5538이다. 이 값은 base stream과 stealing overhead를 포함한다. 2.4835 hybrid tokens/s를 midpoint DRAM die $545.68로 나누면 0.004551 tokens/s/$다. 보드 전력과 전체 시스템 가격은 제외한다.
4. **보드를 제작할 수 있는가?** 아니다. KiCad 범위는 reference/coupon이며 NOT FOR FABRICATION이다.
5. **정확성은 무엇으로 보장했는가?** 실제 MatVec 통합, live scoreboard, exact-once RTL 테스트가 구조적 근거다.
6. **4채널 DDR3L은 확정인가?** 용량 산술과 모델 후보이며 4-MIG placement/timing은 Vivado gate가 남았다.
7. **그림의 숫자 출처는?** `publication_assets/evidence_map.csv`가 각 그림을 원본 CSV로 연결한다.
8. **빈 heatmap cell은?** `not_run`이며 보간하지 않는다.
""", encoding="utf-8")
    (PRES/"demo_script.md").write_text("""# 60초 데모 스크립트

0–10초: 슬라이드 3에서 RTL 블록과 조건부 물리 블록의 색을 구분한다.

10–25초: `publication_assets/flow/work_stealing_animation.mp4`를 재생해 idle cluster가 victim을 고르고, migration cost를 분석 모델로 부과한 뒤 functional MatVec harness에서 exact-once로 완료하는 순서를 보여준다. DMA/link/DDR 전송은 미통합이다.

25–40초: 슬라이드 8에서 mixed workload의 p50/p95/p99, MAC duty와 reservation occupancy, remote traffic을 비교한다.

40–50초: 슬라이드 10에서 ERC/limited DRC 0 badge와 NOT FOR FABRICATION 문구를 함께 지목한다.

50–60초: 슬라이드 12에서 Vivado/MIG, full-token hardware execution, whole-board energy calibration, procurement-time price refresh를 다음 gate로 제시한다.

실행:

```bash
ffplay -autoexit publication_assets/flow/work_stealing_animation.mp4
```
""", encoding="utf-8")


def make_index(specs: list[FigureSpec]) -> None:
    lines = [
        "# VARP K26–Memory FPGA Publication Assets",
        "",
        "This package is generated from committed CSV evidence and bounded KiCad exports.",
        "Orange means analytical model; blue means generated/functional RTL; green means "
        "native KiCad or datasheet check; gray means BLOCKED/not-run.",
        "",
        "## Core paper figures",
        "",
        "| ID | Title | Evidence | Claim boundary |",
        "|---|---|---|---|",
    ]
    for s in specs[:6]:
        lines.append(f"| [{s.id}](figures/{s.id}/figure.png) | {s.title_en} | "
                     f"{s.evidence} | {s.blocked} |")
    lines += ["", "## Technical figures", ""]
    for s in specs[5:]:
        lines.append(f"- [{s.id} — {s.title_en}](figures/{s.id}/figure.png): {s.evidence}")
    lines += [
        "", "## Flow and animation",
        "",
        "- F01 research decision flow",
        "- F02 token/model data flow",
        "- F03 work-stealing sequence",
        "- F04 cycle timeline",
        "- F05 physical data path",
        "- F06 ONNX graph and Runtime evidence boundary",
        "- 8-page 16:9 storyboard",
        "- MP4/GIF illustrative animation",
        "",
        "## Gemma 3 1B modeled evidence",
        "",
        "- 183 graph-derived projection nodes per token",
        "- graph-derived projection replay plus measured Y700 host fallback (hybrid-modeled)",
        "- decode-32 S3: 2.4835 hybrid tokens/s",
        "- S3 estimated dynamic energy sensitivity: 0.2919 / 0.5538 / 1.1850 J/token",
        "- midpoint $545.68 DRAM-die-only denominator: 0.004551 tokens/s/$",
        "",
        "## Explicit blocked evidence",
        "",
        "- cycle-accurate Gemma per-stage event Gantt and end-to-end hardware measurement",
        "- calibrated whole-board energy including refresh/idle/controller/PHY",
        "- whole-system price including FPGA/PCB/power/cooling/assembly/software",
        "- Vivado 4-MIG placement, synthesis, timing closure, and board measurement",
        "",
        "These gaps are rendered visibly and excluded from quantitative claims.",
    ]
    (PUB/"INDEX.md").write_text("\n".join(lines)+"\n", encoding="utf-8")
    # PDF index rendered from deterministic image pages.
    page = Image.new("RGB", (2480, 3508), WHITE)
    d = ImageDraw.Draw(page)
    label(d, (160, 180), "VARP Publication Assets Index", 72, NAVY, True)
    label(d, (160, 280), "Evidence-bounded visual package", 38, MID)
    y = 430
    for s in specs:
        badge(d, 160, y, s.id, BLUE if s.id.startswith("F") else PURPLE, 180)
        label(d, (390, y+28), s.title_en, 34, INK, True, "lm")
        label(d, (390, y+78), s.evidence, 25, MID, False, "lm")
        y += 240
    blocked_panel(d, (160, y+40, 2320, y+600), "Blocked evidence",
                  "full-token hardware execution · whole-board calibrated energy · "
                  "procurement-time price refresh · Vivado/MIG timing · physical board measurement")
    save_pdf_images([page], PUB/"INDEX.pdf", 150)


def validate(specs: list[FigureSpec]) -> dict[str, object]:
    errors = []
    for s in specs:
        out = PUB/"figures"/s.id
        for name in ("figure.png", "figure.pdf", "figure.svg", "figure_bw.png",
                     "figure_16x9.png", "source.csv", "generate.py", "metadata.md"):
            if not (out/name).is_file() or (out/name).stat().st_size == 0:
                errors.append(f"{s.id}: missing {name}")
        if (out/"figure.png").is_file():
            im = Image.open(out/"figure.png")
            if im.size != (W, H):
                errors.append(f"{s.id}: wrong dimension {im.size}")
    for required in (
        PUB/"INDEX.md", PUB/"INDEX.pdf", PUB/"captions_ko.md", PUB/"captions_en.md",
        PUB/"evidence_map.csv", PUB/"flow/storyboard_16x9.pdf",
        PRES/"presentation.pptx", PRES/"presentation.pdf",
        PRES/"speaker_notes.md", PRES/"qa_evidence_index.md", PRES/"demo_script.md",
    ):
        if not required.is_file() or required.stat().st_size == 0:
            errors.append(f"missing {required.relative_to(ROOT)}")
    prs = Presentation(PRES/"presentation.pptx")
    if len(prs.slides) != 12:
        errors.append(f"PPTX slide count {len(prs.slides)} != 12")
    if prs.slide_width / prs.slide_height < 1.77:
        errors.append("PPTX is not 16:9")
    if not all(slide.notes_slide.notes_text_frame.text.strip() for slide in prs.slides):
        errors.append("one or more PPTX notes are empty")
    blocked_text = (PUB/"INDEX.md").read_text(encoding="utf-8")
    if "Gemma 3 1B" not in blocked_text or "hybrid" not in blocked_text:
        errors.append("Gemma hybrid evidence boundary is not explicit")
    if "whole-board energy" not in blocked_text or "BLOCKED" not in blocked_text:
        errors.append("remaining physical evidence gap is not explicit")
    report = {
        "status": "PASS" if not errors else "FAIL",
        "figure_count": len(specs),
        "core_figure_count": 6,
        "flow_count": 5,
        "presentation_slides": len(prs.slides),
        "canvas_pixels": [W, H],
        "minimum_designed_text_px": 23,
        "declared_output_dpi": 150,
        "minimum_equivalent_pt_at_declared_dpi": round(23*72/150, 2),
        "readability_gate_pt": 8,
        "note": "SOURCE-ASSET READABILITY ONLY: designed font sizes are generally "
                "23–72 px. At the declared 150 dpi source density, the 23 px minimum "
                "is 11.04 pt (>8 pt). Manuscript placement is not checked here; use "
                "paper/final/final_placement_manifest.csv and the paper build gate. "
                "Contact sheets are QA previews, not publication-size placements.",
        "errors": errors,
    }
    (PUB/"validation_report.json").write_text(
        json.dumps(report, indent=2, sort_keys=True)+"\n", encoding="utf-8")
    if errors:
        raise SystemExit("validation failed: " + "; ".join(errors))
    return report


def checksums() -> None:
    checksum_path = PUB/"checksums.sha256"
    files = []
    for base in (PUB, PRES):
        for p in base.rglob("*"):
            if p.is_file() and p != checksum_path and ".tmp_kicad" not in p.parts:
                files.append(p)
    lines = []
    for p in sorted(files, key=lambda x: x.as_posix()):
        digest = hashlib.sha256(p.read_bytes()).hexdigest()
        lines.append(f"{digest}  {p.relative_to(ROOT).as_posix()}")
    checksum_path.write_text("\n".join(lines)+"\n", encoding="utf-8")


def clean_outputs() -> None:
    for path in (PUB, PRES):
        if path.exists():
            shutil.rmtree(path)
    PUB.mkdir(parents=True)
    PRES.mkdir(parents=True)


def main() -> None:
    os.environ.setdefault("SOURCE_DATE_EPOCH", "946684800")
    clean_outputs()
    specs = make_figures()
    flows = make_flows()
    make_storyboard_and_animation()
    captions(specs)
    make_evidence_map(specs)
    make_tables_and_images()
    make_contact_sheets(specs, flows)
    make_index(specs)
    make_presentation(specs)
    report = validate(specs)
    if (PUB/".tmp_kicad").exists():
        shutil.rmtree(PUB/".tmp_kicad")
    checksums()
    print(json.dumps(report, indent=2, sort_keys=True))


if __name__ == "__main__":
    main()
