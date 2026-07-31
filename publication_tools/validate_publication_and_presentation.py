#!/usr/bin/env python3
"""Standalone package validation without regenerating artifacts."""

from __future__ import annotations

import hashlib
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
PUB = ROOT / "build/publication_assets"
PRES = ROOT / "build/presentation"


def main() -> None:
    errors: list[str] = []
    for figure_id in ("F01", "F02", "F03", "F04", "F05", "F06", "T01", "T02", "T03", "T04"):
        base = PUB / "figures" / figure_id
        for name in ("figure.svg", "figure.pdf", "figure.png", "figure_bw.png",
                     "figure_16x9.png", "source.csv", "generate.py", "metadata.md"):
            path = base / name
            if not path.is_file() or path.stat().st_size == 0:
                errors.append(f"missing {path.relative_to(ROOT)}")
        if (base/"figure.png").is_file() and Image.open(base/"figure.png").size != (3200, 1800):
            errors.append(f"wrong figure dimensions: {figure_id}")
    for stem in ("F01_research_decision_flow", "F02_token_dataflow",
                 "F03_work_stealing_sequence", "F04_cycle_timeline",
                 "F05_physical_data_path"):
        for suffix in (".svg", ".pdf", ".png"):
            if not (PUB/"flow"/f"{stem}{suffix}").is_file():
                errors.append(f"missing flow {stem}{suffix}")
    for path in (PUB/"INDEX.md", PUB/"INDEX.pdf", PUB/"evidence_map.csv",
                 PUB/"captions_ko.md", PUB/"captions_en.md",
                 PUB/"contact_sheets/figures_contact_sheet.png",
                 PUB/"contact_sheets/flows_contact_sheet.png",
                 PUB/"flow/storyboard_16x9.pdf",
                 PUB/"tables/scheduler_core_metrics.csv",
                 PUB/"tables/blocked_evidence.csv",
                 PUB/"images/README.md",
                 PRES/"presentation.pptx", PRES/"presentation.pdf",
                 PRES/"speaker_notes.md", PRES/"qa_evidence_index.md",
                 PRES/"demo_script.md"):
        if not path.is_file() or path.stat().st_size == 0:
            errors.append(f"missing {path.relative_to(ROOT)}")
    if (PUB/"flow/work_stealing_animation.mp4").is_file():
        if (PUB/"flow/work_stealing_animation.mp4").stat().st_size < 10_000:
            errors.append("MP4 is unexpectedly small")
    prs = Presentation(PRES/"presentation.pptx")
    if len(prs.slides) != 12:
        errors.append(f"expected 12 slides, got {len(prs.slides)}")
    if not all(s.notes_slide.notes_text_frame.text.strip() for s in prs.slides):
        errors.append("speaker notes missing in PPTX")
    text = (PUB/"INDEX.md").read_text(encoding="utf-8")
    for token in ("Gemma 3 1B", "BLOCKED", "energy", "Vivado"):
        if token not in text:
            errors.append(f"INDEX missing boundary token {token}")
    checksum_file = PUB/"checksums.sha256"
    if checksum_file.is_file():
        for line in checksum_file.read_text(encoding="utf-8").splitlines():
            digest, rel = line.split("  ", 1)
            path = ROOT/rel
            if not path.is_file() or hashlib.sha256(path.read_bytes()).hexdigest() != digest:
                errors.append(f"checksum mismatch: {rel}")
    else:
        errors.append("checksums.sha256 missing")
    report = {"status": "PASS" if not errors else "FAIL", "errors": errors}
    print(json.dumps(report, indent=2, sort_keys=True))
    raise SystemExit(1 if errors else 0)


if __name__ == "__main__":
    main()
